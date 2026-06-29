#!/usr/bin/env python3
"""
Impact Slide Preprocessor - Step 1

Scans a folder of documents, images, spreadsheets, and PowerPoint/PDF decks and
creates structured intermediate files for the Impact Slide Analyst Custom GPT.

Outputs include:
- file_inventory.json / file_inventory.md
- extracted_documents.md
- excel_profile.json
- pptx_slide_audit.json
- brand_style_summary.json
- image_ocr_summary.json
- asset_manifest.json
- evidence_register_seed.json
- preprocessor_summary.md

Usage:
  python step1_preprocessor.py --input "C:\\path\\to\\materials" --output "C:\\path\\to\\output" --ocr
"""

from __future__ import annotations

import argparse
import csv
import datetime as dt
import hashlib
import io
import json
import math
import mimetypes
import os
import re
import statistics
import sys
import traceback
import warnings
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

# ----------------------------- optional imports -----------------------------

IMPORT_WARNINGS: Dict[str, str] = {}


def _optional_import(name: str, import_name: Optional[str] = None):
    try:
        return __import__(import_name or name)
    except Exception as exc:  # noqa: BLE001 - show install guidance, keep running
        IMPORT_WARNINGS[name] = str(exc)
        return None


pd = _optional_import("pandas")
np = _optional_import("numpy")
fitz = _optional_import("pymupdf", "fitz")
pdfplumber = _optional_import("pdfplumber")
docx = _optional_import("python-docx", "docx")
mammoth = _optional_import("mammoth")
pptx_pkg = _optional_import("python-pptx", "pptx")
PIL = _optional_import("pillow", "PIL")
pytesseract = _optional_import("pytesseract")
filetype = _optional_import("filetype")

try:
    import orjson  # type: ignore
except Exception as exc:  # noqa: BLE001
    IMPORT_WARNINGS["orjson"] = str(exc)
    orjson = None

try:
    import openpyxl  # type: ignore
    from openpyxl.utils import get_column_letter as _openpyxl_get_column_letter  # type: ignore
except Exception as exc:  # noqa: BLE001
    IMPORT_WARNINGS["openpyxl"] = str(exc)
    openpyxl = None
    _openpyxl_get_column_letter = None

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    from pptx.enum.dml import MSO_COLOR_TYPE  # type: ignore
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None
    MSO_COLOR_TYPE = None

try:
    from PIL import Image  # type: ignore
except Exception:
    Image = None


# ----------------------------- constants -----------------------------

DOC_EXTS = {".pdf", ".docx", ".doc", ".txt", ".md", ".markdown", ".rtf"}
SPREADSHEET_EXTS = {".xlsx", ".xlsm", ".xls", ".xlsb", ".ods", ".csv", ".tsv"}
PRESENTATION_EXTS = {".pptx", ".ppt", ".ppsx"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif"}

IGNORE_DIRS = {
    ".git",
    ".svn",
    "node_modules",
    "__pycache__",
    ".venv",
    "venv",
    "env",
    "ImpactSlide_Preprocessor_Output",
}

KEYWORD_WEIGHTS = {
    "risk": 3,
    "urgent": 3,
    "opportunity": 3,
    "increase": 2,
    "decrease": 2,
    "growth": 2,
    "decline": 2,
    "cost": 2,
    "revenue": 2,
    "profit": 2,
    "margin": 2,
    "customer": 2,
    "client": 2,
    "strategy": 2,
    "priority": 2,
    "roadmap": 2,
    "process": 2,
    "implementation": 2,
    "timeline": 2,
    "decision": 2,
    "approve": 2,
    "recommend": 2,
    "must": 2,
    "should": 1,
    "kpi": 2,
    "metric": 2,
    "target": 2,
}


# ----------------------------- helpers -----------------------------


def now_iso() -> str:
    return dt.datetime.now().replace(microsecond=0).isoformat()


def safe_relpath(path: Path, base: Path) -> str:
    try:
        return str(path.relative_to(base))
    except Exception:
        return str(path)


def clean_text(text: Any) -> str:
    if text is None:
        return ""
    text = str(text).replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[TRUNCATED after {max_chars:,} characters]"


def json_default(obj: Any) -> Any:
    if isinstance(obj, (dt.datetime, dt.date)):
        return obj.isoformat()
    if isinstance(obj, Path):
        return str(obj)
    if hasattr(obj, "item"):
        try:
            return obj.item()
        except Exception:
            pass
    if isinstance(obj, float) and (math.isnan(obj) or math.isinf(obj)):
        return None
    return str(obj)


def to_jsonable(obj: Any) -> Any:
    if obj is None or isinstance(obj, (str, int, bool)):
        return obj
    if isinstance(obj, float):
        if math.isnan(obj) or math.isinf(obj):
            return None
        return obj
    if isinstance(obj, Path):
        return str(obj)
    if isinstance(obj, (dt.datetime, dt.date)):
        return obj.isoformat()
    if isinstance(obj, dict):
        return {str(k): to_jsonable(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple, set)):
        return [to_jsonable(v) for v in obj]
    if hasattr(obj, "item"):
        try:
            return to_jsonable(obj.item())
        except Exception:
            pass
    return str(obj)


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    data = to_jsonable(data)
    if orjson is not None:
        path.write_bytes(orjson.dumps(data, option=orjson.OPT_INDENT_2 | orjson.OPT_SORT_KEYS))
    else:
        path.write_text(json.dumps(data, indent=2, ensure_ascii=False, default=json_default, sort_keys=True), encoding="utf-8")


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def sha256_file(path: Path, block_size: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(block_size), b""):
            h.update(block)
    return h.hexdigest()


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def get_column_letter(index_1_based: int) -> str:
    if _openpyxl_get_column_letter is not None:
        return _openpyxl_get_column_letter(index_1_based)
    # fallback implementation
    result = ""
    n = index_1_based
    while n:
        n, rem = divmod(n - 1, 26)
        result = chr(65 + rem) + result
    return result


def excel_addr(row_1_based: int, col_1_based: int) -> str:
    return f"{get_column_letter(col_1_based)}{row_1_based}"


def sniff_delimiter(path: Path) -> str:
    try:
        sample = path.read_text(encoding="utf-8", errors="ignore")[:4096]
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
        return dialect.delimiter
    except Exception:
        return "\t" if path.suffix.lower() == ".tsv" else ","


def markdown_table(headers: List[str], rows: List[List[Any]]) -> str:
    def esc(x: Any) -> str:
        return str(x).replace("|", "\\|").replace("\n", "<br>")

    out = ["| " + " | ".join(headers) + " |", "|" + "|".join(["---"] * len(headers)) + "|"]
    for row in rows:
        padded = list(row) + [""] * (len(headers) - len(row))
        out.append("| " + " | ".join(esc(v) for v in padded[: len(headers)]) + " |")
    return "\n".join(out)


def color_to_hex_tuple(rgb: Tuple[int, int, int]) -> str:
    return "#%02X%02X%02X" % rgb


def emu_to_inches(value: Any) -> Optional[float]:
    try:
        return round(float(value) / 914400.0, 3)
    except Exception:
        return None


def safe_stat(values: Iterable[float], fn) -> Optional[float]:
    vals = [float(v) for v in values if v is not None and not math.isnan(float(v))]
    if not vals:
        return None
    try:
        return round(float(fn(vals)), 4)
    except Exception:
        return None


def compact_value(value: Any, max_len: int = 120) -> str:
    s = clean_text(value)
    return s if len(s) <= max_len else s[: max_len - 3] + "..."


# ----------------------------- classification -----------------------------


def classify_category(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in DOC_EXTS:
        return "document"
    if ext in SPREADSHEET_EXTS:
        return "spreadsheet"
    if ext in PRESENTATION_EXTS:
        return "presentation"
    if ext in IMAGE_EXTS:
        return "image"
    return "other"


def supports_for_category(category: str, ext: str) -> Dict[str, bool]:
    return {
        "strategy_narrative": category in {"document", "presentation"},
        "evidence_data": category in {"spreadsheet", "document", "presentation", "image"},
        "visuals_assets": category in {"image", "presentation"},
        "brand_style": category in {"presentation", "image"},
        "existing_slide_updates": ext.lower() in PRESENTATION_EXTS or ext.lower() == ".pdf",
    }


def detect_mime(path: Path) -> Dict[str, Optional[str]]:
    guessed_mime, guessed_encoding = mimetypes.guess_type(str(path))
    detected = None
    if filetype is not None:
        try:
            kind = filetype.guess(str(path))
            if kind is not None:
                detected = kind.mime
        except Exception:
            detected = None
    return {"mimetypes_guess": guessed_mime, "encoding_guess": guessed_encoding, "filetype_guess": detected}


def classify_potential_use(text: str) -> str:
    t = text.lower()
    why_terms = ["risk", "urgent", "gap", "pain", "decline", "loss", "problem", "challenge", "why"]
    what_terms = ["strategy", "pillar", "priority", "solution", "approach", "model", "what"]
    how_terms = ["process", "workflow", "implementation", "roadmap", "phase", "step", "timeline", "how"]
    now_terms = ["decision", "approve", "next", "action", "owner", "deadline", "now", "recommend"]
    scores = {
        "Why": sum(term in t for term in why_terms),
        "What": sum(term in t for term in what_terms),
        "How": sum(term in t for term in how_terms),
        "Now": sum(term in t for term in now_terms),
    }
    best = max(scores, key=scores.get)
    return best if scores[best] else "What"


# ----------------------------- extraction utilities -----------------------------


def read_text_file(path: Path, max_chars: int) -> Dict[str, Any]:
    encodings = ["utf-8", "utf-8-sig", "cp1252", "latin-1"]
    last_error = None
    for enc in encodings:
        try:
            text = path.read_text(encoding=enc, errors="strict")
            return {"status": "ok", "encoding": enc, "text": truncate(clean_text(text), max_chars)}
        except Exception as exc:  # noqa: BLE001
            last_error = str(exc)
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return {"status": "partial", "encoding": "utf-8-ignore", "text": truncate(clean_text(text), max_chars)}
    except Exception as exc:  # noqa: BLE001
        return {"status": "error", "error": last_error or str(exc), "text": ""}


def split_candidate_sentences(text: str) -> List[str]:
    parts = re.split(r"(?<=[.!?])\s+|\n+|\r+", text)
    return [clean_text(p) for p in parts if len(clean_text(p)) >= 25]


def sentence_score(sentence: str) -> int:
    s = sentence.lower()
    score = 0
    if re.search(r"\b\d+(?:\.\d+)?\s?%", s):
        score += 4
    if re.search(r"[$€£¥]\s?\d", s):
        score += 4
    if re.search(r"\b\d{4}\b|\bq[1-4]\b|\bfy\d{2,4}\b", s):
        score += 2
    if re.search(r"\b\d+(?:\.\d+)?\b", s):
        score += 1
    for kw, weight in KEYWORD_WEIGHTS.items():
        if kw in s:
            score += weight
    if len(sentence) > 240:
        score -= 2
    return score


def top_candidate_sentences(text: str, max_items: int = 15) -> List[str]:
    sentences = split_candidate_sentences(text)
    scored = [(sentence_score(s), s) for s in sentences]
    scored = [(score, s) for score, s in scored if score > 0]
    scored.sort(key=lambda x: x[0], reverse=True)
    seen = set()
    result = []
    for _, sent in scored:
        key = sent[:120].lower()
        if key in seen:
            continue
        seen.add(key)
        result.append(sent)
        if len(result) >= max_items:
            break
    return result


# ----------------------------- main class -----------------------------


class ImpactSlidePreprocessor:
    def __init__(
        self,
        input_path: Path,
        output_dir: Path,
        recursive: bool = True,
        include_hidden: bool = False,
        enable_ocr: bool = False,
        hash_files: bool = False,
        max_text_chars_per_file: int = 120_000,
        max_sheet_rows: int = 10_000,
        max_evidence_per_file: int = 20,
    ) -> None:
        self.input_path = input_path.resolve()
        self.output_dir = output_dir.resolve()
        self.recursive = recursive
        self.include_hidden = include_hidden
        self.enable_ocr = enable_ocr
        self.hash_files = hash_files
        self.max_text_chars_per_file = max_text_chars_per_file
        self.max_sheet_rows = max_sheet_rows
        self.max_evidence_per_file = max_evidence_per_file

        self.inventory: List[Dict[str, Any]] = []
        self.documents: List[Dict[str, Any]] = []
        self.excel_profiles: List[Dict[str, Any]] = []
        self.pptx_audits: List[Dict[str, Any]] = []
        self.image_summaries: List[Dict[str, Any]] = []
        self.asset_manifest: List[Dict[str, Any]] = []
        self.evidence: List[Dict[str, Any]] = []
        self.errors: List[Dict[str, Any]] = []

    # ----------------------------- run -----------------------------

    def run(self) -> None:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        files = self.gather_files()
        self.inventory = self.build_file_inventory(files)

        for item in self.inventory:
            path = Path(item["absolute_path"])
            if item["access_status"] == "not_readable":
                continue
            category = item["category"]
            try:
                if category == "document":
                    self.documents.append(self.extract_document(path, item))
                elif category == "spreadsheet":
                    self.excel_profiles.append(self.extract_spreadsheet(path, item))
                elif category == "presentation":
                    self.pptx_audits.append(self.extract_presentation(path, item))
                    self.asset_manifest.extend(self.extract_pptx_media_manifest(path, item))
                elif category == "image":
                    img = self.extract_image(path, item)
                    self.image_summaries.append(img)
                    self.asset_manifest.append(self.image_asset_record(path, item, img))
            except Exception as exc:  # noqa: BLE001
                error = {
                    "file_id": item["file_id"],
                    "file_name": item["file_name"],
                    "category": category,
                    "error": str(exc),
                    "traceback": traceback.format_exc(limit=5),
                }
                self.errors.append(error)

        brand_summary = self.build_brand_style_summary()
        self.build_evidence_register(brand_summary)
        self.write_outputs(brand_summary)

    # ----------------------------- files/inventory -----------------------------

    def gather_files(self) -> List[Path]:
        if self.input_path.is_file():
            return [self.input_path]
        pattern = "**/*" if self.recursive else "*"
        files: List[Path] = []
        for path in self.input_path.glob(pattern):
            if not path.is_file():
                continue
            if not self.include_hidden and any(part.startswith(".") for part in path.parts):
                continue
            if any(part in IGNORE_DIRS for part in path.parts):
                continue
            try:
                # avoid ingesting our own output if output directory is inside input
                if self.output_dir in path.resolve().parents:
                    continue
            except Exception:
                pass
            files.append(path)
        return sorted(files, key=lambda p: str(p).lower())

    def build_file_inventory(self, files: List[Path]) -> List[Dict[str, Any]]:
        inventory: List[Dict[str, Any]] = []
        base = self.input_path if self.input_path.is_dir() else self.input_path.parent
        for idx, path in enumerate(files, 1):
            ext = path.suffix.lower()
            category = classify_category(path)
            stat = None
            access_status = "readable"
            limitations = ""
            try:
                stat = path.stat()
                with path.open("rb") as f:
                    f.read(1)
            except Exception as exc:  # noqa: BLE001
                access_status = "not_readable"
                limitations = str(exc)

            mime = detect_mime(path)
            row: Dict[str, Any] = {
                "file_id": f"F{idx:04d}",
                "file_name": path.name,
                "relative_path": safe_relpath(path, base),
                "absolute_path": str(path.resolve()),
                "extension": ext,
                "category": category,
                "size_bytes": stat.st_size if stat else None,
                "modified_time": dt.datetime.fromtimestamp(stat.st_mtime).isoformat() if stat else None,
                "mime": mime,
                "access_status": access_status,
                "key_usable_content": self.default_key_usable_content(category, ext),
                "limitations_notes": limitations,
                "supports": supports_for_category(category, ext),
            }
            if self.hash_files and access_status == "readable":
                try:
                    row["sha256"] = sha256_file(path)
                except Exception as exc:  # noqa: BLE001
                    row["sha256_error"] = str(exc)
            inventory.append(row)
        return inventory

    @staticmethod
    def default_key_usable_content(category: str, ext: str) -> str:
        if category == "document":
            return "Narrative, claims, quotes, requirements, context, source evidence"
        if category == "spreadsheet":
            return "Metrics, trends, comparisons, tables, chart-ready data"
        if category == "presentation":
            return "Existing slide structure, slide text, layouts, brand cues, embedded assets"
        if category == "image":
            return "Visual assets, OCR text, screenshots, charts, style/brand cues"
        return "Unsupported or miscellaneous file; inventory only"

    # ----------------------------- document extraction -----------------------------

    def extract_document(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        ext = path.suffix.lower()
        if ext == ".pdf":
            return self.extract_pdf(path, item)
        if ext in {".docx", ".doc"}:
            return self.extract_docx(path, item)
        return self.extract_plain_text(path, item)

    def extract_plain_text(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        result = read_text_file(path, self.max_text_chars_per_file)
        return {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "file_type": path.suffix.lower(),
            "extraction_method": "text_read",
            "status": result["status"],
            "encoding": result.get("encoding"),
            "text": result.get("text", ""),
            "error": result.get("error"),
        }

    def extract_pdf(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        pages: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []
        method = []
        error_messages = []

        if fitz is not None:
            try:
                doc = fitz.open(str(path))
                method.append("pymupdf")
                char_budget = self.max_text_chars_per_file
                for page_idx in range(len(doc)):
                    text = clean_text(doc[page_idx].get_text("text"))
                    if char_budget <= 0:
                        text = "[TRUNCATED: file text budget reached]"
                    elif len(text) > char_budget:
                        text = truncate(text, char_budget)
                    char_budget -= len(text)
                    pages.append(
                        {
                            "page_number": page_idx + 1,
                            "text": text,
                            "char_count": len(text),
                        }
                    )
                doc.close()
            except Exception as exc:  # noqa: BLE001
                error_messages.append(f"pymupdf: {exc}")

        if pdfplumber is not None:
            try:
                with pdfplumber.open(str(path)) as pdf:
                    method.append("pdfplumber")
                    for page_idx, page in enumerate(pdf.pages[: min(len(pdf.pages), 50)], 1):
                        try:
                            page_tables = page.extract_tables() or []
                            for table_idx, table in enumerate(page_tables[:5], 1):
                                tables.append(
                                    {
                                        "page_number": page_idx,
                                        "table_number": table_idx,
                                        "row_count": len(table),
                                        "column_count": max((len(r) for r in table), default=0),
                                        "sample_rows": table[:5],
                                    }
                                )
                        except Exception:
                            continue
                    if not pages:
                        # fallback text extraction if pymupdf unavailable/failed
                        char_budget = self.max_text_chars_per_file
                        for page_idx, page in enumerate(pdf.pages, 1):
                            text = clean_text(page.extract_text() or "")
                            if char_budget <= 0:
                                text = "[TRUNCATED: file text budget reached]"
                            elif len(text) > char_budget:
                                text = truncate(text, char_budget)
                            char_budget -= len(text)
                            pages.append({"page_number": page_idx, "text": text, "char_count": len(text)})
            except Exception as exc:  # noqa: BLE001
                error_messages.append(f"pdfplumber: {exc}")

        return {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "file_type": ".pdf",
            "extraction_method": "+".join(method) if method else "none",
            "status": "ok" if pages else "error",
            "page_count": len(pages),
            "pages": pages,
            "tables": tables,
            "errors": error_messages,
        }

    def extract_docx(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        paragraphs: List[Dict[str, Any]] = []
        tables: List[Dict[str, Any]] = []
        mammoth_markdown = None
        errors = []

        if docx is not None and path.suffix.lower() == ".docx":
            try:
                document = docx.Document(str(path))
                char_count = 0
                for idx, para in enumerate(document.paragraphs, 1):
                    text = clean_text(para.text)
                    if not text:
                        continue
                    if char_count >= self.max_text_chars_per_file:
                        break
                    text = truncate(text, max(0, self.max_text_chars_per_file - char_count))
                    char_count += len(text)
                    paragraphs.append(
                        {
                            "paragraph_number": idx,
                            "style": getattr(para.style, "name", None),
                            "text": text,
                        }
                    )
                for table_idx, table in enumerate(document.tables, 1):
                    rows = []
                    for row in table.rows[:10]:
                        rows.append([clean_text(cell.text) for cell in row.cells])
                    tables.append(
                        {
                            "table_number": table_idx,
                            "row_count": len(table.rows),
                            "column_count": len(table.columns),
                            "sample_rows": rows,
                        }
                    )
            except Exception as exc:  # noqa: BLE001
                errors.append(f"python-docx: {exc}")

        if mammoth is not None and path.suffix.lower() == ".docx":
            try:
                with path.open("rb") as f:
                    result = mammoth.convert_to_markdown(f)
                mammoth_markdown = truncate(clean_text(result.value), self.max_text_chars_per_file)
            except Exception as exc:  # noqa: BLE001
                errors.append(f"mammoth: {exc}")

        if path.suffix.lower() == ".doc" and not paragraphs:
            errors.append("Legacy .doc extraction is not supported directly. Convert to .docx or PDF for best results.")

        return {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "file_type": path.suffix.lower(),
            "extraction_method": "python-docx+mammoth",
            "status": "ok" if paragraphs or mammoth_markdown else "error",
            "paragraph_count": len(paragraphs),
            "paragraphs": paragraphs,
            "tables": tables,
            "mammoth_markdown": mammoth_markdown,
            "errors": errors,
        }

    # ----------------------------- spreadsheet extraction -----------------------------

    def extract_spreadsheet(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        if pd is None:
            return {
                "file_id": item["file_id"],
                "file_name": item["file_name"],
                "status": "error",
                "error": "pandas is not installed",
            }

        ext = path.suffix.lower()
        result: Dict[str, Any] = {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "file_type": ext,
            "status": "ok",
            "sheets": [],
            "workbook_metadata": {},
            "errors": [],
        }

        if ext in {".csv", ".tsv"}:
            try:
                delimiter = sniff_delimiter(path)
                raw = pd.read_csv(path, sep=delimiter, header=None, nrows=self.max_sheet_rows, dtype=object, engine="python")
                result["workbook_metadata"] = {"delimiter": delimiter, "sheet_count": 1}
                result["sheets"].append(self.profile_dataframe(raw, sheet_name="CSV", source_type="csv"))
            except Exception as exc:  # noqa: BLE001
                result["status"] = "error"
                result["errors"].append(str(exc))
            return result

        try:
            xls = pd.ExcelFile(path)
            result["workbook_metadata"]["sheet_names"] = xls.sheet_names
            result["workbook_metadata"]["sheet_count"] = len(xls.sheet_names)
            for sheet_name in xls.sheet_names:
                try:
                    raw = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=self.max_sheet_rows, dtype=object)
                    result["sheets"].append(self.profile_dataframe(raw, sheet_name=sheet_name, source_type="excel"))
                except Exception as exc:  # noqa: BLE001
                    result["sheets"].append({"sheet_name": sheet_name, "status": "error", "error": str(exc)})
        except Exception as exc:  # noqa: BLE001
            result["status"] = "error"
            result["errors"].append(str(exc))

        if ext in {".xlsx", ".xlsm"}:
            result["openpyxl_metadata"] = self.extract_openpyxl_metadata(path)

        return result

    def detect_header_row(self, raw: Any, search_rows: int = 20) -> Optional[int]:
        if pd is None or raw is None or raw.empty:
            return None
        search = min(search_rows, len(raw))
        best_row = None
        best_score = -1.0
        for idx in range(search):
            row = raw.iloc[idx].tolist()
            values = [v for v in row if not pd.isna(v) and clean_text(v)]
            if len(values) < 2:
                continue
            text_like = 0
            numeric_like = 0
            unique_values = len(set(clean_text(v).lower() for v in values))
            for v in values:
                s = clean_text(v)
                try:
                    float(s.replace(",", ""))
                    numeric_like += 1
                except Exception:
                    text_like += 1
            score = len(values) + text_like * 1.5 + unique_values * 0.2 - numeric_like * 0.5
            if score > best_score:
                best_score = score
                best_row = idx
        return best_row

    @staticmethod
    def make_unique_columns(values: List[Any]) -> List[str]:
        cols = []
        counts: Dict[str, int] = defaultdict(int)
        for idx, value in enumerate(values, 1):
            base = clean_text(value) or f"Column {get_column_letter(idx)}"
            base = base[:80]
            counts[base] += 1
            cols.append(base if counts[base] == 1 else f"{base}_{counts[base]}")
        return cols

    def profile_dataframe(self, raw: Any, sheet_name: str, source_type: str) -> Dict[str, Any]:
        if pd is None:
            return {"sheet_name": sheet_name, "status": "error", "error": "pandas missing"}

        raw_rows, raw_cols = raw.shape
        header_row = self.detect_header_row(raw)
        truncated = raw_rows >= self.max_sheet_rows

        if header_row is not None and header_row < raw_rows - 1:
            columns = self.make_unique_columns(raw.iloc[header_row].tolist())
            data = raw.iloc[header_row + 1 :].copy()
            data.columns = columns
            header_excel_row = header_row + 1
        else:
            columns = [f"Column {get_column_letter(i + 1)}" for i in range(raw_cols)]
            data = raw.copy()
            data.columns = columns
            header_excel_row = None

        non_empty_cells = int(raw.notna().sum().sum())
        total_cells = int(raw_rows * raw_cols)
        numeric_profiles = []
        date_profiles = []
        categorical_profiles = []
        findings = []

        for col_idx, col_name in enumerate(columns):
            series = data.iloc[:, col_idx]
            non_null = series.dropna()
            non_null_count = int(non_null.shape[0])
            if non_null_count == 0:
                continue

            # Numeric profiling
            numeric = pd.to_numeric(series.astype(str).str.replace(",", "", regex=False), errors="coerce")
            numeric_count = int(numeric.notna().sum())
            numeric_ratio = numeric_count / max(non_null_count, 1)
            if numeric_count >= 2 and numeric_ratio >= 0.6:
                desc_values = [float(v) for v in numeric.dropna().tolist()]
                max_idx = numeric.idxmax()
                min_idx = numeric.idxmin()
                max_value = numeric.loc[max_idx]
                min_value = numeric.loc[min_idx]
                max_cell = excel_addr(int(max_idx) + 1, col_idx + 1)
                min_cell = excel_addr(int(min_idx) + 1, col_idx + 1)
                profile = {
                    "column": col_name,
                    "source_column": get_column_letter(col_idx + 1),
                    "numeric_count": numeric_count,
                    "non_null_count": non_null_count,
                    "min": safe_stat(desc_values, min),
                    "max": safe_stat(desc_values, max),
                    "mean": safe_stat(desc_values, statistics.mean),
                    "median": safe_stat(desc_values, statistics.median),
                    "sum": safe_stat(desc_values, sum),
                    "max_cell": max_cell,
                    "max_value": json_default(max_value),
                    "min_cell": min_cell,
                    "min_value": json_default(min_value),
                }
                numeric_profiles.append(profile)
                findings.append(
                    {
                        "type": "numeric_range",
                        "text": f"{sheet_name}: '{col_name}' ranges from {profile['min']} to {profile['max']} (max at {max_cell}, min at {min_cell}).",
                        "location": f"Sheet '{sheet_name}', column {get_column_letter(col_idx + 1)}",
                    }
                )
                continue

            # Date profiling
            try:
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore", UserWarning)
                    dates = pd.to_datetime(series, errors="coerce")
                date_count = int(dates.notna().sum())
                date_ratio = date_count / max(non_null_count, 1)
                if date_count >= 2 and date_ratio >= 0.6:
                    date_profiles.append(
                        {
                            "column": col_name,
                            "source_column": get_column_letter(col_idx + 1),
                            "date_count": date_count,
                            "min_date": dates.min().isoformat() if not pd.isna(dates.min()) else None,
                            "max_date": dates.max().isoformat() if not pd.isna(dates.max()) else None,
                        }
                    )
                    continue
            except Exception:
                pass

            # Categorical profiling
            unique_count = int(non_null.astype(str).nunique())
            if unique_count <= 30 and unique_count / max(non_null_count, 1) <= 0.6:
                top_values = non_null.astype(str).value_counts().head(10)
                categorical_profiles.append(
                    {
                        "column": col_name,
                        "source_column": get_column_letter(col_idx + 1),
                        "unique_count": unique_count,
                        "top_values": [
                            {"value": str(idx), "count": int(count)} for idx, count in top_values.items()
                        ],
                    }
                )

        sample_rows = []
        try:
            for raw_idx, row in data.head(5).iterrows():
                sample_rows.append(
                    {
                        "excel_row": int(raw_idx) + 1,
                        "values": {str(k): compact_value(v) for k, v in row.to_dict().items()},
                    }
                )
        except Exception:
            sample_rows = []

        return {
            "sheet_name": sheet_name,
            "status": "ok",
            "source_type": source_type,
            "raw_shape": {"rows": int(raw_rows), "columns": int(raw_cols)},
            "truncated_at_max_rows": truncated,
            "header_row_guess": header_excel_row,
            "columns": columns,
            "non_empty_cells": non_empty_cells,
            "total_cells_scanned": total_cells,
            "empty_cell_ratio": round(1 - (non_empty_cells / max(total_cells, 1)), 4),
            "numeric_profiles": numeric_profiles,
            "date_profiles": date_profiles,
            "categorical_profiles": categorical_profiles,
            "findings": findings[:20],
            "sample_rows": sample_rows,
            "recommended_chart_types": self.recommend_chart_types(numeric_profiles, date_profiles, categorical_profiles),
        }

    @staticmethod
    def recommend_chart_types(numeric_profiles: List[Dict[str, Any]], date_profiles: List[Dict[str, Any]], categorical_profiles: List[Dict[str, Any]]) -> List[str]:
        recs = []
        if date_profiles and numeric_profiles:
            recs.append("Line chart for trends over time")
        if categorical_profiles and numeric_profiles:
            recs.append("Grouped or stacked bar chart for category comparisons")
        if len(numeric_profiles) >= 3:
            recs.append("KPI card cluster or metric dashboard for headline numbers")
        if len(categorical_profiles) >= 2 and len(numeric_profiles) >= 1:
            recs.append("Heatmap or conditional-format table for matrix comparison")
        if not recs and numeric_profiles:
            recs.append("Simple bar chart or KPI callout for numeric comparison")
        return recs

    def extract_openpyxl_metadata(self, path: Path) -> Dict[str, Any]:
        if openpyxl is None:
            return {"status": "skipped", "reason": "openpyxl not installed"}
        try:
            wb = openpyxl.load_workbook(str(path), read_only=False, data_only=False)
            sheets = []
            formula_samples = []
            formula_count = 0
            for ws in wb.worksheets:
                sheet_meta = {
                    "sheet_name": ws.title,
                    "sheet_state": ws.sheet_state,
                    "max_row": ws.max_row,
                    "max_column": ws.max_column,
                    "merged_ranges": [str(r) for r in list(ws.merged_cells.ranges)[:25]],
                    "table_names": list(ws.tables.keys())[:25] if hasattr(ws, "tables") else [],
                }
                sheets.append(sheet_meta)
                for row in ws.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.startswith("="):
                            formula_count += 1
                            if len(formula_samples) < 100:
                                formula_samples.append(
                                    {
                                        "sheet": ws.title,
                                        "cell": cell.coordinate,
                                        "formula": cell.value,
                                    }
                                )
            wb.close()
            return {
                "status": "ok",
                "sheets": sheets,
                "formula_count": formula_count,
                "formula_samples": formula_samples,
            }
        except Exception as exc:  # noqa: BLE001
            return {"status": "error", "error": str(exc)}

    # ----------------------------- presentation extraction -----------------------------

    def extract_presentation(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        ext = path.suffix.lower()
        if ext != ".pptx":
            return {
                "file_id": item["file_id"],
                "file_name": item["file_name"],
                "file_type": ext,
                "status": "unsupported",
                "error": "Only .pptx extraction is supported directly. Convert .ppt/.ppsx to .pptx for best results.",
            }
        if Presentation is None:
            return {
                "file_id": item["file_id"],
                "file_name": item["file_name"],
                "file_type": ext,
                "status": "error",
                "error": "python-pptx is not installed",
            }
        prs = Presentation(str(path))
        slides = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slides.append(self.extract_slide(slide, slide_idx))
        return {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "file_type": ext,
            "status": "ok",
            "slide_width_inches": emu_to_inches(prs.slide_width),
            "slide_height_inches": emu_to_inches(prs.slide_height),
            "slide_count": len(slides),
            "slides": slides,
        }

    def extract_slide(self, slide: Any, slide_idx: int) -> Dict[str, Any]:
        texts: List[str] = []
        font_records: List[Dict[str, Any]] = []
        colors: List[str] = []
        shapes_info = []
        tables = []
        charts = []
        pictures = []

        title = ""
        try:
            if slide.shapes.title is not None:
                title = clean_text(slide.shapes.title.text)
        except Exception:
            title = ""

        for shape_idx, shape in enumerate(slide.shapes, 1):
            shape_type = self.get_shape_type_name(shape)
            info = {
                "shape_index": shape_idx,
                "shape_type": shape_type,
                "name": getattr(shape, "name", None),
                "left_in": emu_to_inches(getattr(shape, "left", None)),
                "top_in": emu_to_inches(getattr(shape, "top", None)),
                "width_in": emu_to_inches(getattr(shape, "width", None)),
                "height_in": emu_to_inches(getattr(shape, "height", None)),
            }

            # text and fonts
            if getattr(shape, "has_text_frame", False):
                try:
                    text = clean_text(shape.text)
                    if text:
                        texts.append(text)
                    fr, c = self.extract_text_frame_fonts(shape.text_frame)
                    font_records.extend(fr)
                    colors.extend(c)
                except Exception:
                    pass

            # shape fill/line colors
            for c in self.extract_shape_colors(shape):
                colors.append(c)

            if getattr(shape, "has_table", False):
                try:
                    table = shape.table
                    sample = []
                    for r in range(min(len(table.rows), 5)):
                        sample.append([clean_text(table.cell(r, c).text) for c in range(min(len(table.columns), 8))])
                    tables.append(
                        {
                            "shape_index": shape_idx,
                            "row_count": len(table.rows),
                            "column_count": len(table.columns),
                            "sample_rows": sample,
                        }
                    )
                except Exception:
                    pass

            if getattr(shape, "has_chart", False):
                try:
                    chart = shape.chart
                    chart_info = {
                        "shape_index": shape_idx,
                        "chart_type": str(getattr(chart, "chart_type", "unknown")),
                        "has_title": bool(getattr(chart, "has_title", False)),
                        "title": clean_text(chart.chart_title.text_frame.text) if getattr(chart, "has_title", False) else "",
                        "series": [],
                    }
                    try:
                        for series in chart.series:
                            chart_info["series"].append({"name": clean_text(series.name)})
                    except Exception:
                        pass
                    charts.append(chart_info)
                except Exception:
                    pass

            if self.is_picture_shape(shape):
                pictures.append(info.copy())

            shapes_info.append(info)

        if not title and texts:
            title = texts[0].split("\n")[0][:120]

        notes_text = ""
        try:
            if getattr(slide, "has_notes_slide", False):
                notes_text = clean_text(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes_text = ""

        all_text = clean_text("\n".join(texts))
        word_count = len(re.findall(r"\b\w+\b", all_text))
        font_sizes = [r.get("size_pt") for r in font_records if r.get("size_pt")]
        small_font_count = len([s for s in font_sizes if s and s < 14])
        density_flags = []
        if word_count > 120:
            density_flags.append("High word count; consider splitting or reducing text")
        if len(shapes_info) > 25:
            density_flags.append("Many shapes; visual clutter risk")
        if small_font_count:
            density_flags.append(f"{small_font_count} text runs below 14pt; readability risk")
        if len(tables) and word_count > 80:
            density_flags.append("Table plus high text volume; consider appendix or simplification")

        color_counts = Counter([c for c in colors if c])
        font_counts = Counter([r.get("font_name") for r in font_records if r.get("font_name")])

        return {
            "slide_number": slide_idx,
            "layout_name": getattr(getattr(slide, "slide_layout", None), "name", None),
            "title": title,
            "text": all_text,
            "word_count": word_count,
            "notes_text": notes_text,
            "shape_count": len(shapes_info),
            "picture_count": len(pictures),
            "table_count": len(tables),
            "chart_count": len(charts),
            "shapes": shapes_info,
            "tables": tables,
            "charts": charts,
            "pictures": pictures,
            "top_fonts": [{"font": k, "count": v} for k, v in font_counts.most_common(10)],
            "top_colors": [{"color": k, "count": v} for k, v in color_counts.most_common(15)],
            "font_size_range_pt": {
                "min": round(min(font_sizes), 2) if font_sizes else None,
                "max": round(max(font_sizes), 2) if font_sizes else None,
            },
            "density_flags": density_flags,
            "recommended_action_hint": self.recommend_slide_action_hint(word_count, len(shapes_info), len(charts), len(tables), density_flags),
        }

    @staticmethod
    def get_shape_type_name(shape: Any) -> str:
        try:
            return str(shape.shape_type).split(".")[-1]
        except Exception:
            return "unknown"

    @staticmethod
    def is_picture_shape(shape: Any) -> bool:
        try:
            if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
        except Exception:
            pass
        return "PICTURE" in str(getattr(shape, "shape_type", "")).upper()

    def extract_text_frame_fonts(self, text_frame: Any) -> Tuple[List[Dict[str, Any]], List[str]]:
        records = []
        colors = []
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    font = run.font
                    color = self.color_format_to_string(getattr(font, "color", None))
                    if color:
                        colors.append(color)
                    records.append(
                        {
                            "text_sample": compact_value(run.text, 80),
                            "font_name": getattr(font, "name", None),
                            "size_pt": round(font.size.pt, 2) if getattr(font, "size", None) is not None else None,
                            "bold": getattr(font, "bold", None),
                            "italic": getattr(font, "italic", None),
                            "color": color,
                        }
                    )
                except Exception:
                    continue
        return records, colors

    def extract_shape_colors(self, shape: Any) -> List[str]:
        colors = []
        for attr in ["fill", "line"]:
            try:
                obj = getattr(shape, attr)
                if attr == "fill":
                    color_obj = obj.fore_color
                else:
                    color_obj = obj.color
                c = self.color_format_to_string(color_obj)
                if c:
                    colors.append(c)
            except Exception:
                continue
        return colors

    @staticmethod
    def color_format_to_string(color_format: Any) -> Optional[str]:
        if color_format is None:
            return None
        try:
            # RGB color
            rgb = getattr(color_format, "rgb", None)
            if rgb is not None:
                return f"#{str(rgb).upper()}"
        except Exception:
            pass
        try:
            theme_color = getattr(color_format, "theme_color", None)
            if theme_color is not None:
                return f"theme:{theme_color}"
        except Exception:
            pass
        return None

    @staticmethod
    def recommend_slide_action_hint(word_count: int, shape_count: int, chart_count: int, table_count: int, density_flags: List[str]) -> str:
        if word_count > 180 or shape_count > 35:
            return "Split or Convert"
        if density_flags:
            return "Revise or Brand Refresh"
        if chart_count or table_count:
            return "Keep data, improve visual clarity"
        return "Keep or minor polish"

    # ----------------------------- image extraction -----------------------------

    def extract_image(self, path: Path, item: Dict[str, Any]) -> Dict[str, Any]:
        if Image is None:
            return {
                "file_id": item["file_id"],
                "file_name": item["file_name"],
                "status": "error",
                "error": "Pillow is not installed",
            }
        result: Dict[str, Any] = {
            "file_id": item["file_id"],
            "file_name": item["file_name"],
            "status": "ok",
            "ocr_enabled": self.enable_ocr,
        }
        try:
            with Image.open(str(path)) as img:
                result.update(
                    {
                        "format": img.format,
                        "mode": img.mode,
                        "width": img.width,
                        "height": img.height,
                        "aspect_ratio": round(img.width / img.height, 4) if img.height else None,
                        "top_colors": self.extract_image_top_colors(img),
                    }
                )
        except Exception as exc:  # noqa: BLE001
            result["status"] = "error"
            result["error"] = str(exc)
            return result

        if self.enable_ocr:
            result["ocr"] = self.run_ocr(path)
        else:
            result["ocr"] = {"status": "skipped", "reason": "Use --ocr to enable OCR"}
        result["reuse_recommendation_hint"] = self.image_reuse_hint(result)
        return result

    @staticmethod
    def extract_image_top_colors(img: Any, max_colors: int = 8) -> List[Dict[str, Any]]:
        try:
            image = img.convert("RGB")
            image.thumbnail((200, 200))
            # Quantize for stable top colors
            q = image.quantize(colors=16, method=0).convert("RGB")
            colors = q.getcolors(maxcolors=200 * 200) or []
            total = sum(count for count, _ in colors) or 1
            ranked = sorted(colors, key=lambda x: x[0], reverse=True)[:max_colors]
            return [
                {
                    "hex": color_to_hex_tuple(rgb),
                    "count": int(count),
                    "percent": round(count / total, 4),
                }
                for count, rgb in ranked
            ]
        except Exception:
            return []

    def run_ocr(self, path: Path) -> Dict[str, Any]:
        if pytesseract is None or Image is None:
            return {"status": "error", "error": "pytesseract or Pillow is not installed"}
        try:
            with Image.open(str(path)) as img:
                text = pytesseract.image_to_string(img)
            return {"status": "ok", "text": clean_text(text), "char_count": len(clean_text(text))}
        except Exception as exc:  # noqa: BLE001
            return {"status": "error", "error": str(exc)}

    @staticmethod
    def image_reuse_hint(summary: Dict[str, Any]) -> str:
        if summary.get("status") != "ok":
            return "Avoid until readable"
        width = summary.get("width") or 0
        height = summary.get("height") or 0
        ocr_text = clean_text(summary.get("ocr", {}).get("text", ""))
        if width >= 1200 and height >= 700:
            return "Reuse or crop if visually relevant"
        if ocr_text:
            return "Use as evidence/reference or redraw cleanly"
        return "Use as reference or small asset; check resolution before reuse"

    def image_asset_record(self, path: Path, item: Dict[str, Any], img_summary: Dict[str, Any]) -> Dict[str, Any]:
        return {
            "asset_id": f"A-{item['file_id']}",
            "source_file_id": item["file_id"],
            "source_file_name": item["file_name"],
            "asset_type": "standalone_image",
            "path": str(path),
            "width": img_summary.get("width"),
            "height": img_summary.get("height"),
            "format": img_summary.get("format"),
            "top_colors": img_summary.get("top_colors", []),
            "reuse_recommendation_hint": img_summary.get("reuse_recommendation_hint"),
        }

    # ----------------------------- asset extraction -----------------------------

    def extract_pptx_media_manifest(self, path: Path, item: Dict[str, Any]) -> List[Dict[str, Any]]:
        assets = []
        if path.suffix.lower() != ".pptx":
            return assets
        try:
            with zipfile.ZipFile(str(path), "r") as zf:
                for idx, name in enumerate([n for n in zf.namelist() if n.startswith("ppt/media/")], 1):
                    data = zf.read(name)
                    rec: Dict[str, Any] = {
                        "asset_id": f"A-{item['file_id']}-M{idx:03d}",
                        "source_file_id": item["file_id"],
                        "source_file_name": item["file_name"],
                        "asset_type": "pptx_embedded_media",
                        "internal_path": name,
                        "size_bytes": len(data),
                        "sha256": sha256_bytes(data),
                    }
                    if Image is not None:
                        try:
                            with Image.open(io.BytesIO(data)) as img:
                                rec.update(
                                    {
                                        "format": img.format,
                                        "width": img.width,
                                        "height": img.height,
                                        "top_colors": self.extract_image_top_colors(img),
                                    }
                                )
                        except Exception:
                            pass
                    assets.append(rec)
        except Exception as exc:  # noqa: BLE001
            self.errors.append({"file_id": item["file_id"], "file_name": item["file_name"], "asset_error": str(exc)})
        return assets

    # ----------------------------- brand summary -----------------------------

    def build_brand_style_summary(self) -> Dict[str, Any]:
        pptx_files = [p for p in self.pptx_audits if p.get("status") == "ok"]
        color_counter: Counter[str] = Counter()
        font_counter: Counter[str] = Counter()
        layout_counter: Counter[str] = Counter()
        title_positions = []
        chart_count = 0
        table_count = 0
        picture_positions = []

        for deck in pptx_files:
            for slide in deck.get("slides", []):
                if slide.get("layout_name"):
                    layout_counter[slide["layout_name"]] += 1
                for c in slide.get("top_colors", []):
                    color_counter[c["color"]] += c["count"]
                for f in slide.get("top_fonts", []):
                    font_counter[f["font"]] += f["count"]
                chart_count += slide.get("chart_count", 0)
                table_count += slide.get("table_count", 0)
                # approximate title box position from first text shape if available
                for shp in slide.get("shapes", [])[:3]:
                    if shp.get("shape_type") and "PLACEHOLDER" in shp.get("shape_type", "").upper():
                        title_positions.append({"left_in": shp.get("left_in"), "top_in": shp.get("top_in")})
                        break
                for pic in slide.get("pictures", []):
                    picture_positions.append(
                        {
                            "left_in": pic.get("left_in"),
                            "top_in": pic.get("top_in"),
                            "width_in": pic.get("width_in"),
                            "height_in": pic.get("height_in"),
                        }
                    )

        image_color_counter: Counter[str] = Counter()
        for img in self.image_summaries:
            for c in img.get("top_colors", []):
                image_color_counter[c["hex"]] += c["count"]

        summary = {
            "available": bool(pptx_files or self.image_summaries),
            "source_decks": [d.get("file_name") for d in pptx_files],
            "colors": {
                "pptx_top_colors": [{"color": k, "count": v} for k, v in color_counter.most_common(15)],
                "image_top_colors": [{"color": k, "count": v} for k, v in image_color_counter.most_common(15)],
                "interpretation": "Top colors are extracted from visible slide elements/images. Theme colors may appear as theme:* when exact RGB is unavailable.",
            },
            "typography": {
                "top_fonts": [{"font": k, "count": v} for k, v in font_counter.most_common(15)],
                "interpretation": "Fonts are extracted from visible text runs. Missing fonts may use theme defaults not exposed in runs.",
            },
            "layout_grid": {
                "slide_sizes": [
                    {
                        "file_name": d.get("file_name"),
                        "width_in": d.get("slide_width_inches"),
                        "height_in": d.get("slide_height_inches"),
                    }
                    for d in pptx_files
                ],
                "layout_names": [{"layout": k, "count": v} for k, v in layout_counter.most_common(15)],
                "title_position_samples": title_positions[:20],
            },
            "charts_tables": {
                "chart_count": chart_count,
                "table_count": table_count,
                "implication": "Use observed chart/table density and style as guidance; validate data visuals manually before final deck generation.",
            },
            "logo_footer": {
                "picture_position_samples": picture_positions[:20],
                "interpretation": "Repeated small pictures near corners may indicate logos or footer assets; confirm visually.",
            },
            "overall_feel_hint": self.infer_overall_feel(color_counter, font_counter, chart_count, table_count),
        }
        return summary

    @staticmethod
    def infer_overall_feel(color_counter: Counter[str], font_counter: Counter[str], chart_count: int, table_count: int) -> str:
        if not color_counter and not font_counter:
            return "No brand system detected; use clean neutral default unless user provides brand assets."
        if chart_count + table_count >= 5:
            return "Data-oriented / corporate; favor dashboards, charts, and clear hierarchy."
        top_colors = [c for c, _ in color_counter.most_common(5)]
        if any(c.upper() in {"#000000", "#FFFFFF"} for c in top_colors):
            return "Minimal / high-contrast; use whitespace and restrained accent colors."
        return "Brand cues detected; use extracted colors/fonts/layouts as design constraints."

    # ----------------------------- evidence -----------------------------

    def add_evidence(
        self,
        source_file: str,
        exact_location: str,
        evidence_type: str,
        key_finding: str,
        potential_slide_use: Optional[str] = None,
        confidence: str = "medium",
        metadata: Optional[Dict[str, Any]] = None,
    ) -> None:
        eid = f"E{len(self.evidence) + 1:04d}"
        self.evidence.append(
            {
                "evidence_id": eid,
                "source_file": source_file,
                "exact_location": exact_location,
                "evidence_type": evidence_type,
                "key_finding": clean_text(key_finding),
                "potential_slide_use": potential_slide_use or classify_potential_use(key_finding),
                "confidence": confidence,
                "metadata": metadata or {},
            }
        )

    def build_evidence_register(self, brand_summary: Dict[str, Any]) -> None:
        # Documents
        for doc in self.documents:
            if doc.get("status") not in {"ok", "partial"}:
                continue
            source = doc.get("file_name", "")
            ext = doc.get("file_type")
            count = 0
            if ext == ".pdf":
                for page in doc.get("pages", []):
                    for sent in top_candidate_sentences(page.get("text", ""), self.max_evidence_per_file):
                        self.add_evidence(source, f"Page {page.get('page_number')}", "claim_or_metric", sent, confidence="medium")
                        count += 1
                        if count >= self.max_evidence_per_file:
                            break
                    if count >= self.max_evidence_per_file:
                        break
                for table in doc.get("tables", [])[:5]:
                    self.add_evidence(
                        source,
                        f"Page {table.get('page_number')}, table {table.get('table_number')}",
                        "table",
                        f"Detected table with {table.get('row_count')} rows and {table.get('column_count')} columns.",
                        potential_slide_use="What",
                        confidence="medium",
                        metadata={"sample_rows": table.get("sample_rows")},
                    )
            elif ext in {".docx", ".doc"}:
                text_blob = "\n".join(p.get("text", "") for p in doc.get("paragraphs", [])) or doc.get("mammoth_markdown", "")
                for sent in top_candidate_sentences(text_blob, self.max_evidence_per_file):
                    self.add_evidence(source, "DOCX paragraphs", "claim_or_quote", sent, confidence="medium")
            else:
                for sent in top_candidate_sentences(doc.get("text", ""), self.max_evidence_per_file):
                    self.add_evidence(source, "Text file", "claim_or_metric", sent, confidence="medium")

        # Spreadsheets
        for workbook in self.excel_profiles:
            source = workbook.get("file_name", "")
            for sheet in workbook.get("sheets", []):
                sheet_name = sheet.get("sheet_name", "")
                for finding in sheet.get("findings", [])[: self.max_evidence_per_file]:
                    self.add_evidence(
                        source,
                        finding.get("location", f"Sheet '{sheet_name}'"),
                        "metric",
                        finding.get("text", ""),
                        potential_slide_use="What",
                        confidence="high",
                    )
                for rec in sheet.get("recommended_chart_types", [])[:3]:
                    self.add_evidence(
                        source,
                        f"Sheet '{sheet_name}' profile",
                        "visual_recommendation",
                        f"Recommended data visual: {rec}",
                        potential_slide_use="How",
                        confidence="medium",
                    )

        # Presentations
        for deck in self.pptx_audits:
            source = deck.get("file_name", "")
            for slide in deck.get("slides", []):
                title = slide.get("title") or f"Slide {slide.get('slide_number')}"
                if title:
                    self.add_evidence(
                        source,
                        f"Slide {slide.get('slide_number')}",
                        "existing_slide_content",
                        f"Existing slide title/content: {title}",
                        potential_slide_use="What",
                        confidence="high",
                    )
                for flag in slide.get("density_flags", []):
                    self.add_evidence(
                        source,
                        f"Slide {slide.get('slide_number')}",
                        "existing_slide_issue",
                        flag,
                        potential_slide_use="Brand",
                        confidence="high",
                    )

        # Images/OCR
        for img in self.image_summaries:
            source = img.get("file_name", "")
            ocr = img.get("ocr", {})
            if ocr.get("status") == "ok" and clean_text(ocr.get("text", "")):
                for sent in top_candidate_sentences(ocr.get("text", ""), 8):
                    self.add_evidence(source, "Image OCR", "visual_text", sent, potential_slide_use="What", confidence="medium")
            if img.get("top_colors"):
                self.add_evidence(
                    source,
                    "Image dominant colors",
                    "brand_cue",
                    "Dominant image colors: " + ", ".join(c.get("hex", "") for c in img.get("top_colors", [])[:5]),
                    potential_slide_use="Brand",
                    confidence="medium",
                )

        # Brand summary cues
        if brand_summary.get("available"):
            colors = brand_summary.get("colors", {}).get("pptx_top_colors") or brand_summary.get("colors", {}).get("image_top_colors") or []
            fonts = brand_summary.get("typography", {}).get("top_fonts") or []
            if colors:
                self.add_evidence(
                    "Brand Style Summary",
                    "Aggregated files",
                    "brand_cue",
                    "Observed top colors: " + ", ".join(c.get("color", "") for c in colors[:8]),
                    potential_slide_use="Brand",
                    confidence="medium",
                )
            if fonts:
                self.add_evidence(
                    "Brand Style Summary",
                    "Aggregated PPTX text runs",
                    "brand_cue",
                    "Observed top fonts: " + ", ".join(f.get("font", "") for f in fonts[:8]),
                    potential_slide_use="Brand",
                    confidence="medium",
                )

    # ----------------------------- outputs -----------------------------

    def write_outputs(self, brand_summary: Dict[str, Any]) -> None:
        write_json(self.output_dir / "file_inventory.json", self.inventory)
        write_text(self.output_dir / "file_inventory.md", self.render_file_inventory_md())
        write_text(self.output_dir / "extracted_documents.md", self.render_extracted_documents_md())
        write_json(self.output_dir / "excel_profile.json", self.excel_profiles)
        write_json(self.output_dir / "pptx_slide_audit.json", self.pptx_audits)
        write_json(self.output_dir / "brand_style_summary.json", brand_summary)
        write_json(self.output_dir / "image_ocr_summary.json", self.image_summaries)
        write_json(self.output_dir / "asset_manifest.json", self.asset_manifest)
        write_json(self.output_dir / "evidence_register_seed.json", self.evidence)
        write_json(self.output_dir / "processing_errors.json", self.errors)
        write_text(self.output_dir / "preprocessor_summary.md", self.render_summary_md(brand_summary))

    def render_file_inventory_md(self) -> str:
        rows = []
        for f in self.inventory:
            rows.append(
                [
                    f.get("file_id"),
                    f.get("file_name"),
                    f.get("category"),
                    f.get("access_status"),
                    f.get("key_usable_content"),
                    f.get("limitations_notes", ""),
                ]
            )
        return "# File Inventory\n\n" + markdown_table(
            ["File ID", "File Name", "Type", "Access Status", "Key Usable Content", "Limitations / Notes"], rows
        )

    def render_extracted_documents_md(self) -> str:
        parts = ["# Extracted Documents\n"]
        if not self.documents:
            parts.append("No document files were extracted.\n")
            return "\n".join(parts)
        for doc in self.documents:
            parts.append(f"\n---\n\n## {doc.get('file_name')}\n")
            parts.append(f"- **File ID:** {doc.get('file_id')}\n")
            parts.append(f"- **Status:** {doc.get('status')}\n")
            if doc.get("file_type") == ".pdf":
                for page in doc.get("pages", []):
                    parts.append(f"\n### Page {page.get('page_number')}\n\n")
                    parts.append(page.get("text", "") + "\n")
            elif doc.get("file_type") in {".docx", ".doc"}:
                for para in doc.get("paragraphs", []):
                    style = para.get("style") or "Paragraph"
                    parts.append(f"\n**{style} {para.get('paragraph_number')}:** {para.get('text')}\n")
                if doc.get("mammoth_markdown"):
                    parts.append("\n### Mammoth Markdown\n\n")
                    parts.append(doc.get("mammoth_markdown") + "\n")
            else:
                parts.append("\n")
                parts.append(doc.get("text", "") + "\n")
        return "\n".join(parts)

    def render_summary_md(self, brand_summary: Dict[str, Any]) -> str:
        category_counts = Counter(f.get("category") for f in self.inventory)
        output_files = [
            "file_inventory.json / file_inventory.md",
            "extracted_documents.md",
            "excel_profile.json",
            "pptx_slide_audit.json",
            "brand_style_summary.json",
            "image_ocr_summary.json",
            "asset_manifest.json",
            "evidence_register_seed.json",
            "processing_errors.json",
        ]
        parts = [
            "# Impact Slide Preprocessor Summary",
            "",
            f"Generated: {now_iso()}",
            f"Input: `{self.input_path}`",
            f"Output: `{self.output_dir}`",
            "",
            "## File Counts",
            "",
        ]
        for cat, count in category_counts.most_common():
            parts.append(f"- **{cat}:** {count}")
        parts.extend(
            [
                "",
                "## Outputs Created",
                "",
                *[f"- `{name}`" for name in output_files],
                "",
                "## Evidence Seed",
                "",
                f"- Evidence items created: **{len(self.evidence)}**",
                "",
                "## Brand Summary",
                "",
                f"- Brand cues available: **{brand_summary.get('available')}**",
                f"- Overall feel hint: {brand_summary.get('overall_feel_hint')}",
                "",
                "## Processing Warnings",
                "",
            ]
        )
        if IMPORT_WARNINGS:
            parts.append("### Missing/optional package warnings")
            for pkg, warning in sorted(IMPORT_WARNINGS.items()):
                # Only show packages that matter if missing; many optional imports may not be installed
                parts.append(f"- `{pkg}`: {warning}")
        if self.errors:
            parts.append("\n### File processing errors")
            for err in self.errors:
                parts.append(f"- `{err.get('file_name')}`: {err.get('error') or err.get('asset_error')}")
        if not IMPORT_WARNINGS and not self.errors:
            parts.append("No warnings or errors captured.")
        parts.extend(
            [
                "",
                "## Recommended Next Step",
                "",
                "Upload/paste these files into **Impact Slide Analyst**:",
                "",
                "1. `file_inventory.md`",
                "2. `extracted_documents.md`",
                "3. `excel_profile.json`",
                "4. `pptx_slide_audit.json`",
                "5. `brand_style_summary.json`",
                "6. `image_ocr_summary.json`",
                "7. `evidence_register_seed.json`",
                "",
                "Ask Impact Slide Analyst to create the Alignment Summary and Slide Update Plan.",
            ]
        )
        return "\n".join(parts) + "\n"


# ----------------------------- CLI -----------------------------


def default_output_dir(input_path: Path) -> Path:
    stamp = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    if input_path.is_dir():
        return input_path / f"ImpactSlide_Preprocessor_Output_{stamp}"
    return input_path.parent / f"ImpactSlide_Preprocessor_Output_{stamp}"


def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Preprocess presentation source files for Impact Slide Analyst.")
    parser.add_argument("--input", "-i", required=True, help="Input file or folder containing source materials.")
    parser.add_argument("--output", "-o", default=None, help="Output folder. Defaults to a timestamped folder near input.")
    parser.add_argument("--no-recursive", action="store_true", help="Do not scan subfolders.")
    parser.add_argument("--include-hidden", action="store_true", help="Include hidden dotfiles/folders.")
    parser.add_argument("--ocr", action="store_true", help="Enable OCR for images using pytesseract/Tesseract.")
    parser.add_argument("--hash-files", action="store_true", help="Compute SHA-256 hashes for source files. Slower for large folders.")
    parser.add_argument("--max-text-chars-per-file", type=int, default=120_000, help="Max extracted text chars per document file.")
    parser.add_argument("--max-sheet-rows", type=int, default=10_000, help="Max rows to scan per spreadsheet sheet.")
    parser.add_argument("--max-evidence-per-file", type=int, default=20, help="Max seeded evidence items per file/type.")
    return parser.parse_args(argv)


def main(argv: Optional[List[str]] = None) -> int:
    args = parse_args(argv)
    input_path = Path(args.input).expanduser()
    if not input_path.exists():
        print(f"ERROR: input path does not exist: {input_path}", file=sys.stderr)
        return 2
    output_dir = Path(args.output).expanduser() if args.output else default_output_dir(input_path)

    pre = ImpactSlidePreprocessor(
        input_path=input_path,
        output_dir=output_dir,
        recursive=not args.no_recursive,
        include_hidden=args.include_hidden,
        enable_ocr=args.ocr,
        hash_files=args.hash_files,
        max_text_chars_per_file=args.max_text_chars_per_file,
        max_sheet_rows=args.max_sheet_rows,
        max_evidence_per_file=args.max_evidence_per_file,
    )
    pre.run()
    print(f"Done. Outputs written to: {output_dir}")
    print("Start with preprocessor_summary.md")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
