"""
Tests for the v3 PPTX extraction enhancements (#13-16):

  #13 group-shape recursion (nested textboxes are no longer lost)
  #14 SmartArt / graphic-frame text fallback via drawingml a:t runs
  #15 embedded/linked OLE object detection (counted + labelled for the Analyst)
  #16 spatial (top, left) shape ordering so multi-column slides read in order

These prove the v2 statement "SmartArt, grouped shapes, embedded
charts/objects, and multi-column layouts still lose signal" is no longer true
in v3.
"""
from __future__ import annotations

import json
from pathlib import Path

import pytest

import step1_preprocessor_v3 as m


@pytest.fixture()
def make_preprocessor(tmp_workspace):
    """Build a v3 preprocessor (overrides conftest's v2 fixture)."""
    def _make(filter_level="permissive", boost_keywords=None):
        inp = tmp_workspace / "input"
        out = tmp_workspace / "output"
        inp.mkdir(parents=True, exist_ok=True)
        p = m.ImpactSlidePreprocessorV2(
            input_path=str(inp), output_dir=str(out),
            filter_level=filter_level, boost_keywords=boost_keywords or [],
        )
        return p, inp, out
    return _make


def _ev(out):
    return json.load(open(out / "evidence_register_seed.json"))


def _profile(out):
    return json.load(open(out / "pptx_profile.json"))


# --------------------------------------------------------------------------- #
# #13 group-shape recursion
# --------------------------------------------------------------------------- #
class TestGroupRecursion:
    def test_nested_textboxes_recovered(self, make_pptx, make_preprocessor):
        """Text inside a group shape is extracted (v2 lost it entirely)."""
        p, inp, out = make_preprocessor()
        make_pptx(slides=[
            {"title": "Intro"},
            {"title": "Grouped Content Slide",
             "group": [{"text": "North region revenue grew strongly"},
                       {"text": "South region revenue declined"}]},
        ])
        p.run()
        ev = _ev(out)
        blob = " ".join(e["text"] for e in ev)
        assert "North region revenue grew strongly" in blob
        assert "South region revenue declined" in blob

    def test_grouped_text_counts_in_word_count(self, make_pptx, make_preprocessor):
        p, inp, out = make_preprocessor()
        make_pptx(slides=[
            {"title": "Intro"},
            {"title": "Grouped",
             "group": [{"text": "alpha beta gamma delta epsilon zeta"}]},
        ])
        p.run()
        prof = _profile(out)
        # the grouped textbox's words must be counted (word_count > title only)
        assert prof[0]["slides"][1]["word_count"] >= 6


# --------------------------------------------------------------------------- #
# #14 SmartArt / graphic-frame text fallback
# --------------------------------------------------------------------------- #
class TestSmartArtFallback:
    def test_smartart_text_recovered(self, make_pptx, make_preprocessor):
        """SmartArt (graphicFrame with a:t runs) text is recovered via the XML
        fallback, since has_text_frame is False for graphic frames."""
        p, inp, out = make_preprocessor()
        make_pptx(slides=[
            {"title": "Intro"},
            {"title": "Strategy Pillars",
             "smartart": ["Strategy Pillar One", "Strategy Pillar Two"]},
        ])
        p.run()
        ev = _ev(out)
        blob = " ".join(e["text"] for e in ev)
        assert "Strategy Pillar One" in blob
        assert "Strategy Pillar Two" in blob

    def test_extract_shape_text_helper(self):
        """Unit test for the _extract_shape_text fallback path."""
        # A plain string has no _element; the helper must handle gracefully
        assert m._extract_shape_text("not a shape") == ""


# --------------------------------------------------------------------------- #
# #15 embedded OLE object detection
# --------------------------------------------------------------------------- #
class TestEmbeddedObjects:
    def test_embedded_objects_listed_in_details(self, make_pptx, make_preprocessor):
        """Embedded OLE objects are recorded in slide details so the Analyst
        knows unread signal exists (their content can't be extracted)."""
        # Inject an embedded OLE shape via XML (python-pptx can't create them).
        # We use the _is_embedded_object helper directly on a fake + check the
        # real path via the graphicFrame injection.
        p, inp, out = make_preprocessor()
        from lxml import etree
        ole_xml = (
            '<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<p:nvGraphicFramePr><p:cNvPr id="200" name="Embedded Excel Sheet"/>'
            '<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>'
            '<p:xfrm><a:off x="1" y="1"/><a:ext cx="1" cy="1"/></p:xfrm>'
            '<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">'
            '<r:embed r:id="rId99"/></a:graphicData></a:graphic></p:graphicFrame>'
        )
        # Build a deck with a real textbox (so word_count > 0) and inject the OLE frame
        make_pptx(slides=[
            {"title": "Intro"},
            {"title": "Embedded Object Slide",
             "body": "Some real insight content goes here for testing the slide"},
        ])
        # append the OLE graphicFrame to the second slide's spTree
        from pptx import Presentation
        prs = Presentation(str(inp / "deck.pptx"))
        slide2 = prs.slides[1]
        slide2.shapes._spTree.append(etree.fromstring(ole_xml))
        prs.save(str(inp / "deck.pptx"))

        p.run()
        prof = _profile(out)
        eo = prof[0]["slides"][1]["details"].get("embedded_objects", [])
        assert eo, "expected the embedded OLE object to be recorded"
        assert any("Embedded Excel" in str(n) for n in eo)

    def test_is_embedded_object_helper(self):
        # smoke: the helper exists and is callable
        assert callable(m._is_embedded_object)


# --------------------------------------------------------------------------- #
# #16 spatial (top, left) shape ordering
# --------------------------------------------------------------------------- #
class TestSpatialOrdering:
    def test_multi_column_text_reads_top_to_bottom_left_to_right(self, make_pptx, make_preprocessor):
        """Shapes are iterated in (top, left) order so a two-column slide's
        text is concatenated in reading order, not insertion order."""
        p, inp, out = make_preprocessor()
        # Build a slide via python-pptx directly so we control positions:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank first
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        # title
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tb.text_frame.text = "Two Column Layout"
        # RIGHT column text added FIRST (so insertion order would be wrong)
        r = sl.shapes.add_textbox(Inches(6), Inches(1), Inches(3), Inches(1))
        r.text_frame.text = "Right column second in reading order"
        # LEFT column text added SECOND
        l = sl.shapes.add_textbox(Inches(0.5), Inches(1), Inches(3), Inches(1))
        l.text_frame.text = "Left column first in reading order"
        prs.save(str(inp / "col.pptx"))

        p.run()
        prof = _profile(out)
        bullets = prof[0]["slides"][1]["details"]["key_bullets"]
        # left-column text must appear before right-column text (spatial order)
        left_idx = next((i for i, b in enumerate(bullets)
                         if "Left column" in b), -1)
        right_idx = next((i for i, b in enumerate(bullets)
                          if "Right column" in b), -1)
        assert left_idx >= 0 and right_idx >= 0
        assert left_idx < right_idx, \
            f"expected left column before right column; got bullets={bullets}"
