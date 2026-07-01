"""Shared pytest fixtures for the Impact Slides test suite."""
from __future__ import annotations

import json
import shutil
import sys
import tempfile
from pathlib import Path

import pytest

# Make the repo root importable so `import step1_preprocessor_v2_full` works
# regardless of where pytest is invoked from.
REPO_ROOT = Path(__file__).resolve().parent.parent
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))


# --------------------------------------------------------------------------- #
# Temp directories
# --------------------------------------------------------------------------- #
@pytest.fixture()
def tmp_workspace():
    """
    A temp directory that cleans up gracefully on Windows (where openpyxl/pandas
    can briefly hold file handles, causing PermissionError on rmtree).

    Yields a Path; caller writes inputs into it.
    """
    d = Path(tempfile.mkdtemp(prefix="impact_test_"))
    try:
        yield d
    finally:
        shutil.rmtree(d, ignore_errors=True)


# --------------------------------------------------------------------------- #
# Preprocessor factory
# --------------------------------------------------------------------------- #
@pytest.fixture()
def make_preprocessor(tmp_workspace):
    """Factory: build a preprocessor pointing at a fresh input/output pair."""
    from step1_preprocessor_v2_full import ImpactSlidePreprocessorV2

    def _make(filter_level="moderate", boost_keywords=None, enable_ocr=False):
        inp = tmp_workspace / "input"
        out = tmp_workspace / "output"
        inp.mkdir(parents=True, exist_ok=True)
        p = ImpactSlidePreprocessorV2(
            input_path=str(inp),
            output_dir=str(out),
            filter_level=filter_level,
            boost_keywords=boost_keywords or [],
        )
        p.enable_ocr = enable_ocr
        return p, inp, out

    return _make


# --------------------------------------------------------------------------- #
# Synthetic file builders
# --------------------------------------------------------------------------- #
@pytest.fixture()
def make_excel(tmp_workspace):
    import pandas as pd

    def _make(name="sample.xlsx", df=None, sheets=None):
        path = tmp_workspace / "input" / name
        path.parent.mkdir(parents=True, exist_ok=True)
        if sheets:
            with pd.ExcelWriter(path) as w:
                for sn, frame in sheets.items():
                    frame.to_excel(w, sheet_name=sn, index=False)
        else:
            df.to_excel(path, index=False)
        return path

    return _make


@pytest.fixture()
def make_pptx(tmp_workspace):
    """Build a .pptx with configurable slides. Each slide spec is a dict."""
    def _make(name="deck.pptx", slides=None):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slides = slides or []
        for spec in slides:
            layout_idx = spec.get("layout", 6)  # 6 = blank
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            if spec.get("title"):
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                tb.text_frame.text = spec["title"]
            if spec.get("body"):
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
                tb.text_frame.text = spec["body"]
            # table
            if spec.get("table"):
                rows, cols = 2, 2
                tbl_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(2))
                data = spec["table"]
                for r in range(rows):
                    for c in range(cols):
                        tbl_shape.table.cell(r, c).text = str(data[r][c])
            # chart
            if spec.get("chart"):
                from pptx.chart.data import CategoryChartData
                from pptx.enum.chart import XL_CHART_TYPE
                cd = CategoryChartData()
                cd.categories = spec["chart"].get("categories", ["A", "B"])
                for series_name, values in spec["chart"].get("series", {}).items():
                    cd.add_series(series_name, values)
                ct = {
                    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE,
                }.get(spec["chart"].get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
                gframe = slide.shapes.add_chart(ct, Inches(1), Inches(2), Inches(6), Inches(3), cd)
                if spec["chart"].get("title"):
                    gframe.chart.has_title = True
                    gframe.chart.chart_title.text_frame.text = spec["chart"]["title"]
        path = tmp_workspace / "input" / name
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        return path

    return _make


@pytest.fixture()
def scanned_pdf(tmp_workspace):
    """
    A truly scanned (image-only) PDF: render text onto an image and embed it,
    so page.get_text() returns "". Falls back to the real sample if available.
    """
    path = tmp_workspace / "input" / "scanned.pdf"
    path.parent.mkdir(parents=True, exist_ok=True)
    try:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        # Draw text as vector-ish shapes won't fool get_text; instead insert an
        # image of text. We render a PIL image of text and place it full-page.
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new("RGB", (1240, 1754), "white")
        d = ImageDraw.Draw(img)
        d.text((60, 200), "StrategicERP Company Profile\nFounded 2005 in Mumbai\nERP CRM Automation",
               fill="black")
        import io
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        page.insert_image(page.rect, stream=buf.getvalue())
        doc.save(path)
        doc.close()
        return path
    except Exception:
        pytest.skip("Could not synthesize a scanned PDF")


@pytest.fixture()
def real_scanned_pdf():
    """A real scanned PDF from the user's Documents folder, if present."""
    candidates = [
        Path(r"C:\Users\Ag1Le\Documents\PublicWaterMassMailing.pdf"),
        Path(r"C:\Users\Ag1Le\Documents\Company-Profile-Strategicerp.pdf"),
    ]
    for c in candidates:
        if c.exists():
            return c
    pytest.skip("No real scanned PDF available on this machine")
