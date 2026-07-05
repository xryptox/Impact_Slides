"""CLI entry point + console helpers (v4 Phase 4 extraction).

Contains:
  - ``main(argv)``        — argparse entry point; resolves CLI > YAML > default
    config, dispatches to the preprocessor (or standalone --emit-schema).
  - ``test_preprocessor()`` — built-in smoke test (sample Excel + PPTX).
  - ``inspect_register()`` — readable top-N console summary for manual review.

Extracted from the monolith. Depends on the trunk class
(``ImpactSlidePreprocessorV4``), the config leaf, and the schemas leaf — all
imported from the package, so no module-level coupling to the old monolith.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import pandas as pd

from .preprocessor import ImpactSlidePreprocessorV4, _HAS_PYDANTIC, __version__
from .config import load_config, merge_config, validate_config
from .schemas import EvidenceEntry


def test_preprocessor():
    """Quick test covering both Excel + rich PPTX evidence extraction."""
    import tempfile
    import os

    print("\n=== Running Full Pipeline Test ===")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)

        # --- Sample Excel ---
        excel_file = tmp_path / "test_sales.xlsx"
        pd.DataFrame({
            "S.No": list(range(1, 21)),
            "Region": ["North", "South", "East", "West"] * 5,
            "Revenue": [1000 + i*50 for i in range(20)],
            "Status": ["Active", "Active", "Inactive", "Active"] * 5
        }).to_excel(excel_file, index=False)

        # --- Sample PPTX (minimal) ---
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            # Slide 1: Title
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            txBox.text_frame.paragraphs[0].text = "Q3 Impact Review"
            # Slide 2: Data-rich conclusion slide (should generate rich evidence)
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_box.text_frame.paragraphs[0].text = "Key Findings & Recommendations"

            content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(3))
            tf = content_box.text_frame
            tf.text = "• North region grew +23% YoY\n• Total revenue reached $4.2M (record high)\n• Recommendation: Expand North operations by 40%"
            prs.save(tmp_path / "test_deck.pptx")
        except Exception as e:
            print(f"Could not create sample PPTX: {e}")

        # Run preprocessor
        preprocessor = ImpactSlidePreprocessorV4(
            input_path=str(tmp_path),
            output_dir=str(tmp_path / "output"),
            filter_level="moderate"
        )
        preprocessor.run()

        # Verify outputs
        out_dir = tmp_path / "output"
        excel_p = out_dir / "excel_profile.json"
        pptx_p = out_dir / "pptx_profile.json"
        ev_reg = out_dir / "evidence_register_seed.json"

        print("\n=== Test Results ===")
        print(f"Excel profile exists: {excel_p.exists()}")
        print(f"PPTX profile exists: {pptx_p.exists()}")

        if ev_reg.exists():
            with open(ev_reg) as f:
                evidence = json.load(f)
            print(f"Total evidence entries: {len(evidence)}")
            pptx_ev = [e for e in evidence if "pptx" in e.get("insight_type", "") or e.get("insight_type", "").endswith("_insight")]
            print(f"PPTX-derived evidence: {len(pptx_ev)}")
            if pptx_ev:
                print("Sample PPTX evidence:", pptx_ev[0]["text"][:100])

        print("\nTest completed successfully!")


def inspect_register(output_dir: str, top_n: int = 15) -> None:
    """Print a readable top-N summary of the Evidence Register to the console.

    Useful for a quick manual quality review after running the preprocessor:
    shows the highest-priority insights, a breakdown by insight type and source
    file, and flags the known anti-patterns (noise at the top, cross-file
    flooding, missing framework stages).

    Args:
        output_dir: folder where the preprocessor wrote its outputs.
        top_n: how many top-priority evidence entries to print.
    """
    from collections import Counter
    out = Path(output_dir)
    reg_path = out / "evidence_register_seed.json"
    if not reg_path.exists():
        print(f"[inspect] No evidence_register_seed.json found in {out}")
        return

    ev = json.load(open(reg_path, encoding="utf-8"))
    if not ev:
        print("[inspect] Evidence register is empty.")
        return

    stages = {"Why", "What", "How", "Now"}
    print("\n" + "=" * 78)
    print(f" EVIDENCE REGISTER INSPECTION  —  {len(ev)} entries  (showing top {min(top_n, len(ev))})")
    print("=" * 78)

    # Breakdowns
    by_type = Counter(e.get("insight_type", "?") for e in ev)
    by_src = Counter(e.get("source_file", "?") for e in ev)
    avg_prio = sum(e.get("priority_score", 0) for e in ev) / len(ev)
    stages_used = set()
    for e in ev:
        stages_used.update(e.get("suggested_narrative_use", []))

    print(f" Average priority: {avg_prio:.3f}")
    print(f" Framework stages covered: {sorted(stages_used & stages)} "
          f"| MISSING: {sorted(stages - stages_used) or 'none'}")
    print(f" Sources: {len(by_src)}  | Insight types: {len(by_type)}")
    print()
    print(" By source file:")
    for s, c in by_src.most_common():
        print(f"   {c:>4}  {s}")
    print()
    print(" By insight type:")
    for t, c in by_type.most_common():
        print(f"   {c:>4}  {t}")

    # Top-N detail
    print()
    print("-" * 78)
    print(f" TOP {min(top_n, len(ev))} HIGHEST-PRIORITY EVIDENCE")
    print("-" * 78)
    for e in ev[:top_n]:
        eid = e.get("evidence_id", "?")
        prio = e.get("priority_score", 0)
        itype = e.get("insight_type", "?")
        src = e.get("source_file", "?")
        loc = e.get("source_location", "?")
        use = "/".join(e.get("suggested_narrative_use", []))
        text = e.get("text", "").replace("\n", " ")
        print(f"\n[{eid}] p={prio:.2f}  {itype}  ({src} @ {loc})")
        print(f"   stages: {use}   | {text[:120]}")

    # Quality flags
    print()
    print("-" * 78)
    print(" QUALITY FLAGS")
    print("-" * 78)
    flags = []

    noise_markers = ("222.127", "Mozilla", "HTTP/1.1", "GET /", "/cgi-bin/")
    top_noise = [e for e in ev[:5] if any(n in e.get("text", "") for n in noise_markers)]
    if top_noise:
        flags.append(f"  ! {len(top_noise)} noise-looking entries in the top 5 "
                     f"(IPs/URLs/user-agents should not rank this high)")

    numeric_cross = [e for e in ev if e.get("insight_type") == "cross_file_metric"
                     and "Numeric value" in e.get("text", "")]
    if len(numeric_cross) > 3:
        flags.append(f"  ! {len(numeric_cross)} numeric cross-file entries "
                     f"(likely false positives between unrelated files)")

    missing_stages = stages - stages_used
    if missing_stages:
        flags.append(f"  ! Framework stages with no evidence: {sorted(missing_stages)} "
                     f"(e.g. 'Now' needs conclusion/recommendation content)")

    low_conf = sum(1 for e in ev if e.get("confidence") == "medium" and e.get("priority_score", 0) > 0.8)
    if low_conf > len(ev) * 0.4:
        flags.append(f"  ! {low_conf} high-priority entries are only 'medium' confidence "
                     f"(check OCR quality / source reliability)")

    if not flags:
        print("  (no known anti-patterns detected)")
    else:
        for f in flags:
            print(f)

    print("\n" + "=" * 78)
    print(f" Full register: {reg_path}")
    print(f" Summary report: {out / 'preprocessor_summary.md'}")
    print("=" * 78 + "\n")


def main(argv=None):
    """CLI entry point. argv defaults to sys.argv[1:] so the function is
    callable from tests without touching sys.argv."""
    parser = argparse.ArgumentParser(description="Impact Slide Preprocessor v4")
    parser.add_argument("--input", required=False, help="Input folder path")
    parser.add_argument("--output", required=False, help="Output folder path")
    parser.add_argument("--filter-level", default="conservative", choices=["conservative", "moderate", "permissive"],
                        help="Filtering strictness level (default: conservative)")
    parser.add_argument("--boost-keywords", nargs="*", default=[],
                        help="Keywords to boost in evidence (Item 4.3). Example: --boost-keywords recommend critical")
    parser.add_argument("--verbose", action="store_true", default=False,
                        help="Enable detailed console logging (Item 4.5)")
    parser.add_argument("--export-md", action="store_true", default=False,
                        help="Export Evidence Register as Markdown (Item 4.6)")
    parser.add_argument("--export-csv", action="store_true", default=False,
                        help="Export Evidence Register as CSV (Item 4.6)")
    parser.add_argument("--enable-ocr", action="store_true", default=False,
                        help="Enable OCR for scanned PDFs (Phase 2)")
    parser.add_argument("--tesseract-cmd", default=None,
                        help="Path to the Tesseract OCR binary (auto-detected if omitted)")
    parser.add_argument("--pdf-table-engine", choices=["auto", "pdfplumber", "pymupdf"],
                        default="auto",
                        help="PDF table detection backend (default: auto = prefer pdfplumber, fall back to PyMuPDF)")
    parser.add_argument("--dedup-engine", choices=["auto", "embeddings", "tfidf", "fuzzy"],
                        default="auto",
                        help="Semantic near-dup dedup engine (default: auto = prefer sentence-transformers embeddings, fall back to rapidfuzz char-similarity). "
                             "embeddings/tfidf/fuzzy force a tier (graceful fallback if unavailable). "
                             "Note: tfidf is opt-in for prose-heavy registers; empirical testing showed it over-merges templated evidence.")
    parser.add_argument("--inspect", action="store_true", default=False,
                        help="Print a readable top-N Evidence Register summary to the console after running")
    parser.add_argument("--inspect-top", type=int, default=15,
                        help="Number of top-priority evidence entries to show with --inspect (default: 15)")
    parser.add_argument("--emit-schema", action="store_true", default=False,
                        help="Write evidence_schema.json (the Analyst GPT contract) to --output and exit, without processing files")
    parser.add_argument("--focus-areas", type=int, default=5,
                        help="v4 #26: number of ranked focus areas to surface in the Analyst Briefing (default: 5)")
    parser.add_argument("--config", default=None,
                        help="Path to a YAML config file. CLI flags override YAML; YAML overrides defaults. "
                             "(optional; requires PyYAML). Keys mirror the CLI flags in snake_case "
                             "(e.g. filter_level, dedup_engine, boost_keywords).")
    args = parser.parse_args(argv)

    # --- v3 #21: layered config resolution (CLI > YAML > default) ----------
    # Load the YAML config (if --config given), then merge with CLI args so
    # every downstream branch reads from a single resolved `cfg` dict. Errors
    # in the config (missing file, bad value, PyYAML absent) fail fast with a
    # clear message and a non-zero exit code.
    try:
        args.config_data = load_config(args.config)
        cfg = merge_config(parser, args)
        validate_config(cfg)
    except (FileNotFoundError, RuntimeError, ValueError) as e:
        print(f"Config error: {e}")
        return 1

    # Standalone schema emission: writes the JSON Schema for EvidenceEntry so it
    # can be embedded into the Analyst GPT prompt, without running the pipeline.
    if cfg["emit_schema"] and cfg["output"]:
        out_dir = Path(cfg["output"])
        out_dir.mkdir(parents=True, exist_ok=True)
        if not _HAS_PYDANTIC:
            print("pydantic is not installed; cannot emit schema. pip install pydantic")
        else:
            with open(out_dir / "evidence_schema.json", "w") as f:
                json.dump(EvidenceEntry.model_json_schema(), f, indent=2)
            print(f"Evidence schema written to {out_dir / 'evidence_schema.json'}")
    elif cfg["input"] and cfg["output"]:
        preprocessor = ImpactSlidePreprocessorV4(
            input_path=cfg["input"],
            output_dir=cfg["output"],
            filter_level=cfg["filter_level"],
            boost_keywords=cfg["boost_keywords"],
        )
        preprocessor.verbose = cfg["verbose"]
        preprocessor.export_md = cfg["export_md"]
        preprocessor.export_csv = cfg["export_csv"]
        preprocessor.enable_ocr = cfg["enable_ocr"]
        preprocessor.tesseract_cmd = cfg["tesseract_cmd"]
        preprocessor.pdf_table_engine = cfg["pdf_table_engine"]
        preprocessor.dedup_engine = cfg["dedup_engine"]
        preprocessor.focus_areas_count = cfg["focus_areas"]
        # v4 #26: apply optional briefing config from YAML (weights + business
        # keywords). CLI defaults are 5 / built-ins; YAML can override without a
        # dedicated flag. Validated above in validate_config().
        br = cfg.get("briefing") or {}
        preprocessor.briefing_readiness_weights = br.get("readiness_weights")
        preprocessor.briefing_focus_weights = br.get("focus_weights")
        preprocessor.briefing_business_keywords = br.get("business_keywords")
        # v3 #23: snapshot the resolved config (after all attribute sets) so
        # run_metadata.json records exactly which options produced this output.
        preprocessor.config_snapshot = dict(cfg)
        # v3 #24: (re)build stage rules now that config_snapshot is set, so any
        # user stage_rules overrides take effect before run() uses them.
        preprocessor.stage_rules = preprocessor._build_stage_rules()
        preprocessor.run()
        if cfg["inspect"]:
            inspect_register(cfg["output"], top_n=cfg["inspect_top"])
    else:
        print("No input/output provided. Running built-in test...")
        test_preprocessor()
    return 0


if __name__ == "__main__":
    sys.exit(main())
