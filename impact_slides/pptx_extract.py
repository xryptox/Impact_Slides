"""PPTX shape-helper functions (v3 #13–#16).

Group recursion, SmartArt/diagram XML text fallback, embedded-OLE detection,
and spatial (top→left) shape ordering so multi-column slides read in order.

Pure helpers: no `self`, no preprocessor state. The ``extract_pptx`` and
``classify_slide`` orchestrators stay on the trunk (they thread self.pptx_profiles
and self.stage_rules).
"""
from __future__ import annotations

# v3 #13-16: richer PPTX extraction helpers.
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn as _qn
    _A_T = _qn('a:t')   # drawingml text-run element
except ImportError:  # pragma: no cover - python-pptx optional
    _MSO_SHAPE_TYPE = None
    _qn = None
    _A_T = '{http://schemas.openxmlformats.org/drawingml/2006/main}t'


def _is_group(shape) -> bool:
    """True if the shape is a group (has child shapes to recurse into)."""
    try:
        return int(shape.shape_type) == 6  # MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _is_embedded_object(shape) -> bool:
    """True for embedded/linked OLE objects (embedded Excel sheets, PDFs, …).
    These carry data the preprocessor can't read directly, so the Analyst is
    at least told they exist (see embedded_objects in slide details)."""
    try:
        st = int(shape.shape_type)
        return st in (7, 10)  # EMBEDDED_OLE_OBJECT, LINKED_OLE_OBJECT
    except Exception:
        return False


def _shape_position(shape):
    """Return (top, left) for spatial sorting, or (0, 0) if unavailable."""
    try:
        return (int(shape.top or 0), int(shape.left or 0))
    except Exception:
        return (0, 0)


def _iter_shapes_deep(shapes):
    """Yield shapes, recursing into groups so nested text boxes are not lost.
    Top-level shapes (and group children) are yielded in spatial (top, left)
    order so multi-column slides read top-to-bottom, left-to-right."""
    try:
        ordered = sorted(shapes, key=_shape_position)
    except Exception:
        ordered = list(shapes)
    for shape in ordered:
        if _is_group(shape):
            try:
                yield from _iter_shapes_deep(shape.shapes)
            except Exception:
                pass
        else:
            yield shape


def _extract_shape_text(shape) -> str:
    """Extract text from a shape, with a richer fallback.

    First tries the normal text_frame.text. If that's empty/absent (as for
    SmartArt/diagram graphic frames, whose text lives in nested drawingml XML,
    not a text_frame), fall back to collecting every <a:t> text-run under the
    shape's XML element. This recovers SmartArt node labels and diagram text
    that the simple path silently drops."""
    try:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if txt and txt.strip():
                return txt.strip()
    except Exception:
        pass
    # Fallback: pull all drawingml text runs from the shape's XML. Catches
    # SmartArt (<dgm:>) and diagram graphic frames where has_text_frame=False.
    try:
        el = shape._element
        runs = el.findall('.//' + _A_T)
        parts = [r.text for r in runs if r.text and r.text.strip()]
        if parts:
            return ' '.join(p.strip() for p in parts)
    except Exception:
        pass
    return ""
