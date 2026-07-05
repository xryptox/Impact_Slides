#!/usr/bin/env python3
"""
Impact Slide Preprocessor v4
=============================
v4 builds on v3 with everything v3 added (twelve insight-quality
enhancements, Pydantic schema contract, richer PPTX extraction, merged
pdfplumber/PyMuPDF table detection, fuzzy/abbreviation cross-file entity
matching, tiered semantic dedup, optional YAML config, always-on time
profiling, centralized logging with run_metadata.json, configurable
Why/What/How/Now stage mapping, and IQR outlier/correlation/period-trend
analytics) PLUS the new Analyst Briefing Generator (v4 #26):
  - Narrative Readiness Score (0-100 overall + per stage), a 5-component
    weighted composite (coverage balance, priority quality, cross-file
    connectivity, recommendation strength, signal ratio).
  - Ranked Suggested Focus Areas with sophisticated multi-factor scoring
    (avg priority + cross-file strength + insight-quality boost + source
    diversity + business-relevance signals) over multi-signal theme
    detection (column names, "X by Y" patterns, cross-file entities, business
    keywords, derived-insight boosts, near-duplicate theme merging).
  - Top cross-file relationships surfaced clearly.
  - Narrative gap analysis + slide-building recommendations.
  - Emits `analyst_briefing.md` + `analyst_briefing.json` unconditionally for
    a tighter, higher-signal handoff to the Impact Slide Analyst GPT (Step 2).

v3 also added five insight-quality enhancements toward the goal of
seeding a richer, priority-ordered Evidence Register for the Impact Slide
Analyst GPT:
  1. Trend / delta insights across time-ordered spreadsheet sheets
  2. Cross-sheet consolidation of repeated numeric ranges
  3. Per-bullet insight-language ranking (no more flat-0.75 bullet scores)
  4. Coverage map (Why/What/How/Now + per-source) emitted as a handoff file
  5. Computed group-by aggregate insights (not just suggestions)

Usage:
    python step1_preprocessor_v3.py --input /path/to/source_files --output ./output
"""

import argparse
import json
import logging
import math
import os
import platform
import re
import subprocess
import sys
import time
import statistics
import warnings
from collections import defaultdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Iterable

import pandas as pd

# v3 #23: explicit version constant (previously implicit in docstrings).
# Single source of truth consumed by run_metadata.json + the logger startup
# banner + a future --version flag.
__version__ = "4.0.0"

# v3: Pydantic schemas are the single source of truth for the output contracts.
# Optional at runtime: if pydantic isn't installed, validation is skipped (the
# preprocessor still runs), but emitting the JSON schema requires it.
# v4 refactor: the models now live in impact_slides.schemas (the package leaf).
try:
    from impact_slides.schemas import (EvidenceEntry, FileInventoryItem, CoverageMap,
                         EntitiesSummaryItem, AnalystBriefing, NarrativeReadiness,
                         FocusArea, StageScore)
    _HAS_PYDANTIC = True
except ImportError:
    EvidenceEntry = FileInventoryItem = CoverageMap = EntitiesSummaryItem = None
    AnalystBriefing = NarrativeReadiness = FocusArea = StageScore = None
    _HAS_PYDANTIC = False

# v4 #26: Analyst Briefing Generator. Optional at runtime — if the module
# somehow cannot be imported, the briefing artefacts are simply skipped (the
# pipeline still runs and every other output is produced).
try:
    from impact_slides.analyst_briefing import (
        AnalystBriefingGenerator, READINESS_WEIGHTS, FOCUS_WEIGHTS,
        DEFAULT_BUSINESS_KEYWORDS, render_briefing_markdown,
    )
    _HAS_BRIEFING = True
except ImportError:  # pragma: no cover - graceful degradation
    AnalystBriefingGenerator = None
    READINESS_WEIGHTS = {}
    FOCUS_WEIGHTS = {}
    DEFAULT_BUSINESS_KEYWORDS = set()
    render_briefing_markdown = None
    _HAS_BRIEFING = False

# ====================== OPTIONAL IMPORTS ======================

# PDF support (3.1)
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# v3: pdfplumber for superior PDF table detection (optional; graceful
# degradation to PyMuPDF's find_tables() when absent). PyMuPDF stays the
# primary engine for text/layout/OCR rendering; pdfplumber is used only for
# table cell detection where it handles merged/spanning/unruled tables better.
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

# DOCX support (3.3)
try:
    from docx import Document
except ImportError:
    Document = None

# OCR support (Phase 1)
try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

# v3: semantic dedup clustering. numpy powers the pure-Python TF-IDF + cosine
# tier (no sklearn dependency); sentence-transformers is an optional heavier
# tier for true embedding-based similarity. Both are optional — the dedup pass
# falls back to rapidfuzz char-similarity when neither is available.
try:
    import numpy as np
except ImportError:
    np = None

# v3 #21: optional YAML config file support. PyYAML is optional; when absent,
# --config is rejected with a clear message and the preprocessor runs CLI-only
# (identical to the existing behavior). Matches the codebase's graceful-
# degradation pattern for optional dependencies.
try:
    import yaml
    _HAS_YAML = True
except ImportError:
    yaml = None
    _HAS_YAML = False

# v3 #23: centralized logging. structlog gives leveled, context-rich,
# JSON-capable logs; stdlib logging is the always-available fallback so the
# preprocessor works with zero optional deps (same graceful-degradation
# pattern as pdfplumber/pydantic/rapidfuzz/numpy/yaml).
try:
    import structlog
    _HAS_STRUCTLOG = True
except ImportError:
    structlog = None
    _HAS_STRUCTLOG = False


# ====================== VERSION + GIT PROVENANCE (v3 #23) =====================
# Read-only helpers that stamp "which code produced this output?" into
# run_metadata.json for reproducibility. These NEVER commit, stage, push, or
# write anything to git — they only read .git/refs to get the current commit
# hash and whether the working tree has uncommitted edits (dirty).
#
# v4 refactor: the implementations now live in the modular `impact_slides`
# package (logging_setup.py, text_utils.py, heuristics.py, text_analysis.py,
# config.py). The monolith re-exports them under their original module-level
# names so every existing call site + the 430-test suite keep working
# unchanged. The original monolith bodies were relocated verbatim.
from impact_slides.logging_setup import (
    git_commit as _pkg_git_commit, git_dirty as _pkg_git_dirty,
    get_logger as _pkg_get_logger, _StdlibLogAdapter,
    _structlog_console_renderer as _pkg_structlog_console_renderer,
)
from impact_slides.text_utils import (
    clean_text, get_column_letter, excel_addr, safe_stat, compact_value,
    confidence_for_method, _METHOD_RELIABILITY,
)
from impact_slides.heuristics import (
    sheet_time_rank, _parse_numeric_token, make_unique_columns,
    is_likely_identifier_column, is_generic_system_column, _looks_like_noise_cell,
    _NOISE_CELL_PATTERNS, _MONTH_RE, _QUARTER_RE, _YEAR_RE,
)
from impact_slides.text_analysis import (
    extract_advanced_metrics, contains_insight_language, insight_priority_boost,
    calculate_evidence_priority_score, _INSIGHT_KEYWORDS,
)
from impact_slides.stage_mapping import (
    DEFAULT_INSIGHT_TYPE_STAGES, DEFAULT_KEYWORD_STAGE_OVERRIDES,
    DEFAULT_SLIDE_TYPE_KEYWORDS, DEFAULT_SLIDE_TYPE_STAGES,
    DEFAULT_CONCLUSION_BULLET_STAGES,
)
from impact_slides.config import (
    CONFIG_DEFAULTS, CONFIG_CHOICES, load_config, merge_config, validate_config,
    _cli_was_set, _HAS_YAML,
)
# (schemas + analyst_briefing names already imported at the top of this file
# under the _HAS_PYDANTIC / _HAS_BRIEFING guards.)


def git_commit() -> Optional[str]:
    return _pkg_git_commit()


def git_dirty() -> Optional[bool]:
    return _pkg_git_dirty()


# ====================== LOGGER FACTORY (v3 #23) =====================
# Centralized leveled logging replacing the ad-hoc print() pattern. structlog
# is preferred (structured key-value context, JSON-rendered file logs);
# stdlib logging is the always-available fallback. The verbose flag lowers
# the console level to DEBUG instead of gating scattered `if self.verbose`
# blocks.

_LOG = None  # singleton logger for the preprocessor instance


def get_logger(name: str = "preprocessor", log_file: Optional[Path] = None,
               verbose: bool = False, run_id: Optional[str] = None):
    return _pkg_get_logger(name=name, log_file=log_file, verbose=verbose,
                          run_id=run_id, version=__version__)


class _StdlibLogAdapter:
    """Thin adapter over stdlib logging so call sites can use the structlog-style
    kwarg API (log.info("msg", key=value, ...)) even when structlog isn't
    installed. Extra kwargs are folded into the message as key=value pairs."""

    def __init__(self, stdlib_logger, base_ctx: dict):
        self._log = stdlib_logger
        self._ctx = dict(base_ctx)
        self._file_logger = None

    def bind(self, **kwargs):
        new = _StdlibLogAdapter(self._log, {**self._ctx, **kwargs})
        new._file_logger = self._file_logger
        return new

    def _emit(self, level, event, kwargs):
        merged = {**self._ctx, **kwargs}
        extra = " ".join(f"{k}={v!r}" for k, v in merged.items())
        msg = f"{event}" + (f" {extra}" if extra else "")
        self._log.log(level, msg)
        if self._file_logger is not None:
            self._file_logger.log(level, msg)

    def debug(self, event, **kw):   self._emit(logging.DEBUG, event, kw)
    def info(self, event, **kw):    self._emit(logging.INFO, event, kw)
    def warning(self, event, **kw): self._emit(logging.WARNING, event, kw)
    def error(self, event, **kw):   self._emit(logging.ERROR, event, kw)
    def exception(self, event, **kw): self._emit(logging.ERROR, event, kw)


def _structlog_console_renderer():
    return _pkg_structlog_console_renderer()

# ====================== HELPER FUNCTIONS ======================
# (Re-exported from the modular package leaves; kept as module-level names
# so the monolith's existing call sites and tests keep working unchanged.)

# v4 refactor: the dedup engine (v3 #20) + cross-file entity helpers
# (v3 #17-#19) now live in impact_slides.dedup / impact_slides.cross_file.
# Imported under their original module-level names.
from impact_slides.dedup import (
    _text_similarity, _load_sentence_model, _tfidf_vectors, _SemanticDedupEngine,
    _SENTENCE_MODEL, _TFIDF_TOKEN_RE, _SEM_THRESHOLD, np,
)
from impact_slides.cross_file import (
    _ABBREVIATIONS, _entity_aliases, _entity_in_text,
)

# v3 #13-16: richer PPTX extraction helpers.
# Group recursion, SmartArt/diagram XML text fallback, embedded-OLE detection,
# and spatial (top→left) shape ordering so multi-column slides read in order.
# v4 refactor: PPTX shape helpers (v3 #13-#16) moved to
# impact_slides.pptx_extract. Imported under their original names.
from impact_slides.pptx_extract import (
    _is_group, _is_embedded_object, _shape_position, _iter_shapes_deep,
    _extract_shape_text, _MSO_SHAPE_TYPE, _qn, _A_T,
)

# ====================== MAIN PREPROCESSOR CLASS ======================

class ImpactSlidePreprocessorV4:
    def __init__(self, input_path: str, output_dir: str, max_sheet_rows: int = 5000, filter_level: str = "conservative", boost_keywords: List[str] = None):
        """
        filter_level options:
            - "conservative": Strict filtering (recommended for clean Evidence Register)
            - "moderate"    : Balanced filtering
            - "permissive"  : Minimal filtering (keep almost everything)
        """
        self.input_path = Path(input_path)
        self.output_dir = Path(output_dir)
        self.max_sheet_rows = max_sheet_rows
        self.filter_level = filter_level.lower()
        self.inventory = []
        self.excel_profiles = []
        self.pptx_profiles = []
        self.pdf_profiles = []
        self.docx_profiles = []
        self.errors = []
        self.filtered_items = []   # Log of what was filtered and why
        self.boost_keywords = [k.lower().strip() for k in (boost_keywords or [])]   # Item 4.3
        self.verbose = False   # Item 4.5
        self.export_md = False # Item 4.6
        self.export_csv = False # Item 4.6
        self.enable_ocr = False # Phase 2: OCR support
        self.tesseract_cmd = None  # Path to tesseract binary; auto-detected if None
        self._ocr_available = None  # Cache for _ensure_tesseract() result
        self.pdf_table_engine = "auto"  # v3: auto | pdfplumber | pymupdf
        self.dedup_engine = "auto"      # v3: auto | embeddings | tfidf | fuzzy
        # v4 #26: Analyst Briefing settings. focus_areas_count = how many
        # ranked focus areas to surface in the briefing (default 5; override
        # via --focus-areas / YAML). briefing_weights / business_keywords come
        # from YAML (Phase 2) and are None by default → scorer uses built-ins.
        self.focus_areas_count = 5
        self.briefing_readiness_weights = None
        self.briefing_focus_weights = None
        self.briefing_business_keywords = None
        # Uniform ceiling on the `text` field of every evidence entry. The
        # canonical constant lives in schemas.MAX_TEXT_LENGTH; the preprocessor
        # truncates to this length (or a user-lowered value) in _validate_
        # evidence() so the register stays compact and the Analyst GPT token
        # budget is predictable. Set from cfg in cli.main(); the default is
        # the schema ceiling.
        from .schemas import MAX_TEXT_LENGTH as _SCHEMA_MAX_TEXT
        self.max_text_length = _SCHEMA_MAX_TEXT
        # Generated artefacts (set by _generate_analyst_briefing()).
        self.analyst_briefing_md = None
        self.analyst_briefing_json = None
        # v3 #22: structured timing data. Single source of truth that feeds
        # both the always-on console timing block and the preprocessor_summary.md
        # "Processing Time" section, so they never drift apart.
        self.timing = {
            "files": [],          # [{file, category, duration_s, status}, ...]
            "stages": {},         # {discovery, extraction, evidence_build, output}
            "total_seconds": 0.0, # wall-clock end-to-end
        }
        # v3 #23: run identity for run_metadata.json + the logger. Captured
        # at construction so __init__ time is the run's "started_at".
        self.run_id = time.strftime("%Y%m%dT%H%M%S")
        self.run_started_at = time.strftime("%Y-%m-%dT%H:%M:%S")
        self.run_finished_at = None
        self.config_snapshot = None  # set by main() with the resolved cfg dict
        # v3 #23: centralized logger. Lazy-initialized in run() so the file
        # path (output_dir) is known; set to None here so attribute access
        # before run() doesn't crash.
        self.log = None
        # v3 #24: configurable Why/What/How/Now stage mapping. Built from the
        # DEFAULT_* tables + any user overrides in config_snapshot["stage_rules"].
        # Single source of truth for stage assignment — replaces ~20 scattered
        # hardcoded suggested_narrative_use literals.
        self.stage_rules = self._build_stage_rules()

    def _get_filter_thresholds(self):
        """Return filtering thresholds based on configured filter_level."""
        if self.filter_level == "permissive":
            return {
                "min_non_null_ratio": 0.02,
                "max_unique_ratio": 0.98,
                "min_priority": 0.05,
            }
        elif self.filter_level == "moderate":
            return {
                "min_non_null_ratio": 0.05,
                "max_unique_ratio": 0.92,
                "min_priority": 0.15,
            }
        else:  # conservative (default)
            return {
                "min_non_null_ratio": 0.10,
                "max_unique_ratio": 0.90,
                "min_priority": 0.25,
            }

    # ====================== STAGE RULES (v3 #24) =====================

    def _build_stage_rules(self) -> Dict[str, Any]:
        """Build the stage-rules table from DEFAULT_* tables + user overrides.

        User overrides come from self.config_snapshot["stage_rules"] (set by
        main() from the resolved config). Overrides EXTEND/REPLACE defaults —
        the user doesn't have to redefine the whole table. Keyword-overrides are
        pre-compiled into regex objects at build time so the per-evidence match
        loop stays fast (and bad patterns fail fast at startup, not mid-run).
        """
        cfg = getattr(self, "config_snapshot", None) or {}
        user = cfg.get("stage_rules", {}) or {}

        rules = {
            "insight_types":    dict(DEFAULT_INSIGHT_TYPE_STAGES),
            "keyword_overrides": list(DEFAULT_KEYWORD_STAGE_OVERRIDES),
            "slide_type_keywords": {k: list(v) for k, v in DEFAULT_SLIDE_TYPE_KEYWORDS.items()},
            "slide_type_stages": dict(DEFAULT_SLIDE_TYPE_STAGES),
            "conclusion_bullet_stages": list(DEFAULT_CONCLUSION_BULLET_STAGES),
        }
        # Merge user overrides (replace per-key, extend keyword lists).
        for k, v in (user.get("insight_types") or {}).items():
            rules["insight_types"][k] = list(v)
        for entry in (user.get("keyword_overrides") or []):
            rules["keyword_overrides"].append(entry)
        for k, v in (user.get("slide_type_keywords") or {}).items():
            # user can extend (list) or replace the keyword list for a type
            rules["slide_type_keywords"][k] = list(v)
        for k, v in (user.get("slide_type_stages") or {}).items():
            rules["slide_type_stages"][k] = list(v)
        if user.get("conclusion_bullet_stages"):
            rules["conclusion_bullet_stages"] = list(user["conclusion_bullet_stages"])

        # Pre-compile keyword-override regexes (fail fast on bad patterns).
        compiled = []
        for entry in rules["keyword_overrides"]:
            pat = entry["pattern"] if isinstance(entry, dict) else entry[0]
            stages = entry["stages"] if isinstance(entry, dict) else entry[1]
            try:
                compiled.append((re.compile(pat, re.IGNORECASE), list(stages)))
            except re.error as e:
                raise ValueError(
                    f"stage_rules.keyword_overrides: invalid regex {pat!r}: {e}"
                )
        rules["_compiled_keyword_overrides"] = compiled

        # Validate stages against NARRATIVE_STAGES (if pydantic/schemas loaded).
        self._validate_stage_rules(rules)
        return rules

    def _validate_stage_rules(self, rules: Dict[str, Any]) -> None:
        """Validate that every stage in the rules is a valid NARRATIVE_STAGES
        member. Raises ValueError at build time so a bad config fails fast."""
        try:
            from schemas import NARRATIVE_STAGES
        except ImportError:
            return  # can't validate without schemas; skip gracefully
        for section in ("insight_types", "slide_type_stages"):
            for k, stages in rules.get(section, {}).items():
                for s in stages:
                    if s not in NARRATIVE_STAGES:
                        raise ValueError(
                            f"stage_rules.{section}.{k}: invalid stage {s!r}; "
                            f"must be one of {sorted(NARRATIVE_STAGES)}"
                        )
        for s in rules.get("conclusion_bullet_stages", []):
            if s not in NARRATIVE_STAGES:
                raise ValueError(
                    f"stage_rules.conclusion_bullet_stages: invalid stage {s!r}; "
                    f"must be one of {sorted(NARRATIVE_STAGES)}"
                )
        for _, stages in rules.get("_compiled_keyword_overrides", []):
            for s in stages:
                if s not in NARRATIVE_STAGES:
                    raise ValueError(
                        f"stage_rules.keyword_overrides: invalid stage {s!r}; "
                        f"must be one of {sorted(NARRATIVE_STAGES)}"
                    )

    def _stages_for(self, insight_type: str, text: str = "",
                    default: Optional[List[str]] = None) -> List[str]:
        """Return the suggested_narrative_use stages for an insight.

        Lookup order: (1) keyword-override (first regex match in text, if any),
        (2) insight_types table, (3) the ``default`` arg, (4) ["What"].

        This replaces every hardcoded `suggested_narrative_use` literal and is
        the single entry point for stage assignment, so user config overrides
        flow through automatically.
        """
        # Layer 2: keyword overrides (first match wins).
        if text:
            for pat, stages in self.stage_rules.get("_compiled_keyword_overrides", []):
                if pat.search(text):
                    return list(stages)
        # Layer 1: insight-type default.
        stages = self.stage_rules.get("insight_types", {}).get(insight_type)
        if stages:
            return list(stages)
        # Fallback.
        return list(default) if default else ["What"]

    def gather_files(self) -> List[Path]:
        if not self.input_path.exists():
            raise FileNotFoundError(f"Input path not found: {self.input_path}")
        # v4 fix: accept a single file as input, not only a directory.
        if self.input_path.is_file():
            files = [self.input_path] if not self.input_path.name.startswith('.') else []
        else:
            files = [p for p in self.input_path.rglob("*") if p.is_file() and not p.name.startswith('.')]
        return files

    def build_file_inventory(self, files: List[Path]) -> List[Dict]:
        inventory = []
        for idx, file_path in enumerate(files, 1):
            try:
                ext = file_path.suffix.lower()
                if ext in [".xlsx", ".xls", ".csv", ".xlsm"]:
                    category = "spreadsheet"
                elif ext in [".pdf"]:
                    category = "pdf"
                elif ext in [".docx", ".doc"]:
                    category = "docx"
                else:
                    category = "other"

                access_status = "readable"
                try:
                    with open(file_path, "rb") as f:
                        f.read(1)
                except Exception:
                    access_status = "not_readable"

                inventory.append({
                    "file_id": f"F{idx:04d}",
                    "file_name": file_path.name,
                    "absolute_path": str(file_path.resolve()),
                    "category": category,
                    "access_status": access_status,
                })
            except Exception as e:
                self.errors.append({"file": str(file_path), "error": str(e)})
        return inventory

    # ====================== IMPROVED EXCEL PROFILING ======================

    def profile_dataframe(self, raw: Any, sheet_name: str, source_type: str) -> Dict[str, Any]:
        if raw is None or raw.empty:
            return {"sheet_name": sheet_name, "status": "empty"}

        raw_rows, raw_cols = raw.shape
        header_row = self._detect_header_row(raw)
        truncated = raw_rows >= self.max_sheet_rows

        if header_row is not None and header_row < raw_rows - 1:
            columns = make_unique_columns(raw.iloc[header_row].tolist())
            data = raw.iloc[header_row + 1:].copy()
            data.columns = columns
        else:
            columns = [f"Column {get_column_letter(i + 1)}" for i in range(raw_cols)]
            data = raw.copy()
            data.columns = columns

        numeric_profiles = []
        date_profiles = []
        categorical_profiles = []
        findings = []
        column_priority_scores = {}

        for col_idx, col_name in enumerate(columns):
            series = data.iloc[:, col_idx]
            non_null = series.dropna()
            non_null_count = int(non_null.shape[0])
            if non_null_count == 0:
                continue

            non_null_ratio = non_null_count / max(len(series), 1)
            has_business_name = not str(col_name).lower().startswith("column ")

            # ====================== CONFIGURABLE FILTERING ======================
            thresholds = self._get_filter_thresholds()

            # Skip columns with too many missing values
            if non_null_ratio < thresholds["min_non_null_ratio"]:
                self.filtered_items.append({
                    "reason": "high_missing_values",
                    "column": col_name,
                    "non_null_ratio": round(non_null_ratio, 3),
                    "threshold": thresholds["min_non_null_ratio"]
                })
                continue

            # Skip obvious generic/system columns
            if is_generic_system_column(col_name):
                self.filtered_items.append({
                    "reason": "generic_system_column",
                    "column": col_name
                })
                continue

            # --- Identifier Detection ---
            is_identifier = is_likely_identifier_column(col_name, series)
            if is_identifier:
                column_priority_scores[col_name] = 0.15
                numeric = pd.to_numeric(series.astype(str).str.replace(",", "", regex=False), errors="coerce")
                if numeric.notna().sum() >= 2:
                    numeric_profiles.append({
                        "column": col_name,
                        "is_identifier": True,
                        "note": "Likely ID/Serial column - limited business value"
                    })
                continue

            # --- Numeric ---
            numeric = pd.to_numeric(series.astype(str).str.replace(",", "", regex=False), errors="coerce")
            numeric_count = int(numeric.notna().sum())
            numeric_ratio = numeric_count / max(non_null_count, 1)

            if numeric_count >= 2 and numeric_ratio >= 0.6:
                desc_values = [float(v) for v in numeric.dropna().tolist()]
                max_idx = numeric.idxmax()
                min_idx = numeric.idxmin()
                profile = {
                    "column": col_name,
                    "min": safe_stat(desc_values, min),
                    "max": safe_stat(desc_values, max),
                    "mean": safe_stat(desc_values, statistics.mean),
                    "median": safe_stat(desc_values, statistics.median),
                }
                numeric_profiles.append(profile)

                priority = calculate_evidence_priority_score(
                    col_name, "numeric", non_null_ratio=non_null_ratio, has_business_name=has_business_name
                )
                column_priority_scores[col_name] = priority

                findings.append({
                    "type": "numeric_range",
                    "text": f"{sheet_name}: '{col_name}' ranges from {profile['min']} to {profile['max']}.",
                    "priority_score": priority,
                    "location": sheet_name,
                    "column": col_name,
                })

                # v3 #25: IQR outlier detection per numeric column. Values
                # below Q1 - 1.5*IQR or above Q3 + 1.5*IQR are statistical
                # outliers — extremely high-signal for the 'What' stage
                # (anomalies worth investigating) and 'How' (root-cause analysis).
                try:
                    q1 = float(numeric.quantile(0.25))
                    q3 = float(numeric.quantile(0.75))
                    iqr = q3 - q1
                    if iqr > 0:
                        lower_bound = q1 - 1.5 * iqr
                        upper_bound = q3 + 1.5 * iqr
                        outliers = numeric[(numeric < lower_bound) | (numeric > upper_bound)]
                        outlier_count = int(outliers.shape[0])
                        if outlier_count > 0:
                            outlier_pct = round(outlier_count / numeric_count * 100, 1)
                            sample_vals = [clean_text(v) for v in outliers.head(3).tolist()]
                            direction = "high" if outliers.mean() > profile.get("mean", 0) else "low"
                            findings.append({
                                "type": "outlier_insight",
                                "text": (f"{sheet_name}: '{col_name}' has {outlier_count} outlier(s) "
                                         f"({outlier_pct}% of data) outside [{round(lower_bound, 2)}, "
                                         f"{round(upper_bound, 2)}] (IQR method). Examples: {', '.join(sample_vals)}."),
                                "priority_score": 0.82,
                                "location": sheet_name,
                                "column": col_name,
                            })
                            profile["outlier_count"] = outlier_count
                            profile["outlier_bounds"] = [round(lower_bound, 2), round(upper_bound, 2)]
                except Exception:
                    pass
                continue

            # --- Date ---
            try:
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore", UserWarning)
                    dates = pd.to_datetime(series, errors="coerce")
                if dates.notna().sum() >= 2:
                    date_profiles.append({"column": col_name})
                    column_priority_scores[col_name] = calculate_evidence_priority_score(
                        col_name, "date", non_null_ratio=non_null_ratio, has_business_name=has_business_name
                    )
                    continue
            except Exception:
                pass

            # --- Improved Categorical (with conservative high-cardinality filter) ---
            unique_count = int(non_null.astype(str).nunique())
            unique_ratio = unique_count / max(non_null_count, 1)

            # Filter very high cardinality text columns (likely free text / IDs)
            if unique_ratio > thresholds["max_unique_ratio"]:
                self.filtered_items.append({
                    "reason": "high_cardinality_text",
                    "column": col_name,
                    "unique_ratio": round(unique_ratio, 3),
                    "threshold": thresholds["max_unique_ratio"]
                })
                continue

            if unique_count >= 2 and unique_ratio <= 0.85:
                top_values = non_null.astype(str).value_counts().head(6)
                total = non_null_count
                # Percentage covered by the top (up to 3) values. Previously this
                # returned 0.0% for columns with fewer than 3 unique values
                # (e.g. a 2-value Gender column reported '0.0%'). Now sums the
                # actually-available top values.
                top_n = min(3, len(top_values))
                top_3_pct = sum(v for v in top_values.values[:top_n]) / total * 100 if total else 0.0
                top_value_names = [str(k) for k in top_values.index[:5]]

                categorical_profiles.append({
                    "column": col_name,
                    "unique_count": unique_count,
                    "top_values": [{"value": str(k), "count": int(v)} for k, v in top_values.items()]
                })

                priority = calculate_evidence_priority_score(
                    col_name, "categorical", unique_ratio=unique_ratio, non_null_ratio=non_null_ratio, has_business_name=has_business_name
                )
                column_priority_scores[col_name] = priority

                # Include the actual top value names so the evidence is richer for
                # the Analyst AND so cross-file relationship detection can match
                # entities (e.g. "North") that appear in both data and narrative.
                top_names_str = ", ".join(top_value_names)
                findings.append({
                    "type": "categorical_distribution",
                    "text": f"{sheet_name}: '{col_name}' has {unique_count} unique values ({top_names_str}). Top 3 account for {top_3_pct:.1f}% of records.",
                    "priority_score": priority,
                    "location": sheet_name,
                    "column": col_name,
                })

        # v3 #25: correlation hints between numeric columns. Compute Pearson r
        # for each pair of non-ID numeric columns; emit a correlation_insight
        # when |r| >= 0.6 (moderate+) so the analyst sees which metrics move
        # together. Extremely high-signal for the 'How' stage (driver analysis).
        # Cap at top 8 pairs by |r| to keep the register compact.
        correlation_findings = []
        numeric_non_id = [p for p in numeric_profiles if not p.get("is_identifier")]
        if len(numeric_non_id) >= 2:
            corr_pairs = []
            for i, p1 in enumerate(numeric_non_id):
                for p2 in numeric_non_id[i + 1:]:
                    c1, c2 = p1["column"], p2["column"]
                    if c1 not in data.columns or c2 not in data.columns:
                        continue
                    try:
                        s1 = pd.to_numeric(data[c1], errors="coerce")
                        s2 = pd.to_numeric(data[c2], errors="coerce")
                        paired = pd.concat([s1, s2], axis=1).dropna()
                        if len(paired) < 5:
                            continue
                        r = float(paired.iloc[:, 0].corr(paired.iloc[:, 1]))
                        if not pd.isna(r) and abs(r) >= 0.6:
                            corr_pairs.append((abs(r), r, c1, c2, len(paired)))
                    except Exception:
                        pass
            corr_pairs.sort(key=lambda x: x[0], reverse=True)
            for _, r, c1, c2, n in corr_pairs[:8]:
                direction = "positive" if r > 0 else "negative"
                strength = "strong" if abs(r) >= 0.8 else "moderate"
                correlation_findings.append({
                    "type": "correlation_insight",
                    "text": (f"{sheet_name}: '{c1}' and '{c2}' have a {strength} {direction} "
                             f"correlation (r={r:.2f}, n={n}). They move {'together' if r > 0 else 'in opposite directions'}."),
                    "priority_score": round(0.80 + abs(r) * 0.10, 3),
                    "location": sheet_name,
                    "columns": [c1, c2],
                })
        findings.extend(correlation_findings)

        # Basic multi-column insight
        multi_column_insights = []
        if numeric_profiles and categorical_profiles:
            multi_column_insights.append({
                "type": "category_by_metric_suggestion",
                "text": f"Consider analyzing numeric metrics grouped by categorical columns (e.g. {categorical_profiles[0]['column']} vs {numeric_profiles[0]['column']})."
            })

        # v3 #5: compute concrete group-by aggregates (numeric metric grouped by a
        # low-cardinality categorical column). Previously only a *suggestion* was
        # emitted; now we seed actual source-backed aggregate insights (e.g.
        # "Total by Branch: A=..., B=..."), which is exactly the segmented view an
        # analyst wants. Cap pairs and categories to keep the register compact.
        aggregate_findings = []
        numeric_cols = [p["column"] for p in numeric_profiles]
        low_card_cats = [c for c in categorical_profiles if c.get("unique_count", 99) <= 8]
        for cat in low_card_cats[:2]:
            cat_col = cat["column"]
            for num_col in numeric_cols[:3]:
                try:
                    grouped = data.groupby(cat_col)[num_col].agg(["sum", "mean"])
                    if grouped.empty or len(grouped) > 10:
                        continue
                    parts = []
                    for cat_val, row in grouped.iterrows():
                        total = safe_stat([row["sum"]], lambda v: v[0])
                        parts.append(f"{cat_val}={total}")
                    agg_text = f"{sheet_name}: '{num_col}' by '{cat_col}' — " + ", ".join(parts)
                    aggregate_findings.append({
                        "type": "aggregate_insight",
                        "text": agg_text,
                        "priority_score": 0.83,
                        "location": sheet_name,
                        "column": num_col,
                        "group_by": cat_col,
                    })
                except Exception:
                    pass

        # v3 #25: within-sheet YoY/QoQ/MoM period trends. If the sheet has a
        # date column AND numeric columns, detect the period span and compute
        # per-period deltas (e.g. monthly mean revenue Jan vs Feb = +12%).
        # Much more robust than the cross-sheet sheet-name heuristic (#1) —
        # works even when all data is in one sheet with a Date column.
        period_trend_findings = self._detect_period_trends(
            data, date_profiles, numeric_profiles, sheet_name)
        findings.extend(period_trend_findings)

        return {
            "sheet_name": sheet_name,
            "status": "ok",
            "numeric_profiles": numeric_profiles,
            "categorical_profiles": categorical_profiles,
            "findings": findings,
            "aggregate_findings": aggregate_findings,
            "multi_column_insights": multi_column_insights,
            "column_priority_scores": column_priority_scores,
        }

    def _detect_header_row(self, raw: pd.DataFrame, search_rows: int = 15) -> Optional[int]:
        best_row, best_score, best_unique = None, -1, -1.0
        for idx in range(min(search_rows, len(raw))):
            row = raw.iloc[idx].tolist()
            values = [v for v in row if not pd.isna(v) and clean_text(v)]
            if len(values) < 2:
                continue
            text_like = sum(1 for v in values if not str(v).replace(",", "").replace(".", "").isdigit())
            score = len(values) + text_like * 1.3
            # Tie-breaker: prefer rows whose values are all distinct — real header
            # rows are usually unique, while junk/title rows often repeat the same
            # token (e.g. ["junk","junk","junk"]). Previously a tie kept the
            # earliest row, which could pick a junk row over the real header.
            str_values = [str(v) for v in values]
            unique_ratio = len(set(str_values)) / len(str_values)
            if score > best_score or (score == best_score and unique_ratio > best_unique):
                best_score, best_unique, best_row = score, unique_ratio, idx
        return best_row

    def _detect_period_trends(self, data: pd.DataFrame, date_profiles: List[Dict],
                               numeric_profiles: List[Dict], sheet_name: str) -> List[Dict]:
        """v3 #25: Detect within-sheet YoY/QoQ/MoM trends.

        If the sheet has a date column, group numeric metrics by the detected
        period (year, quarter, or month) and compute deltas between consecutive
        periods. Emits 'period_trend_insight' findings — much more robust than
        the cross-sheet heuristic (#1) because it works on a single sheet with
        a Date column (the common real-world pattern).
        """
        if not date_profiles or not numeric_profiles:
            return []

        # Find the best date column (most non-null dates).
        best_date_col = None
        best_date_series = None
        best_date_count = 0
        for dp in date_profiles:
            col = dp["column"]
            if col not in data.columns:
                continue
            try:
                with warnings.catch_warnings():
                    warnings.simplefilter("ignore", UserWarning)
                    ds = pd.to_datetime(data[col], errors="coerce")
                count = int(ds.notna().sum())
                if count > best_date_count and count >= 2:
                    best_date_col = col
                    best_date_series = ds
                    best_date_count = count
            except Exception:
                pass

        if best_date_series is None or best_date_count < 2:
            return []

        # Determine the period: if dates span > 1 year, use YoY; if > 1 quarter
        # within a year, use QoQ; else use MoM.
        dates_clean = best_date_series.dropna()
        try:
            min_date = dates_clean.min()
            max_date = dates_clean.max()
            span_days = (max_date - min_date).days
        except Exception:
            return []

        if span_days < 40:
            return []  # too short for any meaningful trend

        # Period detection: YoY if span > 365, QoQ if span > 90, else MoM.
        if span_days > 365:
            period_label = "YoY"
            period_group = dates_clean.dt.year
        elif span_days > 90:
            period_label = "QoQ"
            period_group = dates_clean.dt.year.astype(str) + "-Q" + dates_clean.dt.quarter.astype(str)
        else:
            period_label = "MoM"
            period_group = dates_clean.dt.year.astype(str) + "-" + dates_clean.dt.month.astype(str).str.zfill(2)

        findings = []
        for np_profile in numeric_profiles:
            if np_profile.get("is_identifier"):
                continue
            col = np_profile["column"]
            if col not in data.columns or col == best_date_col:
                continue
            try:
                vals = pd.to_numeric(data[col], errors="coerce")
                df = pd.DataFrame({"period": period_group, "val": vals})
                df = df.dropna()
                if len(df) < 2:
                    continue
                per_period = df.groupby("period")["val"].mean().sort_index()
                if len(per_period) < 2:
                    continue
                first_p, last_p = per_period.index[0], per_period.index[-1]
                first_v, last_v = float(per_period.iloc[0]), float(per_period.iloc[-1])
                if first_v == 0:
                    continue
                pct = round((last_v - first_v) / abs(first_v) * 100, 1)
                direction = "increase" if pct >= 0 else "decrease"
                findings.append({
                    "type": "period_trend_insight",
                    "text": (f"{sheet_name}: '{col}' shows a {direction} of {abs(pct)}% "
                             f"({period_label}) from {first_p} ({first_v:.2f}) to "
                             f"{last_p} ({last_v:.2f})."),
                    "priority_score": 0.87,
                    "location": sheet_name,
                    "column": col,
                    "period_type": period_label,
                })
            except Exception:
                pass
        return findings[:10]  # cap to keep the register compact

    def extract_spreadsheet(self, path: Path, item: Dict) -> Dict:
        try:
            sheets_data = []
            xls = pd.ExcelFile(path)
            for sheet_name in xls.sheet_names:
                raw = pd.read_excel(xls, sheet_name=sheet_name, header=None, nrows=self.max_sheet_rows)
                profile = self.profile_dataframe(raw, sheet_name, "excel")
                sheets_data.append(profile)

            # v3 #1 & #2: connect time-ordered sheets into trend insights.
            # Detect sheets that look time-ordered (month/quarter/year) and, for
            # each numeric column present across them, compute a delta/trend.
            # Also consolidate the per-sheet numeric_range findings into a single
            # consolidated entry so 3 monthly sheets -> 1 trend instead of 3
            # near-identical range entries.
            trend_findings = []
            timed = [(s.get("sheet_name"), s) for s in sheets_data
                     if s.get("status") == "ok" and sheet_time_rank(s.get("sheet_name", "")) is not None]
            if len(timed) >= 2:
                timed.sort(key=lambda t: sheet_time_rank(t[0]))
                # Build {column: [(sheet, profile), ...]} across timed sheets.
                col_series = defaultdict(list)
                for sn, sprof in timed:
                    for p in sprof.get("numeric_profiles", []):
                        if p.get("is_identifier"):
                            continue
                        col_series[p["column"]].append((sn, p))
                for col, entries in col_series.items():
                    if len(entries) < 2:
                        continue
                    # Use the mean as the comparable per-sheet statistic.
                    seq = [(sn, p.get("mean")) for sn, p in entries if p.get("mean") is not None]
                    if len(seq) < 2:
                        continue
                    first_sheet, first_val = seq[0]
                    last_sheet, last_val = seq[-1]
                    if first_val == 0:
                        continue
                    pct = round((last_val - first_val) / abs(first_val) * 100, 1)
                    direction = "increase" if pct >= 0 else "decrease"
                    trend_findings.append({
                        "type": "trend_insight",
                        "text": (f"{col}: {direction} of {abs(pct)}% from {first_sheet} "
                                 f"({first_val}) to {last_sheet} ({last_val}) across sheets."),
                        "priority_score": 0.88,
                        "location": f"{first_sheet}→{last_sheet}",
                        "column": col,
                    })
            if trend_findings:
                # Attach trends to the first timed sheet's profile so they flow
                # through build_evidence_register alongside that sheet.
                first_timed = timed[0][1] if timed else None
                if first_timed is not None:
                    first_timed.setdefault("findings", []).extend(trend_findings)
                # v3 #2: collapse per-sheet numeric_range findings for columns
                # that produced a trend into a single consolidated entry across the
                # WHOLE file (keep the first sheet's representative range; drop the
                # cross-sheet duplicates that a trend now subsumes). Columns
                # WITHOUT a trend keep all their per-sheet ranges.
                trend_cols = {f["column"] for f in trend_findings}
                seen_trend_cols = set()
                for sprof in sheets_data:
                    if sprof.get("status") != "ok":
                        continue
                    kept = []
                    for f in sprof.get("findings", []):
                        if f.get("type") == "numeric_range" and f.get("column") in trend_cols:
                            if f["column"] not in seen_trend_cols:
                                kept.append(f)  # keep the first representative across the file
                                seen_trend_cols.add(f["column"])
                            # else: drop this sheet's duplicate (trend subsumes it)
                        else:
                            kept.append(f)
                    sprof["findings"] = kept

            return {"file_id": item["file_id"], "file_name": item["file_name"], "sheets": sheets_data, "status": "ok"}
        except Exception as e:
            return {"file_id": item["file_id"], "status": "error", "error": str(e)}

    def run(self):
        # v3 #23: initialize the centralized logger (console + run.log file).
        # Done here (not __init__) because output_dir must exist for the file.
        global _LOG, _LOG_FILE
        _LOG_FILE = self.output_dir / "run.log"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        _LOG = None  # reset singleton so each run() gets a fresh logger
        self.log = get_logger(
            name="preprocessor", log_file=_LOG_FILE,
            verbose=self.verbose, run_id=self.run_id,
        )

        print("\n=== Impact Slide Preprocessor v4 ===")
        print(f"Filter level: {self.filter_level}")
        self.log.info("pipeline_start", filter_level=self.filter_level,
                       input=str(self.input_path), output=str(self.output_dir))
        if self.verbose:
            print(f"Boost keywords: {self.boost_keywords}")
            print(f"Export options: MD={self.export_md}, CSV={self.export_csv}")
            self.log.debug("config", boost_keywords=self.boost_keywords,
                           export_md=self.export_md, export_csv=self.export_csv)

        run_start = time.perf_counter()

        # [1/5] File discovery & inventory.
        t0 = time.perf_counter()
        files = self.gather_files()
        self.inventory = self.build_file_inventory(files)
        self.timing["stages"]["discovery"] = round(time.perf_counter() - t0, 3)
        print(f"[1/5] Discovered {len(files)} files ({len(self.inventory)} readable)")
        self.log.info("files_discovered", total=len(files), readable=len(self.inventory))

        spreadsheet_count = 0
        pptx_count = 0
        t_extract = time.perf_counter()

        for item in self.inventory:
            if item["access_status"] != "readable":
                continue
            path = Path(item["absolute_path"])
            f0 = time.perf_counter()  # per-file wall-clock (not cumulative)
            category = None
            status = "ok"

            if item["category"] == "spreadsheet":
                category = "spreadsheet"
                print(f"[2/5] Processing spreadsheet: {item['file_name']}")
                profile = self.extract_spreadsheet(path, item)
                self.excel_profiles.append(profile)
                spreadsheet_count += 1
                if profile.get("status") == "error":
                    self.errors.append(profile)
                    status = "error"
                    self.log.error("spreadsheet_failed", file=item['file_name'])
                else:
                    sheet_count = len(profile.get("sheets", []))
                    msg = f"       -> Analyzed {sheet_count} sheets"
                    if self.verbose:
                        msg += f" | Time: {time.perf_counter() - f0:.2f}s"
                    print(msg)
                    self.log.info("spreadsheet_processed", file=item['file_name'],
                                  sheets=sheet_count)

            elif item["category"] == "other" and item["file_name"].lower().endswith(".pptx"):
                category = "pptx"
                print(f"[3/5] Processing PPTX: {item['file_name']}")
                pptx_analysis = self.extract_pptx(path)
                self.pptx_profiles.append(pptx_analysis)
                pptx_count += 1
                if pptx_analysis.get("status") == "error":
                    self.errors.append(pptx_analysis)
                    status = "error"
                    self.log.error("pptx_failed", file=item['file_name'])
                else:
                    slide_count = pptx_analysis.get("total_slides", 0)
                    msg = f"       -> Analyzed {slide_count} slides"
                    if self.verbose:
                        msg += f" | Time: {time.perf_counter() - f0:.2f}s"
                    print(msg)
                    self.log.info("pptx_processed", file=item['file_name'],
                                  slides=slide_count)

            elif item["category"] == "pdf":
                category = "pdf"
                print(f"[3/5] Processing PDF: {item['file_name']}")
                pdf_profile = self.extract_pdf(path, use_ocr=self.enable_ocr)
                self.pdf_profiles.append(pdf_profile)
                if pdf_profile.get("status") == "error":
                    status = "error"
                    self.log.error("pdf_failed", file=item['file_name'])
                else:
                    self.log.info("pdf_processed", file=item['file_name'],
                                 pages=pdf_profile.get("total_pages", 0))

            elif item["category"] == "docx":
                category = "docx"
                print(f"[3/5] Processing DOCX: {item['file_name']}")
                docx_profile = self.extract_docx(path)
                self.docx_profiles.append(docx_profile)
                if docx_profile.get("status") == "error":
                    status = "error"
                    self.log.error("docx_failed", file=item['file_name'])
                else:
                    self.log.info("docx_processed", file=item['file_name'])

            # Record this file's own duration (independent of other files —
            # fixes the old cumulative-elapsed-time bug).
            if category:
                self.timing["files"].append({
                    "file": item["file_name"],
                    "category": category,
                    "duration_s": round(time.perf_counter() - f0, 3),
                    "status": status,
                })

        self.timing["stages"]["extraction"] = round(time.perf_counter() - t_extract, 3)
        self.log.info("extraction_complete", duration_s=self.timing["stages"]["extraction"],
                      files=len(self.timing["files"]))

        print(f"[4/5] Building Evidence Register from {spreadsheet_count} Excel + {pptx_count} PPTX files...")

        # Build rich Evidence Register from all processed data.
        t0 = time.perf_counter()
        self.evidence_register = self.build_evidence_register()
        self.timing["stages"]["evidence_build"] = round(time.perf_counter() - t0, 3)
        self.log.info("evidence_built", entries=len(self.evidence_register),
                      duration_s=self.timing["stages"]["evidence_build"])

        # v4 #26: build the Analyst Briefing (analyst_briefing.md + .json) from
        # the final evidence register + coverage map. Done here (after the
        # register is built) and before _save_outputs() so the artefacts are
        # written alongside the other handoff files.
        t0 = time.perf_counter()
        self._generate_analyst_briefing()
        self.timing["stages"]["briefing"] = round(time.perf_counter() - t0, 3)

        t0 = time.perf_counter()
        # Record total-so-far before saving so the summary has a near-complete
        # total; _save_outputs() generates the summary internally, and we
        # refresh it once more below with the exact output-stage duration.
        self.timing["total_seconds"] = round(time.perf_counter() - run_start, 3)
        self._save_outputs()
        self.timing["stages"]["output"] = round(time.perf_counter() - t0, 3)
        self.timing["total_seconds"] = round(time.perf_counter() - run_start, 3)
        # Refresh the summary report with the final, exact timing numbers
        # (the in-_save_outputs copy was written before output/total were set).
        self._generate_summary_report()
        self.run_finished_at = time.strftime("%Y-%m-%dT%H:%M:%S")

        print(f"[5/5] Complete. Evidence entries: {len(self.evidence_register) if hasattr(self, 'evidence_register') else 0}")

        # v3 #23: always emit run_metadata.json (the reproducibility artifact).
        self._emit_run_metadata()

        self.log.info("pipeline_complete",
                      total_seconds=self.timing['total_seconds'],
                      evidence_entries=len(self.evidence_register) if hasattr(self, 'evidence_register') else 0)

        # v3 #22: always-on timing summary (not verbose-gated). One line for
        # stages, one for per-file breakdown (slowest first).
        st = self.timing["stages"]
        n_files = len(self.timing["files"])
        print(f"[Timing] Total: {self.timing['total_seconds']:.2f}s | "
              f"Discovery: {st.get('discovery', 0):.2f}s | "
              f"Extraction: {st.get('extraction', 0):.2f}s ({n_files} files) | "
              f"Evidence build: {st.get('evidence_build', 0):.2f}s | "
              f"Output: {st.get('output', 0):.2f}s")
        if self.timing["files"]:
            per_file = " | ".join(
                f"{f['file']} {f['duration_s']:.2f}s"
                for f in sorted(self.timing["files"], key=lambda x: x["duration_s"], reverse=True)
            )
            print(f"[Timing] Per file: {per_file}")

    def build_evidence_register(self) -> List[Dict]:
        """
        Improved Evidence Register creation (Conservative version).
        - Uses priority_score from improved profiling
        - Applies minimum priority threshold to filter very low-value evidence
        - Adds structured fields (insight_type, column_name, suggested_use)
        - Sorts by priority (highest first)
        """
        evidence = []
        evidence_id = 1
        thresholds = self._get_filter_thresholds()
        MIN_PRIORITY_THRESHOLD = thresholds["min_priority"]

        # Process Excel profiles
        for file_profile in self.excel_profiles:
            if file_profile.get("status") != "ok":
                continue

            source_file = file_profile.get("file_name", "unknown.xlsx")

            for sheet in file_profile.get("sheets", []):
                sheet_name = sheet.get("sheet_name", "Sheet1")

                for finding in sheet.get("findings", []):
                    priority = finding.get("priority_score", 0.5)

                    # Conservative filter: skip very low priority evidence
                    if priority < MIN_PRIORITY_THRESHOLD:
                        continue

                    insight_type = finding.get("type", "unknown")

                    # v3 #24: stage assignment via the centralized table (was
                    # a hardcoded if/elif chain).
                    suggested_use = self._stages_for(insight_type, finding.get("text", ""))

                    ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "sheet_name": sheet_name,
                        "column_name": finding.get("column"),
                        "insight_type": insight_type,
                        "text": finding.get("text", ""),
                        "priority_score": priority,
                        "confidence": "high" if priority > 0.6 else "medium",
                        "suggested_narrative_use": suggested_use,
                        "source_location": finding.get("location", sheet_name),
                    }
                    evidence.append(ev)
                    evidence_id += 1

                # v3 #5: emit computed group-by aggregate insights (actual data,
                # not just suggestions)
                for finding in sheet.get("aggregate_findings", []):
                    if finding.get("priority_score", 0) < MIN_PRIORITY_THRESHOLD:
                        continue
                    ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "sheet_name": sheet_name,
                        "column_name": finding.get("column"),
                        "insight_type": "aggregate_insight",
                        "text": finding.get("text", ""),
                        "priority_score": finding.get("priority_score", 0.83),
                        "confidence": "high",
                        "suggested_narrative_use": self._stages_for("aggregate_insight", finding.get("text", "")),
                        "source_location": finding.get("location", sheet_name),
                    }
                    evidence.append(ev)
                    evidence_id += 1

                # Add multi-column insights if present (they usually have decent value)
                for insight in sheet.get("multi_column_insights", []):
                    ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "sheet_name": sheet_name,
                        "insight_type": "multi_column_suggestion",
                        "text": insight.get("text", ""),
                        "priority_score": 0.65,
                        "confidence": "medium",
                        "suggested_narrative_use": self._stages_for("multi_column_suggestion", insight.get("text", "")),
                        "source_location": sheet_name,
                    }
                    evidence.append(ev)
                    evidence_id += 1

        # ====================== NEW: Integrate high-value PPTX insights + rich extraction ======================
        for pptx_profile in getattr(self, "pptx_profiles", []):
            if pptx_profile.get("status") != "ok":
                continue
            source_file = Path(pptx_profile.get("file", "unknown.pptx")).name

            for slide in pptx_profile.get("slides", []):
                cls = slide.get("classification", {})
                priority = cls.get("priority_for_evidence", 0)
                slide_type = cls.get("type", "")

                # Option A: Allow section dividers even if priority < 0.45
                if priority < 0.45 and slide_type != "section":
                    continue

                slide_title = slide.get("title", "Slide")
                slide_type = cls.get("type", "content_insight")
                details = slide.get("details", {})

                # Base evidence entry
                ev_text = f"{source_file} - Slide {slide['slide_index']}: {slide_title} ({slide_type})"
                ev = {
                    "evidence_id": f"E{evidence_id:04d}",
                    "source_file": source_file,
                    "sheet_name": None,
                    "insight_type": "pptx_slide_insight",
                    "text": ev_text,
                    "priority_score": round(cls.get("priority_for_evidence", 0.5), 3),
                    "confidence": "high" if cls.get("confidence", 0.5) > 0.8 else "medium",
                    "suggested_narrative_use": self._stages_for(None, ev_text,
                        default=cls.get("recommended_evidence_types") or self.stage_rules.get("slide_type_stages", {}).get(cls.get("type"), ["What", "Why"])),
                    "source_location": f"Slide {slide['slide_index']}",
                    "pptx_classification": cls,
                }
                evidence.append(ev)
                evidence_id += 1

                # === Rich evidence from data-heavy slides ===
                if slide_type in ["data_chart", "data_mixed"]:
                    for chart in details.get("chart_details", []):
                        # Item 1.1: Create summary evidence
                        chart_ev = {
                            "evidence_id": f"E{evidence_id:04d}",
                            "source_file": source_file,
                            "insight_type": "chart_insight",
                            "text": f"Chart: '{chart.get('title', 'Untitled')}' with series: {', '.join(chart.get('series', []))}",
                            "priority_score": 0.88,
                            "confidence": "high",
                            "suggested_narrative_use": self._stages_for("chart_insight", chart.get('title', '')),
                            "source_location": f"Slide {slide['slide_index']}",
                        }
                        evidence.append(chart_ev)
                        evidence_id += 1

                        # Item 1.1: Create chart_data_insight entries from actual values
                        categories = chart.get("categories", [])
                        for sdata in chart.get("series_data", []):
                            s_name = sdata.get("name", "Series")
                            values = sdata.get("values", [])
                            for i, val in enumerate(values[:6]):  # limit per series
                                cat = categories[i] if i < len(categories) else f"Point {i+1}"
                                data_text = f"{cat}: {val} ({s_name})"
                                chart_data_ev = {
                                    "evidence_id": f"E{evidence_id:04d}",
                                    "source_file": source_file,
                                    "insight_type": "chart_data_insight",
                                    "text": data_text,
                                    "priority_score": 0.85,
                                    "confidence": "high",
                                    "suggested_narrative_use": self._stages_for("chart_data_insight", data_text),
                                    "source_location": f"Slide {slide['slide_index']}",
                                }
                                evidence.append(chart_data_ev)
                                evidence_id += 1

                if slide_type in ["data_table", "data_mixed", "data_chart"]:
                    for tbl in details.get("table_details", []):
                        if tbl.get("cells"):
                            # Item 1.2: Extract every cell with priority scoring
                            scored_cells = []
                            for cell in tbl["cells"]:
                                val = cell["value"]
                                score = 0.70  # base
                                # Boost for headers
                                if cell.get("is_header"):
                                    score += 0.12
                                # Boost for genuine numeric values only (a real
                                # data point). Previously *any* digit-bearing
                                # string was boosted, which ranked IPs / log
                                # lines / user-agents (e.g. '222.127.111.234',
                                # 'GET /a.jpg HTTP/1.1') above real insights.
                                try:
                                    float(val.replace(",", ""))
                                    score += 0.10
                                except ValueError:
                                    pass
                                # Penalise noisy/technical cell content that is
                                # rarely a business insight: IP addresses, URLs,
                                # user-agent strings, file paths, HTTP requests.
                                noisy = _looks_like_noise_cell(val)
                                if noisy:
                                    score -= 0.35
                                # Boost for longer, more descriptive text (but
                                # not if it's noise)
                                if len(val) > 30 and not noisy:
                                    score += 0.08
                                score = max(0.05, min(0.95, score))
                                scored_cells.append({
                                    "value": val,
                                    "row": cell["row"],
                                    "col": cell["col"],
                                    "priority": round(score, 3)
                                })

                            # Sort by priority descending
                            scored_cells = sorted(scored_cells, key=lambda x: x["priority"], reverse=True)

                            for cell in scored_cells[:15]:  # cap at 15 per table
                                cell_ev = {
                                    "evidence_id": f"E{evidence_id:04d}",
                                    "source_file": source_file,
                                    "insight_type": "table_cell",
                                    "text": f"Table cell (r{cell['row']},c{cell['col']}): {cell['value']}",
                                    "priority_score": cell["priority"],
                                    "confidence": "medium",
                                    "suggested_narrative_use": self._stages_for("table_cell", str(cell.get('value', ''))),
                                    "source_location": f"Slide {slide['slide_index']}",
                                    "table_cell": True,
                                }
                                evidence.append(cell_ev)
                                evidence_id += 1
                        elif tbl.get("rows", 0) > 1:
                            # Fallback old behavior
                            table_ev = {
                                "evidence_id": f"E{evidence_id:04d}",
                                "source_file": source_file,
                                "insight_type": "table_insight",
                                "text": f"Table ({tbl.get('rows')} rows × {tbl.get('cols')} cols)",
                                "priority_score": 0.82,
                                "confidence": "high",
                                "suggested_narrative_use": self._stages_for("table_insight", f"Table {tbl.get('rows')}x{tbl.get('cols')}"),
                                "source_location": f"Slide {slide['slide_index']}",
                            }
                            evidence.append(table_ev)
                            evidence_id += 1

                # Extracted metrics as evidence (prefer rich advanced version when available)
                advanced_metrics = details.get("extracted_advanced_metrics", [])
                if advanced_metrics:
                    for m in advanced_metrics:
                        metric_ev = {
                            "evidence_id": f"E{evidence_id:04d}",
                            "source_file": source_file,
                            "insight_type": "text_metric",
                            "text": f"{m.get('context', '')}",
                            "priority_score": 0.82,
                            "confidence": "high",
                            "suggested_narrative_use": self._stages_for("text_metric", m.get('context', '')),
                            "source_location": f"Slide {slide['slide_index']}",
                            "metric_value": m.get("value"),
                            "metric_type": m.get("type"),
                        }
                        evidence.append(metric_ev)
                        evidence_id += 1
                else:
                    for metric in details.get("extracted_metrics", []):
                        metric_ev = {
                            "evidence_id": f"E{evidence_id:04d}",
                            "source_file": source_file,
                            "insight_type": "text_metric",
                            "text": f"Key metric on slide {slide['slide_index']}: {metric}",
                            "priority_score": 0.78,
                            "confidence": "medium",
                            "suggested_narrative_use": self._stages_for("text_metric", str(metric)),
                            "source_location": f"Slide {slide['slide_index']}",
                        }
                        evidence.append(metric_ev)
                        evidence_id += 1

                # Key bullets from conclusion / insight slides
                if slide_type in ["conclusion", "content_insight"]:
                    # v3 #24: conclusion-bullet stages now come from the
                    # configurable table (conclusion_bullet_stages) instead of
                    # a hardcoded literal. The keyword-override layer in
                    # _stages_for() can still redirect individual bullets by
                    # their text (e.g. a bullet mentioning 'revenue').
                    if slide_type == "conclusion":
                        bullet_use = list(self.stage_rules.get("conclusion_bullet_stages", DEFAULT_CONCLUSION_BULLET_STAGES))
                    else:
                        bullet_use = self._stages_for("bullet_insight", "")
                    for bullet in details.get("key_bullets", [])[:3]:
                        # v3 #3: rank individual bullets by insight-language
                        # density. Previously every bullet scored a flat 0.75,
                        # so "Recommendation: expand" ranked equal to "LogLevel...".
                        # Now insight-bearing bullets get up to +0.25.
                        bp = insight_priority_boost(bullet, 0.75)
                        # v3 #24: per-bullet stage = keyword override (first
                        # match) else the slide-type default (conclusion stages
                        # or bullet_insight default).
                        bullet_stages = list(bullet_use)
                        for pat, stages in self.stage_rules.get("_compiled_keyword_overrides", []):
                            if pat.search(bullet):
                                bullet_stages = list(stages)
                                break
                        bullet_ev = {
                            "evidence_id": f"E{evidence_id:04d}",
                            "source_file": source_file,
                            "insight_type": "bullet_insight",
                            "text": bullet,
                            "priority_score": bp,
                            "confidence": "high" if bp >= 0.85 else "medium",
                            "suggested_narrative_use": bullet_stages,
                            "source_location": f"Slide {slide['slide_index']}",
                        }
                        evidence.append(bullet_ev)
                        evidence_id += 1

                # === NEW: Extract process step texts from diagram/process flow slides (Item 4.3 extension) ===
                if slide_type == "diagram_process" or details.get("has_process_diagram"):
                    for step in details.get("process_steps", [])[:6]:
                        if len(step) > 5:
                            step_ev = {
                                "evidence_id": f"E{evidence_id:04d}",
                                "source_file": source_file,
                                "insight_type": "process_step",
                                "text": step,
                                "priority_score": 0.76,
                                "confidence": "medium",
                                "suggested_narrative_use": self._stages_for("process_step", step),
                                "source_location": f"Slide {slide['slide_index']}",
                            }
                            evidence.append(step_ev)
                            evidence_id += 1

                # Item 1.3: Speaker notes as evidence
                notes = details.get("speaker_notes", "")
                if notes and len(notes) > 20:
                    notes_ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "insight_type": "speaker_notes_insight",
                        "text": notes,
                        "priority_score": 0.85,
                        "confidence": "high",
                        "suggested_narrative_use": self._stages_for("speaker_notes_insight", notes),
                        "source_location": f"Slide {slide['slide_index']}",
                    }
                    evidence.append(notes_ev)
                    evidence_id += 1

                # Item 1.6: Bold / emphasized text as evidence
                for bold in details.get("bold_texts", [])[:4]:
                    bold_ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "insight_type": "emphasized_text",
                        "text": bold,
                        "priority_score": 0.78,
                        "confidence": "medium",
                        "suggested_narrative_use": self._stages_for("emphasized_text", bold),
                        "source_location": f"Slide {slide['slide_index']}",
                    }
                    evidence.append(bold_ev)
                    evidence_id += 1

            # Item 1.5: Section divider evidence (moved outside bold loop)
            if cls.get("type") == "section":
                section_ev = {
                    "evidence_id": f"E{evidence_id:04d}",
                    "source_file": source_file,
                    "insight_type": "section_divider",
                    "text": f"Section divider: {slide_title}",
                    "priority_score": 0.22,
                    "confidence": "medium",
                    "suggested_narrative_use": self._stages_for("section_divider", slide_title),
                    "source_location": f"Slide {slide['slide_index']}",
                }
                evidence.append(section_ev)
                evidence_id += 1

        # ====================== PDF + DOCX EVIDENCE (3.2 + 3.3) ======================
        for pdf_profile in getattr(self, "pdf_profiles", []):
            if pdf_profile.get("status") != "ok":
                continue
            source_file = Path(pdf_profile.get("file", "unknown.pdf")).name

            for page in pdf_profile.get("pages", []):
                text = page.get("text", "").strip()
                is_ocr = page.get("ocr_used", False)

                # Improved PDF evidence logic (Steps 1, 2, 3)
                min_len = 15
                has_insight = contains_insight_language(text) > 0.3

                # Only create evidence if text is long enough OR has strong insight language
                if len(text) >= min_len or has_insight:
                    # Dynamic priority
                    base_priority = 0.62
                    if has_insight:
                        base_priority = 0.78
                    if len(text) > 150:
                        base_priority += 0.05

                    insight_type = "pdf_ocr_page_insight" if is_ocr else "pdf_page_insight"

                    ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "insight_type": insight_type,
                        "text": f"Page {page['page']}: {text}",
                        "priority_score": round(base_priority, 3),
                        "confidence": "medium" if is_ocr else ("high" if has_insight else "medium"),
                        "suggested_narrative_use": self._stages_for(insight_type, text),
                        "source_location": f"Page {page['page']}",
                        "ocr_used": is_ocr,
                    }
                    evidence.append(ev)
                    evidence_id += 1

            for tbl in pdf_profile.get("tables", []):
                # v3: richer PDF table evidence using the merged pdfplumber/PyMuPDF
                # output (header + cols + engine) instead of just row count + preview.
                rows = tbl.get("data", [])
                header = tbl.get("header") or (rows[0] if rows else [])
                cols = tbl.get("cols") or (len(header) if header else 0)
                engine = tbl.get("engine", "pymupdf")
                preview = " | ".join(str(c) for c in header[:4]) if header else ""

                text = f"PDF Table on page {tbl['page']} ({tbl.get('rows', len(rows))} rows × {cols} cols)"
                if preview:
                    text += f": {preview}"

                priority = 0.72
                if len(rows) > 3:
                    priority = 0.78
                # pdfplumber detection is more reliable -> higher confidence
                confidence = "high" if engine == "pdfplumber" else "medium"

                ev = {
                    "evidence_id": f"E{evidence_id:04d}",
                    "source_file": source_file,
                    "insight_type": "pdf_table_insight",
                    "text": text,
                    "priority_score": priority,
                    "confidence": confidence,
                    "suggested_narrative_use": self._stages_for("pdf_table_insight", text),
                    "source_location": f"Page {tbl['page']}",
                }
                evidence.append(ev)
                evidence_id += 1
                # v3: also seed per-cell evidence (capped) so the Analyst gets the
                # actual table contents, not just a one-line summary.
                for r_idx, row in enumerate(rows[:5]):
                    for c_idx, cell in enumerate(row[:6]):
                        if cell and str(cell).strip():
                            cell_ev = {
                                "evidence_id": f"E{evidence_id:04d}",
                                "source_file": source_file,
                                "insight_type": "pdf_table_cell",
                                "text": f"PDF table cell (p{tbl['page']} r{r_idx} c{c_idx}): {cell}",
                                "priority_score": 0.70,
                                "confidence": confidence,
                                "suggested_narrative_use": self._stages_for("pdf_table_cell", str(cell)),
                                "source_location": f"Page {tbl['page']}",
                            }
                            evidence.append(cell_ev)
                            evidence_id += 1

        for docx_profile in getattr(self, "docx_profiles", []):
            if docx_profile.get("status") != "ok":
                continue
            source_file = Path(docx_profile.get("file", "unknown.docx")).name

            for para in docx_profile.get("paragraphs", [])[:10]:
                if len(para) > 20:
                    ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": source_file,
                        "insight_type": "docx_insight",
                        "text": para,
                        "priority_score": 0.60,
                        "confidence": "medium",
                        "suggested_narrative_use": self._stages_for("docx_insight", para),
                        "source_location": "DOCX",
                    }
                    evidence.append(ev)
                    evidence_id += 1

        # Sort by priority_score descending (most important evidence first)
        evidence.sort(key=lambda x: x.get("priority_score", 0), reverse=True)

        # === Item 4.4: Detect cross-file relationships (Excel + PPTX) ===
        cross_file_evidence = self._find_cross_file_relationships(evidence)
        evidence.extend(cross_file_evidence)

        # Re-sort after adding cross-file evidence
        evidence.sort(key=lambda x: x.get("priority_score", 0), reverse=True)

        # === Item 4.2 (+ v3 #10): Evidence Deduplication (exact + semantic) ===
        evidence = self._deduplicate_evidence(evidence)

        # === Item 4.3: Apply configurable boost keywords ===
        evidence = self._apply_boost_rules(evidence)

        # v3 #8: cap per (source_file, insight_type) so one source can't flood
        evidence = self._apply_per_type_caps(evidence)

        # v3 #6 + #9: stamp extraction_method + reliability-based confidence
        evidence = self._apply_provenance_and_confidence(evidence)

        # v3 #12: top-entities summary (per Excel categorical column)
        self.entities_summary = self._build_entities_summary()

        # v3 #4: attach a coverage map so the Analyst GPT gets a structured
        # signal of where evidence is thin (per Why/What/How/Now stage and per
        # source file) before it starts reasoning. Surfaced in the JSON handoff
        # and the summary report.
        self.coverage_map = self._build_coverage_map(evidence)
        if self.entities_summary:
            self.coverage_map["entities_summary_count"] = len(self.entities_summary)
        # v3: surface per-entity mention stats (which entities appear across how
        # many files) into the coverage map so the Analyst GPT gets a structured
        # cross-file signal.
        if getattr(self, "_entity_mention_stats", None):
            self.coverage_map["entity_mentions"] = self._entity_mention_stats

        # v3: validate every evidence entry against the Pydantic contract before
        # it's handed off. Malformed entries are dropped to processing_errors
        # instead of silently shipping bad data to the Analyst GPT.
        evidence = self._validate_evidence(evidence)

        return evidence

    def _validate_evidence(self, evidence: List[Dict]) -> List[Dict]:
        """Validate evidence entries against the EvidenceEntry schema. Bad entries
        are logged to self.errors and dropped. If pydantic isn't installed, this
        is a no-op (the pipeline still runs, just without runtime guarantees).

        Uniformly truncates every entry's `text` field to self.max_text_length
        (default = schemas.MAX_TEXT_LENGTH = 800) *before* validation, so the
        schema's max_length constraint never rejects a too-long string — it
        just sees the already-truncated text. This replaces the per-extractor
        caps (150/200/300) that were applied inconsistently across PDF/PPTX/
        DOCX paths; those extractors now store the full text and this single
        chokepoint enforces the ceiling."""
        cap = self.max_text_length
        for ev in evidence:
            t = ev.get("text")
            if isinstance(t, str) and len(t) > cap:
                # Truncate to cap-1 chars + ellipsis (U+2026) = exactly cap chars.
                ev["text"] = t[:cap - 1] + "\u2026"
        if not _HAS_PYDANTIC:
            return evidence
        kept = []
        for ev in evidence:
            try:
                EvidenceEntry(**ev)
                kept.append(ev)
            except Exception as ex:
                self.errors.append({
                    "file": ev.get("source_file", "?"),
                    "evidence_id": ev.get("evidence_id", "?"),
                    "error": f"schema validation failed: {ex}",
                })
        if len(kept) < len(evidence):
            print(f"       [warn] {len(evidence) - len(kept)} evidence entries failed "
                  f"schema validation and were dropped (see processing_errors.json)")
        return kept

    def _build_coverage_map(self, evidence: List[Dict]) -> Dict[str, Any]:
        """Summarise evidence coverage by narrative stage and source file."""
        from collections import Counter
        stages = Counter()
        for e in evidence:
            for s in e.get("suggested_narrative_use", []):
                stages[s] += 1
        by_source = Counter(e.get("source_file", "?") for e in evidence)
        by_type = Counter(e.get("insight_type", "?") for e in evidence)
        all_stages = {"Why", "What", "How", "Now"}
        return {
            "total_evidence": len(evidence),
            "by_narrative_stage": {s: stages.get(s, 0) for s in sorted(all_stages)},
            "stages_with_no_evidence": sorted(all_stages - set(stages)),
            "by_source_file": dict(by_source.most_common()),
            "by_insight_type": dict(by_type.most_common()),
            "avg_priority": round(sum(e.get("priority_score", 0) for e in evidence) / max(len(evidence), 1), 3),
        }

    def _deduplicate_evidence(self, evidence: List[Dict]) -> List[Dict]:
        """
        Item 4.2 (+ v3 #10, v3 #20): Remove near-duplicate evidence entries.
        Keeps the highest priority version of similar evidence.

        Two passes: (1) exact normalized-prefix dedup (cheap, catches the
        cross-sheet/column repeats), then (2) semantic near-dup clustering.

        v3 #20: Pass 2 now uses a TIERED semantic engine (sentence-transformers
        embeddings -> pure-numpy TF-IDF + cosine -> rapidfuzz char-similarity)
        so it clusters TRUE semantic near-duplicates (sharing few character
        n-grams, e.g. "North America revenue grew 12%" vs "US & Canada sales up
        a tenth") rather than only lexical rephrasings. When a near-dup is
        dropped, its source_file + evidence_id are merged onto the surviving
        entry (dedup_merged_sources / dedup_merged_ids) so source provenance is
        preserved instead of silently lost. The dedup engine is selectable via
        self.dedup_engine (auto/embeddings/tfidf/fuzzy).
        """
        if not evidence:
            return evidence

        # Pass 1: exact normalized prefix.
        seen_texts = {}   # normalized text -> best evidence
        for ev in evidence:
            text = ev.get("text", "").strip().lower()
            if not text:
                continue
            norm_text = re.sub(r'\s+', ' ', text)[:120]
            existing = seen_texts.get(norm_text)
            if existing is None or ev.get("priority_score", 0) > existing.get("priority_score", 0):
                seen_texts[norm_text] = ev
        pass1 = list(seen_texts.values())

        # Pass 2: tiered semantic near-dup clustering. Only short text entries
        # (bullets/paragraphs/cells, <=200 chars) are clustered, where phrasing
        # variance is common; longer entries are kept as-is to avoid merging
        # genuinely distinct long insights. The engine precomputes a similarity
        # matrix once (for embeddings/TF-IDF), so the greedy loop stays cheap.
        candidates = []   # list of pass1 indices that are clusterable
        cand_texts = []
        for i, ev in enumerate(pass1):
            raw = ev.get("text", "")
            text = raw.strip().lower()
            if text and len(raw) <= 200:
                candidates.append(i)
                cand_texts.append(text)
        cand_map = {pi: ci for ci, pi in enumerate(candidates)}  # pass1_idx -> cand_idx
        engine = _SemanticDedupEngine(cand_texts, engine=getattr(self, "dedup_engine", "auto"))

        order = sorted(range(len(pass1)),
                       key=lambda i: pass1[i].get("priority_score", 0), reverse=True)
        kept = []  # list of (pass1_idx, ev)
        for pi in order:
            ev = pass1[pi]
            raw = ev.get("text", "")
            text = raw.strip().lower()
            if not text or len(raw) > 200:
                kept.append((pi, ev))
                continue
            ci = cand_map.get(pi)
            if ci is None:
                kept.append((pi, ev))
                continue
            is_dup = False
            for (kpi, k) in kept:
                kci = cand_map.get(kpi)
                if kci is None:
                    continue
                if len(k.get("text", "")) > 200:
                    continue
                if engine.similar(ci, kci):
                    is_dup = True
                    # Merge source provenance from the dropped near-dup onto the
                    # surviving (higher-priority) entry so the Analyst GPT still
                    # sees that multiple sources backed this insight.
                    merged_src = k.setdefault("dedup_merged_sources", [])
                    merged_ids = k.setdefault("dedup_merged_ids", [])
                    src = ev.get("source_file")
                    if src and src not in merged_src and src != k.get("source_file"):
                        merged_src.append(src)
                    eid = ev.get("evidence_id")
                    if eid and eid != k.get("evidence_id") and eid not in merged_ids:
                        merged_ids.append(eid)
                    break
            if not is_dup:
                kept.append((pi, ev))
        return [ev for (_, ev) in kept]

    def _apply_boost_rules(self, evidence: List[Dict]) -> List[Dict]:
        """
        Item 4.3: Boost evidence that matches user-defined keywords.
        """
        if not self.boost_keywords:
            return evidence

        for ev in evidence:
            text = ev.get("text", "").lower()
            for kw in self.boost_keywords:
                kw = kw.lower().strip()
                if kw and kw in text:
                    # Boost priority (cap at 0.98)
                    current = ev.get("priority_score", 0.5)
                    ev["priority_score"] = round(min(0.98, current + 0.15), 3)
                    ev["boosted_by_rule"] = kw
                    break   # only boost once per entry

        return evidence

    def _apply_provenance_and_confidence(self, evidence: List[Dict]) -> List[Dict]:
        """
        v3 #6 + #9: stamp every evidence with an `extraction_method` (how the
        insight was derived) and a reliability-based `confidence`, so the
        Analyst GPT can weight computed/chart data above OCR'd text. Existing
        confidence values are kept only if they're already stricter than the
        method-based default.
        """
        for ev in evidence:
            itype = ev.get("insight_type", "")
            method = self._method_for_insight(itype, ev)
            ev["extraction_method"] = method
            method_conf = confidence_for_method(method)
            # Downgrade to method confidence if the entry was over-confident
            # for its reliability (e.g. an OCR page shouldn't be 'high').
            current = ev.get("confidence", "medium")
            rank = {"high": 2, "medium": 1, "low": 0}
            if rank.get(method_conf, 1) < rank.get(current, 1):
                ev["confidence"] = method_conf
            elif not current:
                ev["confidence"] = method_conf
        return evidence

    @staticmethod
    def _method_for_insight(itype: str, ev: Dict) -> str:
        """Map an insight_type (and entry flags) to an extraction method."""
        if itype in ("trend_insight", "period_trend_insight", "aggregate_insight",
                     "outlier_insight", "correlation_insight"):
            return "computed"
        if itype == "chart_data_insight":
            return "chart_data"
        if itype == "chart_insight":
            return "chart_data"
        if itype == "numeric_range":
            return "numeric_range"
        if itype == "categorical_distribution":
            return "categorical"
        if itype == "table_cell":
            return "table_cell"
        if itype in ("pdf_ocr_page_insight",) or ev.get("ocr_used"):
            return "ocr"
        if itype == "pdf_page_insight":
            return "text_layer"
        if itype in ("pdf_table_insight",):
            return "text_layer"
        if itype == "pdf_table_cell":
            return "table_cell"
        if itype == "bullet_insight":
            return "bullet"
        if itype == "docx_insight":
            return "paragraph"
        if itype == "cross_file_metric":
            return "cross_file"
        if itype in ("pptx_slide_insight", "section_divider",
                     "speaker_notes_insight", "emphasized_text",
                     "process_step", "text_metric", "table_insight"):
            return "classifier"
        return "unknown"

    def _apply_per_type_caps(self, evidence: List[Dict]) -> List[Dict]:
        """
        v3 #8: cap the number of evidence entries per (source_file, insight_type)
        so a single source can't drown out the register. Keeps the highest-
        priority representatives. Caps only the high-volume types that tend to
        flood (bullets, slide insights, table cells); analytical types
        (trends, aggregates, numeric ranges) are uncapped since each is
        distinct and high-value.
        """
        from collections import defaultdict
        CAPS = {
            "bullet_insight": 20,
            "pptx_slide_insight": 15,
            "table_cell": 12,
            "categorical_distribution": 12,
            "outlier_insight": 10,
            "correlation_insight": 8,
            "period_trend_insight": 10,
        }
        buckets = defaultdict(list)
        for ev in evidence:
            key = (ev.get("source_file", "?"), ev.get("insight_type", "?"))
            buckets[key].append(ev)
        kept = []
        for (src, itype), items in buckets.items():
            cap = CAPS.get(itype)
            if cap is not None and len(items) > cap:
                items = sorted(items, key=lambda x: x.get("priority_score", 0), reverse=True)[:cap]
            kept.extend(items)
        kept.sort(key=lambda x: x.get("priority_score", 0), reverse=True)
        return kept

    def _build_entities_summary(self) -> List[Dict]:
        """
        v3 #12: build a 'top entities' summary block — the top categorical
        values per Excel column — so the Analyst gets ready-made segmentation
        anchors (e.g. top product lines, top regions). Emitted as a separate
        handoff file (`entities_summary.json`) and folded into the coverage map.
        """
        entities = []
        for prof in getattr(self, "excel_profiles", []):
            if prof.get("status") != "ok":
                continue
            source_file = prof.get("file_name", "unknown.xlsx")
            for sheet in prof.get("sheets", []):
                sheet_name = sheet.get("sheet_name", "Sheet1")
                for c in sheet.get("categorical_profiles", []):
                    col = c.get("column")
                    top = c.get("top_values", [])[:5]
                    if not col or not top:
                        continue
                    total = sum(t.get("count", 0) for t in top) or 1
                    entities.append({
                        "source_file": source_file,
                        "sheet": sheet_name,
                        "column": col,
                        "unique_count": c.get("unique_count"),
                        "top_values": [
                            {"value": t.get("value"), "count": t.get("count"),
                             "share_pct": round(t.get("count", 0) / total * 100, 1)}
                            for t in top
                        ],
                    })
        return entities

    def _find_cross_file_relationships(self, all_evidence: List[Dict]) -> List[Dict]:
        """
        Item 4.4: Multi-file relationship insights.
        Detects when the same metric, region, product, or value appears
        in both Excel and PPTX files and creates boosted cross-file evidence.
        """
        cross_evidence = []
        evidence_id = 9000   # start from high number to avoid collision

        # Separate evidence by source type and track the actual filenames so
        # cross-file evidence stays source-backed (the Analyst must be able to
        # trace every insight back to real files).
        excel_files = sorted({e.get("source_file") for e in all_evidence
                              if str(e.get("source_file", "")).lower().endswith(".xlsx")})
        pptx_files = sorted({e.get("source_file") for e in all_evidence
                             if str(e.get("source_file", "")).lower().endswith(".pptx")})
        excel_evidence = [e for e in all_evidence if str(e.get("source_file", "")).lower().endswith(".xlsx")]
        pptx_evidence = [e for e in all_evidence if str(e.get("source_file", "")).lower().endswith(".pptx")]

        if not excel_evidence or not pptx_evidence:
            return cross_evidence

        related = excel_files + pptx_files

        # 1. Find common entities (Region, Product, etc.).
        #    Rather than a fixed 7-word list (which misses most real datasets),
        #    derive candidate entities dynamically from the actual categorical
        #    column values the profiler extracted from the Excel files
        #    (e.g. City -> Naypyitaw, Yangon, Mandalay). Keep a small curated
        #    list as a bonus for decks that reference generic segments.
        keywords = set(["north", "south", "east", "west", "enterprise", "smb", "startup"])
        derived = set()
        for prof in getattr(self, "excel_profiles", []):
            if prof.get("status") != "ok":
                continue
            for sheet in prof.get("sheets", []):
                for c in sheet.get("categorical_profiles", []):
                    for tv in c.get("top_values", []):
                        v = clean_text(tv.get("value")).lower()
                        # Multi-word / longer categorical values are real
                        # entities; skip trivial ones (single letters, digits).
                        # Also keep short values that are known abbreviations
                        # (US, EU, EMEA, ...) so the abbreviation-expansion
                        # matcher can link them to their full forms elsewhere.
                        is_abbrev = v in _ABBREVIATIONS
                        if v and (len(v) >= 3 or is_abbrev) and not v.replace(".", "").isdigit():
                            derived.add(v)
        keywords |= derived

        excel_entities = set()
        pptx_entities = set()
        # v3: track per-entity mention stats so the Analyst GPT gets a
        # structured "this entity appears in N files" signal, not just a
        # binary "in both". Maps entity -> set of source files mentioning it.
        entity_files: Dict[str, set] = defaultdict(set)

        for ev in all_evidence:   # Search in ALL evidence, not just numeric
            text = ev.get("text", "").lower()
            src = str(ev.get("source_file", "")).lower()
            src_name = ev.get("source_file")
            for kw in keywords:
                if kw and _entity_in_text(kw, text):
                    if src.endswith(".xlsx"):
                        excel_entities.add(kw)
                    elif src.endswith(".pptx"):
                        pptx_entities.add(kw)
                    if src_name:
                        entity_files[kw].add(src_name)

        common_entities = excel_entities.intersection(pptx_entities)
        # Cap entity matches so a large shared vocabulary can't flood the register.
        ENTITY_CROSS_CAP = 5
        self._entity_mention_stats = {}  # surfaced into the coverage map
        for entity in sorted(common_entities,
                             key=lambda e: len(entity_files.get(e, set())),
                             reverse=True)[:ENTITY_CROSS_CAP]:
            files_list = sorted(entity_files.get(entity, set()))
            n_files = len(files_list)
            self._entity_mention_stats[entity] = {
                "files": files_list, "file_count": n_files,
                "in_excel": entity in excel_entities,
                "in_pptx": entity in pptx_entities,
            }
            cross_ev = {
                "evidence_id": f"E{evidence_id:04d}",
                "source_file": related[0] if related else None,
                "insight_type": "cross_file_metric",
                "text": (f"'{entity.title()}' mentioned in {n_files} file(s): "
                         f", ".join(files_list) if files_list
                         else f"'{entity.title()}' mentioned in both Excel data and PPTX narrative"),
                "priority_score": 0.90,
                "confidence": "high",
                "suggested_narrative_use": self._stages_for("cross_file_metric", entity),
                "source_location": "Cross-file",
                "related_files": related,
            }
            cross_evidence.append(cross_ev)
            evidence_id += 1

        # 2. Find overlapping NUMERIC values — but only *distinctive* ones, to
        #    avoid false positives between unrelated files. Bare small integers
        #    (5, 6, 10, 17, 42 ...) collide constantly across any two documents
        #    and produce misleading 'cross-file' evidence. We treat a number as
        #    distinctive only if it is either (a) a decimal/decimal-looking value
        #    that reads like a real metric, or (b) an integer >= 100. We also cap
        #    the number of numeric cross-file matches per source so a single
        #    collision-rich file can't flood the register.
        def _is_distinctive(n: str) -> bool:
            if "." in n:
                return True   # decimal value, e.g. 99.96 / 0.6045
            try:
                return int(n) >= 100
            except ValueError:
                return False

        excel_numbers = set()
        for ev in excel_evidence:
            nums = re.findall(r'\d+(?:\.\d+)?', ev.get("text", ""))
            for n in nums:
                if _is_distinctive(n):
                    excel_numbers.add(n)

        numeric_cross_count = 0
        NUMERIC_CROSS_CAP = 3   # max numeric cross-file matches overall
        for ev in pptx_evidence:
            if numeric_cross_count >= NUMERIC_CROSS_CAP:
                break
            nums = re.findall(r'\d+(?:\.\d+)?', ev.get("text", ""))
            for num in nums:
                if numeric_cross_count >= NUMERIC_CROSS_CAP:
                    break
                if num in excel_numbers and _is_distinctive(num):
                    cross_ev = {
                        "evidence_id": f"E{evidence_id:04d}",
                        "source_file": related[0] if related else None,
                        "insight_type": "cross_file_metric",
                        "text": f"Numeric value '{num}' appears in both Excel and PPTX",
                        "priority_score": 0.88,
                        "confidence": "medium",
                        "suggested_narrative_use": self._stages_for("cross_file_metric", str(num)),
                        "source_location": "Cross-file",
                        "related_files": related,
                    }
                    cross_evidence.append(cross_ev)
                    evidence_id += 1
                    break

        return cross_evidence

    def _generate_analyst_briefing(self):
        """v4 #26: build the condensed strategic handoff for the Analyst GPT.

        Consumes the already-built evidence register + coverage map (no
        re-derivation) and produces ``analyst_briefing.md`` +
        ``analyst_briefing.json``. Skipped gracefully if the analyst_briefing
        module is unavailable or there is no evidence register yet (the
        zero-evidence case still produces a 'no_evidence' briefing so the
        Analyst GPT gets an explicit gap signal rather than silence).
        """
        if not _HAS_BRIEFING:
            return
        evidence = getattr(self, "evidence_register", None) or []
        coverage_map = getattr(self, "coverage_map", None) or {}
        # Cross-file relationships = the cross_file_metric entries in the final
        # register (single source of truth — already deduped/capped/validated).
        cross_file = [e for e in evidence
                      if e.get("insight_type") == "cross_file_metric"]
        run_metadata = {
            "run_id": self.run_id,
            "source_folder": str(self.input_path),
        }
        generator = AnalystBriefingGenerator(
            evidence=evidence,
            coverage_map=coverage_map,
            cross_file=cross_file,
            run_metadata=run_metadata,
            boost_keywords=self.boost_keywords,
            business_keywords=self.briefing_business_keywords,
            readiness_weights=self.briefing_readiness_weights,
            focus_weights=self.briefing_focus_weights,
            focus_areas_count=self.focus_areas_count,
        )
        briefing = generator.generate()
        self.analyst_briefing_json = generator.to_json(briefing)
        self.analyst_briefing_md = generator.render_markdown(briefing)
        if self.log:
            self.log.info(
                "analyst_briefing_generated",
                overall_score=briefing["narrative_readiness"]["overall_score"],
                focus_areas=len(briefing["suggested_focus_areas"]),
                quality_flags=len(briefing["quality_flags"]),
            )

    def _save_outputs(self):
        self.output_dir.mkdir(parents=True, exist_ok=True)

        with open(self.output_dir / "file_inventory.json", "w") as f:
            json.dump(self.inventory, f, indent=2)

        with open(self.output_dir / "excel_profile.json", "w") as f:
            json.dump(self.excel_profiles, f, indent=2)

        # v3: emit the Evidence Register JSON Schema (the contract the Analyst
        # GPT can reference). Always written when pydantic is available so the
        # schema stays in sync with the actual data shape.
        if _HAS_PYDANTIC:
            with open(self.output_dir / "evidence_schema.json", "w") as f:
                json.dump(EvidenceEntry.model_json_schema(), f, indent=2)

        if self.pptx_profiles:
            with open(self.output_dir / "pptx_profile.json", "w") as f:
                json.dump(self.pptx_profiles, f, indent=2)
            print(f"PPTX profile saved with {len(self.pptx_profiles)} files")

        # Save improved Evidence Register
        if hasattr(self, "evidence_register") and self.evidence_register:
            with open(self.output_dir / "evidence_register_seed.json", "w") as f:
                json.dump(self.evidence_register, f, indent=2)
            print(f"Evidence Register created with {len(self.evidence_register)} entries")

        # v3 #4: save the coverage map as a standalone handoff file
        if hasattr(self, "coverage_map") and self.coverage_map:
            with open(self.output_dir / "coverage_map.json", "w") as f:
                json.dump(self.coverage_map, f, indent=2)
            print(f"Coverage map saved ({self.coverage_map.get('total_evidence', 0)} entries, "
                  f"stages: {self.coverage_map.get('by_narrative_stage', {})})")

        # v3 #12: save the top-entities summary (per Excel categorical column)
        if getattr(self, "entities_summary", None):
            with open(self.output_dir / "entities_summary.json", "w") as f:
                json.dump(self.entities_summary, f, indent=2)
            print(f"Entities summary saved: {len(self.entities_summary)} columns")

        # v4 #26: write the Analyst Briefing (the condensed strategic handoff
        # for the Impact Slide Analyst GPT). Always written when the briefing
        # module is available — even the zero-evidence case emits a briefing
        # with a 'no_evidence' flag so the GPT gets an explicit gap signal.
        if getattr(self, "analyst_briefing_md", None) is not None:
            (self.output_dir / "analyst_briefing.md").write_text(
                self.analyst_briefing_md, encoding="utf-8")
        if getattr(self, "analyst_briefing_json", None) is not None:
            with open(self.output_dir / "analyst_briefing.json", "w", encoding="utf-8") as f:
                json.dump(self.analyst_briefing_json, f, indent=2)
            if self.analyst_briefing_json.get("narrative_readiness"):
                score = self.analyst_briefing_json["narrative_readiness"]["overall_score"]
                nfa = len(self.analyst_briefing_json.get("suggested_focus_areas", []))
                print(f"Analyst briefing saved (readiness {score}/100, {nfa} focus areas)")

        # Save filtering log (very useful for debugging)
        if self.filtered_items:
            with open(self.output_dir / "filtering_log.json", "w") as f:
                json.dump(self.filtered_items, f, indent=2)
            print(f"Filtering log saved: {len(self.filtered_items)} items filtered")

        if self.errors:
            with open(self.output_dir / "processing_errors.json", "w") as f:
                json.dump(self.errors, f, indent=2)

        # Generate human-readable summary report
        self._generate_summary_report()

        # Item 4.6: Export Evidence Register
        if self.export_md:
            md_path = self.output_dir / "evidence_register.md"
            md_lines = ["# Evidence Register\n"]
            for ev in self.evidence_register:
                md_lines.append(f"### {ev['evidence_id']} — {ev['insight_type']} (Priority: {ev['priority_score']:.2f})")
                md_lines.append(f"**Text:** {ev.get('text', '')}")
                md_lines.append("")
            with open(md_path, "w", encoding="utf-8") as f:
                f.write("\n".join(md_lines))
            print(f"Exported Evidence Register as Markdown: {md_path}")

        if self.export_csv:
            import csv
            csv_path = self.output_dir / "evidence_register.csv"
            # v3 #7: export the full field set (previously only 7 columns —
            # dropped source_file, column_name, ocr_used, extraction_method,
            # related_files, boosted_by_rule). Suggested_narrative_use is a list,
            # so join it for CSV.
            all_fields = []
            seen = set()
            for ev in self.evidence_register:
                for k in ev.keys():
                    if k not in seen:
                        seen.add(k)
                        all_fields.append(k)
            # Put the key fields first in a sensible order, then any extras.
            preferred = ["evidence_id", "source_file", "sheet_name", "column_name",
                         "insight_type", "extraction_method", "text", "priority_score",
                         "confidence", "suggested_narrative_use", "source_location",
                         "ocr_used", "related_files", "boosted_by_rule"]
            ordered = [f for f in preferred if f in seen] + [f for f in all_fields if f not in preferred]
            with open(csv_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.DictWriter(f, fieldnames=ordered, extrasaction="ignore")
                writer.writeheader()
                for ev in self.evidence_register:
                    row = dict(ev)
                    nau = row.get("suggested_narrative_use")
                    if isinstance(nau, list):
                        row["suggested_narrative_use"] = "/".join(nau)
                    rf = row.get("related_files")
                    if isinstance(rf, list):
                        row["related_files"] = ";".join(str(x) for x in rf)
                    writer.writerow(row)
            print(f"Exported Evidence Register as CSV: {csv_path}")

        print(f"Outputs saved to: {self.output_dir}")

    def _emit_run_metadata(self):
        """v3 #23: write run_metadata.json — the reproducibility artifact.

        Always emitted (not gated). Captures preprocessor version, git commit,
        run timestamps, the resolved config snapshot (from #21), per-stage
        timing (from #22), the optional-deps inventory (which fallback tiers
        were active), and high-level counts. So any past run can be traced to
        its exact code + config + environment and reproduced.
        """
        import sys
        metadata = {
            "preprocessor_version": __version__,
            "git": {
                "commit": git_commit(),
                "dirty": git_dirty(),
            },
            "run_id": self.run_id,
            "started_at": self.run_started_at,
            "finished_at": self.run_finished_at,
            "total_seconds": self.timing.get("total_seconds", 0),
            "timing": self.timing,
            "config": self.config_snapshot or {},
            "environment": {
                "python_version": sys.version.split()[0],
                "platform": platform.platform(),
                "optional_deps": {
                    "fitz_pymupdf": fitz is not None,
                    "pdfplumber": pdfplumber is not None,
                    "docx": Document is not None,
                    "pytesseract": pytesseract is not None,
                    "pydantic": _HAS_PYDANTIC,
                    "numpy": np is not None,
                    "rapidfuzz": self._has_rapidfuzz(),
                    "sentence_transformers": self._has_sentence_transformers(),
                    "yaml": _HAS_YAML,
                    "structlog": _HAS_STRUCTLOG,
                },
            },
            "counts": {
                "files_discovered": len(self.inventory),
                "files_processed": len(self.timing.get("files", [])),
                "evidence_entries": len(self.evidence_register) if hasattr(self, "evidence_register") else 0,
                "errors": len(self.errors),
                "filtered": len(self.filtered_items),
            },
        }
        # v4 #26: attach a briefing summary so run_metadata.json captures the
        # narrative-readiness snapshot without needing to open a second file.
        if getattr(self, "analyst_briefing_json", None):
            br = self.analyst_briefing_json
            metadata["briefing"] = {
                "overall_readiness_score": br.get("narrative_readiness", {}).get("overall_score", 0),
                "stage_scores": {
                    s: sc.get("score", 0) for s, sc in
                    br.get("narrative_readiness", {}).get("stage_scores", {}).items()
                },
                "components": br.get("narrative_readiness", {}).get("components", {}),
                "focus_areas": [
                    {"area": fa.get("area"), "score": fa.get("score"),
                     "evidence_count": fa.get("evidence_count")}
                    for fa in br.get("suggested_focus_areas", [])
                ],
                "quality_flags": br.get("quality_flags", []),
                "cross_file_links": len(br.get("top_cross_file_relationships", [])),
            }
        path = self.output_dir / "run_metadata.json"
        with open(path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=2)
        if self.log:
            self.log.info("run_metadata_emitted", path=str(path))

    @staticmethod
    def _has_rapidfuzz() -> bool:
        try:
            import rapidfuzz
            return True
        except ImportError:
            return False

    @staticmethod
    def _has_sentence_transformers() -> bool:
        try:
            import sentence_transformers
            return True
        except ImportError:
            return False

    def _generate_summary_report(self):
        """Generate a human-readable Markdown summary report."""
        lines = []
        lines.append("# Preprocessor Summary Report\n")
        lines.append(f"**Filter Level:** `{self.filter_level}`\n")
        lines.append(f"**Generated at:** {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")

        # v3 #22: Processing Time (always present — runtime is a top-level fact
        # about the run). Placed before File Inventory so total runtime is the
        # first thing the reader sees.
        lines.append("## Processing Time\n")
        st = self.timing.get("stages", {})
        lines.append(f"- **Total runtime:** {self.timing.get('total_seconds', 0):.2f}s\n")
        lines.append(f"- **Discovery:** {st.get('discovery', 0):.2f}s\n")
        n_files = len(self.timing.get("files", []))
        lines.append(f"- **Extraction:** {st.get('extraction', 0):.2f}s ({n_files} files)\n")
        lines.append(f"- **Evidence register build:** {st.get('evidence_build', 0):.2f}s\n")
        lines.append(f"- **Output & report:** {st.get('output', 0):.2f}s\n\n")
        if self.timing.get("files"):
            lines.append("### Per-File Timing\n")
            lines.append("| File | Category | Duration | Status |\n")
            lines.append("|------|----------|----------|--------|\n")
            for f in sorted(self.timing["files"], key=lambda x: x["duration_s"], reverse=True):
                lines.append(f"| {f['file']} | {f['category']} | {f['duration_s']:.2f}s | {f['status']} |\n")
            lines.append("\n")

        # File Inventory Summary
        total_files = len(self.inventory)
        categories = {}
        for item in self.inventory:
            cat = item.get("category", "other")
            categories[cat] = categories.get(cat, 0) + 1

        lines.append("## File Inventory\n")
        lines.append(f"- **Total files discovered:** {total_files}\n")
        for cat, count in sorted(categories.items()):
            lines.append(f"- **{cat.capitalize()}:** {count}\n")
        lines.append("\n")

        # Excel Processing Summary
        if self.excel_profiles:
            lines.append("## Excel Processing\n")
            total_sheets = sum(len(p.get("sheets", [])) for p in self.excel_profiles)
            lines.append(f"- **Excel files processed:** {len(self.excel_profiles)}\n")
            lines.append(f"- **Total sheets analyzed:** {total_sheets}\n\n")

        # PPTX Processing Summary (NEW - Item 1 integration)
        if self.pptx_profiles:
            lines.append("## PPTX Processing\n")
            total_slides = sum(p.get("total_slides", 0) for p in self.pptx_profiles if p.get("status") == "ok")
            high_priority_slides = 0
            data_rich_slides = 0

            for p in self.pptx_profiles:
                if p.get("status") != "ok":
                    continue
                for slide in p.get("slides", []):
                    cls = slide.get("classification", {})
                    if cls.get("priority_for_evidence", 0) >= 0.7:
                        high_priority_slides += 1
                    if cls.get("type") in ["data_chart", "data_table", "data_mixed", "conclusion"]:
                        data_rich_slides += 1

            lines.append(f"- **PPTX files processed:** {len(self.pptx_profiles)}\n")
            lines.append(f"- **Total slides analyzed:** {total_slides}\n")
            lines.append(f"- **High-priority slides (for evidence):** {high_priority_slides}\n")
            lines.append(f"- **Data-rich / Conclusion slides:** {data_rich_slides}\n\n")

        # Evidence Register Summary
        if hasattr(self, "evidence_register") and self.evidence_register:
            ev = self.evidence_register
            lines.append("## Evidence Register\n")
            lines.append(f"- **Total evidence entries:** {len(ev)}\n")

            if ev:
                avg_priority = sum(e.get("priority_score", 0) for e in ev) / len(ev)
                lines.append(f"- **Average priority score:** {avg_priority:.3f}\n")

                # Breakdown by insight type
                from collections import Counter
                types = Counter(e.get("insight_type", "unknown") for e in ev)
                lines.append("\n### Evidence by Type\n")
                for t, count in types.most_common(8):
                    lines.append(f"- `{t}`: {count}\n")

                # Top 5 evidence
                lines.append("\n### Top 5 Highest Priority Evidence\n")
                for e in ev[:5]:
                    lines.append(f"- **{e['evidence_id']}** (Score: {e.get('priority_score', 0):.2f}) — {e.get('text', '')[:120]}...\n")

                # v3 #4: coverage map section
                if hasattr(self, "coverage_map") and self.coverage_map:
                    cm = self.coverage_map
                    lines.append("\n### Coverage Map\n")
                    lines.append("Narrative-stage coverage (Why → What → How → Now):\n")
                    for s, c in cm.get("by_narrative_stage", {}).items():
                        lines.append(f"- `{s}`: {c}\n")
                    missing = cm.get("stages_with_no_evidence", [])
                    if missing:
                        lines.append(f"\n**Stages with no evidence:** {', '.join(missing)}\n")
                    lines.append("\nEvidence by source file:\n")
                    for src, c in cm.get("by_source_file", {}).items():
                        lines.append(f"- `{src}`: {c}\n")
            lines.append("\n")

        # v4 #26: Narrative Readiness section (mirrors analyst_briefing.json).
        if getattr(self, "analyst_briefing_json", None):
            br = self.analyst_briefing_json
            readiness = br.get("narrative_readiness", {})
            lines.append("## Narrative Readiness\n")
            lines.append(f"- **Overall readiness score:** {readiness.get('overall_score', 0)}/100\n")
            lines.append(f"- {readiness.get('explanation', '')}\n")
            ss = readiness.get("stage_scores", {})
            if ss:
                lines.append("\n| Stage | Score | Evidence | Avg Priority |\n")
                lines.append("|-------|------:|---------:|-------------:|\n")
                for s in ("Why", "What", "How", "Now"):
                    sc = ss.get(s, {})
                    lines.append(f"| {s} | {sc.get('score', 0)} | "
                                 f"{sc.get('evidence_count', 0)} | "
                                 f"{sc.get('avg_priority', 0):.3f} |\n")
                lines.append("\n")
            focus = br.get("suggested_focus_areas", [])
            if focus:
                lines.append("### Suggested Focus Areas\n")
                for fa in focus:
                    stages = ", ".join(fa.get("dominant_stages", [])) or "—"
                    lines.append(f"{fa.get('rank', 0)}. **{fa.get('area', '?')}** "
                                 f"(score {fa.get('score', 0):.0f}/100) — "
                                 f"{fa.get('evidence_count', 0)} items, stages: {stages}\n")
                lines.append("\n")
            flags = br.get("quality_flags", [])
            if flags:
                lines.append(f"**Quality flags:** {', '.join('`'+f+'`' for f in flags)}\n\n")

        # Filtering Summary
        if self.filtered_items:
            lines.append("## Filtering Summary\n")
            lines.append(f"- **Total items filtered:** {len(self.filtered_items)}\n\n")

            from collections import Counter
            reasons = Counter(item["reason"] for item in self.filtered_items)
            lines.append("### Reasons for Filtering\n")
            for reason, count in reasons.most_common():
                lines.append(f"- `{reason}`: {count} items\n")
            lines.append("\n")

        # Errors
        if self.errors:
            lines.append("## Errors Encountered\n")
            lines.append(f"- **Total errors:** {len(self.errors)}\n\n")

        # Output Files
        lines.append("## Generated Output Files\n")
        lines.append("- `file_inventory.json`\n")
        lines.append("- `excel_profile.json`\n")
        if _HAS_PYDANTIC:
            lines.append("- `evidence_schema.json` (Analyst GPT contract)\n")
        if self.pptx_profiles:
            lines.append("- `pptx_profile.json`\n")
        if hasattr(self, "evidence_register") and self.evidence_register:
            lines.append("- `evidence_register_seed.json`\n")
        if hasattr(self, "coverage_map") and self.coverage_map:
            lines.append("- `coverage_map.json`\n")
        if getattr(self, "entities_summary", None):
            lines.append("- `entities_summary.json`\n")
        if getattr(self, "analyst_briefing_md", None) is not None:
            lines.append("- `analyst_briefing.md` (v4 #26 strategic handoff)\n")
            lines.append("- `analyst_briefing.json` (structured briefing)\n")
        if self.filtered_items:
            lines.append("- `filtering_log.json`\n")
        lines.append("- `preprocessor_summary.md` (this file)\n")

        # PPTX Classification Summary (Item 1 deliverable)
        if self.pptx_profiles:
            lines.append("\n## PPTX Classification Summary\n")
            lines.append("| File | Slides | High-Priority | Data-Rich/Conclusion | Diagram/Process |\n")
            lines.append("|------|--------|---------------|----------------------|-----------------|\n")
            for p in self.pptx_profiles:
                if p.get("status") != "ok":
                    continue
                file_name = Path(p.get("file", "unknown.pptx")).name
                total = p.get("total_slides", 0)
                high = sum(1 for s in p.get("slides", []) if s.get("classification", {}).get("priority_for_evidence", 0) >= 0.7)
                data_rich = sum(1 for s in p.get("slides", []) if s.get("classification", {}).get("type") in ["data_chart", "data_table", "data_mixed", "conclusion"])
                diagram = sum(1 for s in p.get("slides", []) if s.get("classification", {}).get("type") == "diagram_process")
                lines.append(f"| {file_name} | {total} | {high} | {data_rich} | {diagram} |\n")

        # Write the report
        report_path = self.output_dir / "preprocessor_summary.md"
        with open(report_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))

        print(f"Summary report generated: {report_path}")

    # ====================== PPTX SLIDE CLASSIFICATION (v3 Improved) ======================

    def classify_slide(
        self,
        title: str = "",
        word_count: int = 0,
        chart_count: int = 0,
        table_count: int = 0,
        picture_count: int = 0,
        shape_count: int = 0,
        slide_idx: int = 1,
        total_slides: int = 10,
        has_diagram_like_shapes: bool = False,
        has_arrows_connectors: bool = False,
        has_process_diagram: bool = False,
        text_density: float = 0.0,   # words per shape
        is_section_slide: bool = False,   # Item 1.5
    ) -> Dict[str, Any]:
        """
        Production-grade slide classifier (v3).
        Returns rich classification dict:
          - type (primary category)
          - confidence (0-1)
          - evidence_tags
          - recommended_evidence_types
          - position, priority_for_evidence
        """
        title_lower = (title or "").lower().strip()
        is_first = slide_idx == 1
        is_last = slide_idx == total_slides
        total_visuals = chart_count + table_count + picture_count

        classification = {
            "type": "content_light",
            "confidence": 0.6,
            "evidence_tags": [],
            "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("content_light", ["What", "Why"]),
            "position": "middle",
            "priority_for_evidence": 0.5,
        }

        if is_first:
            classification["position"] = "first"
        elif is_last:
            classification["position"] = "last"

        # --- Strong title slide ---
        if is_first and word_count < 25 and total_visuals == 0:
            classification.update({"type": "title", "confidence": 0.95, "priority_for_evidence": 0.1})
            return classification

        # --- Agenda / Overview ---
        if any(kw in title_lower for kw in self.stage_rules.get("slide_type_keywords", {}).get("agenda", DEFAULT_SLIDE_TYPE_KEYWORDS["agenda"])):
            classification.update({"type": "agenda", "confidence": 0.9, "priority_for_evidence": 0.3,
                                   "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("agenda", ["What", "Why"])})
            return classification

        # --- Diagram / Process flow scoring (computed early so the section-divider
        #     guard below can use it; was previously referenced before assignment — bug fix) ---
        diagram_score = 0
        if has_diagram_like_shapes:
            diagram_score += 2
        if has_arrows_connectors:
            diagram_score += 2
        if has_process_diagram:
            diagram_score += 4   # strong signal
        if shape_count >= 6 and word_count < 70:
            diagram_score += 2

        # --- Section divider (Item 1.5) ---
        if is_section_slide or (any(kw in title_lower for kw in self.stage_rules.get("slide_type_keywords", {}).get("section", DEFAULT_SLIDE_TYPE_KEYWORDS["section"])) and diagram_score < 3):
            classification.update({"type": "section", "confidence": 0.88, "priority_for_evidence": 0.22,
                                   "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("section", ["What"])})
            return classification

        # --- Thank you / Q&A ---
        if is_last and any(kw in title_lower for kw in self.stage_rules.get("slide_type_keywords", {}).get("thank_you", DEFAULT_SLIDE_TYPE_KEYWORDS["thank_you"])):
            classification.update({"type": "thank_you", "confidence": 0.95, "priority_for_evidence": 0.1,
                                   "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("thank_you", ["What"])})
            return classification

        # --- Conclusion / Recommendations ---
        if any(kw in title_lower for kw in self.stage_rules.get("slide_type_keywords", {}).get("conclusion", DEFAULT_SLIDE_TYPE_KEYWORDS["conclusion"])):
            classification.update({"type": "conclusion", "confidence": 0.9, "priority_for_evidence": 0.85,
                                   "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("conclusion", ["What", "How", "Why", "Now"])})
            return classification

        # --- Data-rich slides (high evidence value) ---
        if chart_count >= 2 and table_count >= 1:
            classification.update({"type": "data_mixed", "confidence": 0.9, "priority_for_evidence": 0.95,
                                   "evidence_tags": ["numeric", "categorical"], "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("data_mixed", ["How", "What"])})
            return classification
        if chart_count >= 1:
            classification.update({"type": "data_chart", "confidence": 0.9, "priority_for_evidence": 0.9,
                                   "evidence_tags": ["numeric"], "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("data_chart", ["How", "What"])})
            return classification
        if table_count >= 1:
            classification.update({"type": "data_table", "confidence": 0.85, "priority_for_evidence": 0.85,
                                   "evidence_tags": ["categorical", "numeric"], "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("data_table", ["What", "Why"])})
            return classification

        # --- Diagram / Process flows (Item 4.3 improved) ---
        if diagram_score >= 3:
            priority = round(min(0.72, 0.45 + (diagram_score * 0.05)), 3)
            classification.update({
                "type": "diagram_process",
                "confidence": 0.85,
                "priority_for_evidence": priority,
                "evidence_tags": ["process", "diagram"]
            })
            return classification

        # --- Comparison slides ---
        if any(kw in title_lower for kw in self.stage_rules.get("slide_type_keywords", {}).get("comparison", DEFAULT_SLIDE_TYPE_KEYWORDS["comparison"])):
            classification.update({"type": "comparison", "confidence": 0.8, "priority_for_evidence": 0.75,
                                   "recommended_evidence_types": self.stage_rules.get("slide_type_stages", {}).get("comparison", ["What", "Why"])})

        # --- Quote / Callout ---
        if word_count < 40 and picture_count == 0 and shape_count <= 3:
            classification.update({"type": "quote_callout", "confidence": 0.7, "priority_for_evidence": 0.35})

        # --- Low-value visual-heavy ---
        if word_count < 12 and chart_count == 0 and table_count == 0 and picture_count >= 4:
            classification.update({"type": "low_value", "confidence": 0.75, "priority_for_evidence": 0.15})
            return classification

        # --- Content insight (default high-value text) ---
        if word_count >= 35 or total_visuals > 0:
            classification.update({"type": "content_insight", "confidence": 0.75, "priority_for_evidence": 0.65,
                                   "evidence_tags": ["text", "insight"]})
        else:
            classification.update({"type": "content_light", "confidence": 0.6, "priority_for_evidence": 0.35})

        return classification

    def extract_pptx(self, pptx_path: Path) -> Dict[str, Any]:
        """Scan a .pptx file and classify all slides using the improved classifier.
        Also extracts richer details for evidence generation from data-heavy slides."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        try:
            prs = Presentation(str(pptx_path))
            total_slides = len(prs.slides)
            slide_analyses = []

            for idx, slide in enumerate(prs.slides, 1):
                title = ""
                word_count = 0
                chart_count = 0
                table_count = 0
                picture_count = 0
                shape_count = 0
                text_shapes = 0
                has_connectors = False
                has_process_diagram = False
                diagram_shape_count = 0
                is_section_slide = False   # Item 1.5 (light)

                # Enhanced extraction
                chart_details = []
                table_details = []
                extracted_metrics = []
                extracted_advanced_metrics = []
                key_bullets = []
                process_steps = []   # Item 4.3 extension: texts from diagram shapes
                speaker_notes = ""   # Item 1.3
                bold_texts = []      # Item 1.6: emphasized text
                theme_colors = []    # Item 1.4 (light)

                embedded_objects = []   # v3 #15: embedded/linked OLE objects (counted, content unreadable)
                # v3 #13-16: iterate shapes deeply (recursing into groups) in
                # spatial (top, left) order so multi-column slides read in order
                # and nested/grouped text boxes are not lost.
                for shape in _iter_shapes_deep(slide.shapes):
                    shape_count += 1

                    # v3 #15: detect embedded/linked OLE objects (embedded Excel
                    # sheets, PDFs, …). Their content can't be read here, but we
                    # record them so the Analyst knows unread signal exists.
                    if _is_embedded_object(shape):
                        try:
                            embedded_objects.append(shape.name or "Embedded Object")
                        except Exception:
                            embedded_objects.append("Embedded Object")
                        continue

                    # v3 #14: richer text extraction — falls back to drawingml
                    # <a:t> runs for SmartArt/diagram graphic frames whose text
                    # isn't exposed via text_frame.
                    text = _extract_shape_text(shape)
                    if text:
                        has_tf = False
                        try:
                            has_tf = shape.has_text_frame
                        except Exception:
                            has_tf = False
                        if has_tf or text:
                            word_count += len(text.split())
                            text_shapes += 1

                            # Capture title from first substantial text.
                            # Skip pure page-number textboxes (e.g. '7', '42')
                            # which are commonly the first text on a slide but
                            # are not a title (bug #6 fix).
                            if not title and len(text) < 120 and not text.isdigit():
                                title = text[:80]

                            # Extract key bullets. Besides explicit bullet
                            # markers, also capture substantive plain-text
                            # lines (text-heavy decks rarely use bullet glyphs)
                            # so their content is seeded as evidence (bug #5 fix).
                            # v3 #11: skip lines that duplicate the slide title or
                            # look like navigation/section headers (short, title-
                            # cased) so repeated deck navigation text doesn't
                            # flood the bullet evidence.
                            title_norm = clean_text(title).lower()
                            for line in text.splitlines():
                                line = line.strip()
                                if not line or len(line) <= 10:
                                    continue
                                line_norm = line.lower()
                                # skip if it's (nearly) the slide title repeated
                                if title_norm and (line_norm == title_norm
                                                   or title_norm in line_norm
                                                   or line_norm in title_norm):
                                    continue
                                if line.startswith(("•", "-", "–", "▪", "*")):
                                    key_bullets.append(line.lstrip("•–▪*").strip())
                                elif (
                                    len(line.split()) >= 4
                                    and len(line) <= 200
                                    and not line.isdigit()
                                    and not line.startswith(("http://", "https://"))
                                ):
                                    key_bullets.append(line)

                            # Use advanced metric extraction
                            advanced = extract_advanced_metrics(text)
                            for m in advanced:
                                extracted_advanced_metrics.append(m)
                                extracted_metrics.append(m["value"])

                            # Item 1.6: Detect bold / emphasized text (only when
                            # a real text_frame is available for run-level access)
                            if has_tf:
                                try:
                                    for para in shape.text_frame.paragraphs:
                                        for run in para.runs:
                                            if run.font.bold:
                                                bold_text = run.text.strip()
                                                if bold_text and 5 < len(bold_text) < 120:
                                                    bold_texts.append(bold_text)
                                except Exception:
                                    pass

                    if shape.has_chart:
                        chart_count += 1
                        try:
                            chart = shape.chart
                            chart_title = ""
                            if chart.has_title and chart.chart_title:
                                chart_title = chart.chart_title.text_frame.text.strip()

                            # Item 1.1: Extract real chart data (categories + values)
                            categories = []
                            series_data = []

                            try:
                                plot = chart.plots[0]
                                if hasattr(plot, "categories") and plot.categories:
                                    categories = [str(c) for c in plot.categories][:10]  # limit

                                for series in chart.series:
                                    s_name = str(series.name) if series.name else "Series"
                                    values = []
                                    try:
                                        for val in series.values:
                                            if val is not None:
                                                values.append(float(val))
                                    except:
                                        pass
                                    series_data.append({
                                        "name": s_name,
                                        "values": values[:10]   # limit data points
                                    })
                            except Exception:
                                pass

                            chart_details.append({
                                "title": chart_title or "Chart",
                                "categories": categories,
                                "series_data": series_data,
                                "series": [s["name"] for s in series_data][:3]  # backward compat
                            })
                        except Exception:
                            chart_details.append({"title": "Chart", "series": [], "categories": [], "series_data": []})

                    if shape.has_table:
                        table_count += 1
                        try:
                            table = shape.table
                            cells = []
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    val = clean_text(cell.text)
                                    if val:
                                        cells.append({
                                            "row": row_idx,
                                            "col": col_idx,
                                            "value": val,
                                            "is_header": row_idx == 0
                                        })
                            table_details.append({
                                "rows": len(table.rows),
                                "cols": len(table.columns),
                                "cells": cells   # All cells (Item 1.2)
                            })
                        except Exception:
                            table_details.append({"rows": 0, "cols": 0, "cells": []})

                    if shape.shape_type == 13:
                        picture_count += 1
                    if hasattr(shape, "connector_format"):
                        has_connectors = True

                    # Item 4.3: Enhanced shape type detection for diagrams/process flows
                    try:
                        st = int(shape.shape_type)
                        # AUTO_SHAPE (1), FREEFORM (5), GROUP (6), DIAGRAM (21), and many flowchart shapes
                        if st in (1, 5, 6, 21) or (20 <= st <= 30):
                            diagram_shape_count += 1
                            if diagram_shape_count >= 3:
                                has_process_diagram = True

                            # NEW: Extract text from diagram shapes (process step labels)
                            if shape.has_text_frame:
                                txt = shape.text_frame.text.strip()
                                if txt and 3 < len(txt) < 80:
                                    process_steps.append(txt)

                        if st == 9:  # LINE
                            has_connectors = True
                    except Exception:
                        pass

                text_density = round(word_count / max(text_shapes, 1), 1) if text_shapes > 0 else 0.0

                # Item 4.1: Calculate insight language score for prioritization
                insight_score = contains_insight_language(text) if text_shapes > 0 else 0.0

                # Item 1.3: Extract speaker notes
                try:
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        speaker_notes = notes_slide.notes_text_frame.text.strip()
                except Exception:
                    speaker_notes = ""

                # Item 1.4 (light): Capture theme color hints from text runs
                try:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.color and run.font.color.type == 1:  # RGB
                                        try:
                                            rgb = run.font.color.rgb
                                            if rgb:
                                                theme_colors.append(str(rgb))
                                        except:
                                            pass
                except Exception:
                    pass

                # Item 1.5: Improved section divider detection
                title_lower_check = (title or "").lower()
                section_keywords = any(kw in title_lower_check for kw in ["section", "part ", "chapter", "module", "agenda"])
                looks_like_divider = (word_count < 10 and shape_count >= 5) or section_keywords
                if looks_like_divider:
                    is_section_slide = True

                classification = self.classify_slide(
                    title=title,
                    word_count=word_count,
                    chart_count=chart_count,
                    table_count=table_count,
                    picture_count=picture_count,
                    shape_count=shape_count,
                    slide_idx=idx,
                    total_slides=total_slides,
                    has_arrows_connectors=has_connectors,
                    has_process_diagram=has_process_diagram,
                    text_density=text_density,
                    is_section_slide=is_section_slide,   # Item 1.5
                )

                # Strong boost for data-rich slides with extracted content
                if extracted_metrics or key_bullets or chart_count > 0 or table_count > 0:
                    current_p = classification.get("priority_for_evidence", 0.5)
                    boost = 0.25 if (extracted_metrics or key_bullets) else 0.18
                    classification["priority_for_evidence"] = round(min(0.98, current_p + boost), 3)
                    if classification["type"] in ["content_light", "quote_callout"]:
                        classification["type"] = "content_insight"

                # Item 4.1: Additional boost from insight language
                if insight_score > 0.3:
                    extra = round(insight_score * 0.18, 3)
                    classification["priority_for_evidence"] = round(min(0.98, classification.get("priority_for_evidence", 0.5) + extra), 3)
                    if classification["type"] == "content_light":
                        classification["type"] = "content_insight"

                slide_analyses.append({
                    "slide_index": idx,
                    "title": title,
                    "word_count": word_count,
                    "visual_counts": {
                        "charts": chart_count,
                        "tables": table_count,
                        "pictures": picture_count,
                        "shapes": shape_count
                    },
                    "classification": classification,
                    "details": {
                        "chart_details": chart_details,
                        "table_details": table_details,
                        "extracted_metrics": extracted_metrics[:8],
                        "extracted_advanced_metrics": extracted_advanced_metrics[:8],
                        "key_bullets": key_bullets[:5],
                        "has_process_diagram": has_process_diagram,
                        "diagram_shape_count": diagram_shape_count,
                        "process_steps": process_steps[:8],
                        "speaker_notes": speaker_notes,       # Item 1.3
                        "bold_texts": bold_texts[:6],         # Item 1.6
                        "theme_colors": list(set(theme_colors))[:4],   # Item 1.4 (light)
                        "embedded_objects": embedded_objects,          # v3 #15: OLE objects (content unreadable)
                    }
                })

            # v3 #11: deck-wide navigation/section-header filter. A short text
            # line that appears on many slides is navigation (page footer,
            # section label, running header), not insight content. Drop bullets
            # whose normalized text repeats across >= 30% of slides.
            from collections import Counter
            bullet_freq = Counter()
            for sa in slide_analyses:
                for b in sa.get("details", {}).get("key_bullets", []):
                    bn = re.sub(r"\s+", " ", clean_text(b).lower())
                    if bn:
                        bullet_freq[bn] += 1
            repeat_threshold = max(3, int(len(slide_analyses) * 0.30))
            repeated = {bn for bn, c in bullet_freq.items() if c >= repeat_threshold}
            if repeated:
                for sa in slide_analyses:
                    details = sa.get("details", {})
                    details["key_bullets"] = [
                        b for b in details.get("key_bullets", [])
                        if re.sub(r"\s+", " ", clean_text(b).lower()) not in repeated
                    ]

            return {
                "file": str(pptx_path),
                "total_slides": total_slides,
                "slides": slide_analyses,
                "status": "ok"
            }
        except Exception as e:
            return {"file": str(pptx_path), "status": "error", "error": str(e)}

    # ====================== PDF SUPPORT (3.1 + 3.2) ======================

    def _ensure_tesseract(self) -> bool:
        """
        Make the Tesseract OCR binary reachable by pytesseract.
        Honors an explicit self.tesseract_cmd if set, otherwise auto-detects
        common Windows install locations. Result is cached in self._ocr_available.
        Returns True if Tesseract is usable.
        """
        if self._ocr_available is not None:
            return self._ocr_available
        if not pytesseract:
            self._ocr_available = False
            return False

        # Honor an explicitly configured path first.
        if self.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = self.tesseract_cmd

        try:
            pytesseract.get_tesseract_version()
            self._ocr_available = True
            return True
        except Exception:
            pass

        # Try common install locations (and PATH via shutil.which).
        import shutil
        candidates = [
            shutil.which("tesseract"),
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
            "/usr/bin/tesseract",
            "/usr/local/bin/tesseract",
        ]
        for c in candidates:
            if c and os.path.isfile(c):
                pytesseract.pytesseract.tesseract_cmd = c
                self.tesseract_cmd = c
                try:
                    pytesseract.get_tesseract_version()
                    self._ocr_available = True
                    return True
                except Exception:
                    continue

        self._ocr_available = False
        return False

    def _extract_pdf_tables(self, path: Path, engine: str = "auto") -> List[Dict]:
        """Extract tables from a PDF, merging the best of pdfplumber and PyMuPDF.

        - engine='auto' (default): prefer pdfplumber (better cell detection for
          ruled/unruled/merged tables), fall back to PyMuPDF's find_tables().
        - engine='pdfplumber': use pdfplumber only.
        - engine='pymupdf': use PyMuPDF only.

        Returns a list of {page, rows, cols, header, data, bbox, engine} dicts.
        Enriches over the old output (which only had rows + first 5 raw rows)
        so pdf_table_insight evidence can carry header + cell-level structure.
        """
        results: List[Dict] = []
        use_pdfplumber = (engine == "pdfplumber") or (engine == "auto" and pdfplumber is not None)
        use_pymupdf = fitz is not None and (engine == "pymupdf" or (engine == "auto" and not use_pdfplumber))

        if use_pdfplumber:
            try:
                import pdfplumber as _pp
                with _pp.open(str(path)) as pdf:
                    for pno, page in enumerate(pdf.pages, 1):
                        try:
                            found = page.find_tables()
                        except Exception:
                            found = []
                        for t in found:
                            try:
                                data = t.extract() or []
                            except Exception:
                                data = []
                            # normalize cells to strings and cap rows
                            norm = [[("" if c is None else str(c).strip()) for c in row]
                                    for row in data[:8]]
                            if not norm:
                                continue
                            header = norm[0] if norm else []
                            results.append({
                                "page": pno,
                                "rows": len(data),
                                "cols": max((len(r) for r in norm), default=0),
                                "header": header,
                                "data": norm,
                                "bbox": list(t.bbox) if getattr(t, "bbox", None) else None,
                                "engine": "pdfplumber",
                            })
            except Exception as e:
                if self.verbose:
                    print(f"       pdfplumber table extraction failed: {e}")
                # fall through to PyMuPDF if auto
                if engine == "auto":
                    use_pymupdf = fitz is not None

        if use_pymupdf and fitz is not None:
            try:
                doc = fitz.open(str(path))
                for pno, page in enumerate(doc, 1):
                    try:
                        tabs = page.find_tables()
                    except Exception:
                        tabs = []
                    for t in tabs:
                        try:
                            data = t.extract() or []
                        except Exception:
                            data = []
                        norm = [[("" if c is None else str(c).strip()) for c in row]
                                for row in data[:8]]
                        if not norm:
                            continue
                        header = norm[0] if norm else []
                        results.append({
                            "page": pno,
                            "rows": len(data),
                            "cols": max((len(r) for r in norm), default=0),
                            "header": header,
                            "data": norm,
                            "bbox": list(t.bbox) if getattr(t, "bbox", None) else None,
                            "engine": "pymupdf",
                        })
                doc.close()
            except Exception as e:
                if self.verbose:
                    print(f"       PyMuPDF table extraction failed: {e}")
        return results

    def extract_pdf(self, path: Path, use_ocr: bool = False) -> Dict[str, Any]:
        if fitz is None:
            return {"file": str(path), "status": "error", "error": "PyMuPDF not installed"}

        try:
            doc = fitz.open(str(path))
            total_pages = len(doc)
            pages = []

            # Resolve the Tesseract binary once for the whole document.
            ocr_available = use_ocr and self._ensure_tesseract()
            if use_ocr and not ocr_available and not getattr(self, "_ocr_warned", False):
                msg = "OCR enabled but Tesseract binary not found; scanned pages will yield no text."
                self.errors.append({"file": str(path), "error": msg})
                self._ocr_warned = True
                print(f"       [warn] {msg}")

            for page_num, page in enumerate(doc, 1):
                text = page.get_text().strip()
                ocr_used_this_page = False

                # OCR fallback logic (Phase 2)
                if ocr_available and len(text) < 30:
                    try:
                        # Render page at 300 DPI for good OCR quality
                        pix = page.get_pixmap(dpi=300)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        ocr_text = pytesseract.image_to_string(img).strip()
                        if len(ocr_text) > len(text):
                            text = ocr_text
                            ocr_used_this_page = True
                    except Exception as e:
                        if self.verbose:
                            print(f"       OCR failed on page {page_num}: {e}")
                        # OCR failed, keep original (possibly empty) text

                pages.append({
                    "page": page_num,
                    "text": text[:5000],
                    "ocr_used": ocr_used_this_page
                })

            doc.close()

            # Strip running headers/footers and isolated page numbers from
            # every page's text. A running header is a short line that
            # appears at the top (first 2 lines) or bottom (last 2 lines) of
            # >= 30% of pages — e.g. "Table of Contents" on 193/232 pages of
            # a 10-K. Isolated page-number lines (pure digits, F-33, roman
            # numerals, <= 6 chars) in the first/last 2 lines are also
            # stripped. Mirrors the v3 #11 PPTX deck-wide navigation filter.
            pages = self._strip_pdf_running_headers(pages)

            # v3: table extraction via the merged pdfplumber/PyMuPDF engine.
            # PyMuPDF stays primary for text/layout/OCR; pdfplumber (optional)
            # is preferred for table cell detection where it handles
            # merged/spanning/unruled tables better than PyMuPDF's default.
            tables = self._extract_pdf_tables(path, engine=self.pdf_table_engine)

            return {
                "file": str(path),
                "total_pages": total_pages,
                "pages": pages,
                "tables": tables,
                "status": "ok"
            }
        except Exception as e:
            return {"file": str(path), "status": "error", "error": str(e)}

    # ----- running header/footer stripping (PDF) -------------------------- #

    # A line is considered a page number if it is short (<= 6 chars) and
    # matches one of: pure digits, letter-dash-digits (F-33, A-12), or
    # roman numerals (i, ii, iii, iv, v, vi, ... up to xx).
    _PAGE_NUM_RE = re.compile(
        r"^(?:\d{1,5}|[A-Z]-\d{1,4}|[ivxlcmIVXLCM]{1,6})$"
    )

    def _strip_pdf_running_headers(self, pages: List[Dict]) -> List[Dict]:
        """Detect and strip running headers/footers + isolated page numbers
        from each page's text.

        A running header/footer is a short text line that appears at the top
        (first 2 lines) or bottom (last 2 lines) of >= 30% of pages (minimum
        3 pages). These are PDF boilerplate (e.g. "Table of Contents" on
        193/232 pages of a 10-K) that PyMuPDF includes in the page text but
        that carries no insight value.

        Isolated page-number lines (pure digits, F-33, roman numerals) in
        the first/last 2 lines are also stripped from every page.

        This mirrors the v3 #11 PPTX deck-wide navigation/section-header
        filter and is a pure function of the pages list.
        """
        if len(pages) < 3:
            return pages  # too few pages to detect running headers

        from collections import Counter

        # --- 1. detect running headers/footers by frequency -------------- #
        # Check the first 2 and last 2 non-empty lines of each page.
        first_lines: Counter = Counter()
        last_lines: Counter = Counter()
        page_line_cache: List[List[str]] = []  # cache non-empty lines per page

        for pg in pages:
            lines = [ln.strip() for ln in pg.get("text", "").split("\n")
                     if ln.strip()]
            page_line_cache.append(lines)
            for ln in lines[:2]:
                first_lines[ln] += 1
            for ln in lines[-2:]:
                last_lines[ln] += 1

        threshold = max(3, int(len(pages) * 0.30))
        running_headers = {ln for ln, c in first_lines.items() if c >= threshold}
        running_footers = {ln for ln, c in last_lines.items() if c >= threshold}
        boilerplate = running_headers | running_footers

        # --- 2. strip boilerplate + isolated page numbers ---------------- #
        for pg, lines in zip(pages, page_line_cache):
            if not lines:
                continue
            kept = []
            for i, ln in enumerate(lines):
                # Strip running headers/footers (any position).
                if ln in boilerplate:
                    continue
                # Strip isolated page numbers, but only in the first 2 or
                # last 2 lines (body page references like "See Note 5" must
                # survive).
                is_edge = i < 2 or i >= len(lines) - 2
                if is_edge and len(ln) <= 6 and self._PAGE_NUM_RE.match(ln):
                    continue
                kept.append(ln)
            pg["text"] = "\n".join(kept)[:5000]

        return pages

    # ====================== DOCX SUPPORT (3.3) ======================

    def extract_docx(self, path: Path) -> Dict[str, Any]:
        if Document is None:
            return {"file": str(path), "status": "error", "error": "python-docx not installed"}

        try:
            doc = Document(str(path))
            paragraphs = []
            tables = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                tables.append({"rows": len(rows), "data": rows[:5]})

            return {
                "file": str(path),
                "total_paragraphs": len(paragraphs),
                "paragraphs": paragraphs[:20],
                "tables": tables,
                "status": "ok"
            }
        except Exception as e:
            return {"file": str(path), "status": "error", "error": str(e)}

